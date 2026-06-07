"""
document_renderer/providers.py

Domain Layer - Data classes and rendering logic for document generation.

Contains:
- Enums: OutputFormat, PageSize, ChartType, TableStyle, etc.
- Data classes: RenderConfig, Section, ChartData, TableData, etc.
- Renderers: PDFRenderer, DOCXRenderer, PPTXRenderer
- Builders: ChartBuilder, TableBuilder

Version: 1.0.0
"""

from __future__ import annotations

import base64
import io
import logging
import os
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple, Union

logger = logging.getLogger(__name__)


# ============================================================================
# Enums
# ============================================================================


class OutputFormat(str, Enum):
    """Supported output formats."""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    MARKDOWN = "md"
    HTML = "html"


class PageSize(str, Enum):
    """Standard page sizes."""
    A4 = "A4"
    A3 = "A3"
    LETTER = "Letter"
    LEGAL = "Legal"
    TABLOID = "Tabloid"


class Orientation(str, Enum):
    """Page orientation."""
    PORTRAIT = "portrait"
    LANDSCAPE = "landscape"


class ChartType(str, Enum):
    """Supported chart types."""
    BAR = "bar"
    BAR_HORIZONTAL = "bar_horizontal"
    BAR_STACKED = "bar_stacked"
    LINE = "line"
    LINE_AREA = "area"
    PIE = "pie"
    DONUT = "donut"
    SCATTER = "scatter"
    BUBBLE = "bubble"
    RADAR = "radar"
    WATERFALL = "waterfall"
    HEATMAP = "heatmap"
    TREEMAP = "treemap"


class TableStyle(str, Enum):
    """Table styling options."""
    PROFESSIONAL = "professional"
    MINIMAL = "minimal"
    BORDERED = "bordered"
    STRIPED = "striped"
    COLORFUL = "colorful"


class ImagePosition(str, Enum):
    """Image positioning options."""
    INLINE = "inline"
    FLOAT_LEFT = "float_left"
    FLOAT_RIGHT = "float_right"
    CENTER = "center"
    FULL_WIDTH = "full_width"


class RenderStatus(str, Enum):
    """Rendering status."""
    SUCCESS = "success"
    PARTIAL = "partial"
    FAILED = "failed"


# ============================================================================
# Configuration Data Classes
# ============================================================================


@dataclass
class Margins:
    """Page margins in cm."""
    top: float = 2.5
    bottom: float = 2.5
    left: float = 2.5
    right: float = 2.5


@dataclass
class FontConfig:
    """Font configuration."""
    family: str = "Arial"
    size: int = 11
    color: str = "#000000"
    bold: bool = False
    italic: bool = False


@dataclass
class ColorScheme:
    """Color scheme for documents."""
    primary: str = "#1a73e8"
    secondary: str = "#5f6368"
    accent: str = "#ea4335"
    background: str = "#ffffff"
    text: str = "#202124"
    heading: str = "#1a73e8"
    link: str = "#1a73e8"
    border: str = "#dadce0"
    chart_colors: List[str] = field(default_factory=lambda: [
        "#1a73e8", "#ea4335", "#fbbc04", "#34a853",
        "#ff6d01", "#46bdc6", "#7baaf7", "#f07b72"
    ])


@dataclass 
class HeaderFooter:
    """Header/footer configuration."""
    text: str = ""
    include_page_number: bool = True
    include_date: bool = False
    include_title: bool = False
    font: FontConfig = field(default_factory=FontConfig)


@dataclass
class DocumentStyle:
    """Complete document styling."""
    page_size: PageSize = PageSize.A4
    orientation: Orientation = Orientation.PORTRAIT
    margins: Margins = field(default_factory=Margins)
    colors: ColorScheme = field(default_factory=ColorScheme)
    body_font: FontConfig = field(default_factory=FontConfig)
    heading_font: FontConfig = field(default_factory=lambda: FontConfig(size=14, bold=True))
    header: Optional[HeaderFooter] = None
    footer: Optional[HeaderFooter] = None
    line_spacing: float = 1.15

    def to_dict(self) -> Dict[str, Any]:
        return {
            "page_size": self.page_size.value,
            "orientation": self.orientation.value,
            "line_spacing": self.line_spacing,
        }


# ============================================================================
# Content Data Classes
# ============================================================================


@dataclass
class ChartData:
    """Data for chart generation."""
    chart_type: ChartType
    title: str = ""
    labels: List[str] = field(default_factory=list)
    datasets: List[Dict[str, Any]] = field(default_factory=list)
    x_label: str = ""
    y_label: str = ""
    show_legend: bool = True
    show_grid: bool = True
    width: int = 600
    height: int = 400
    options: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "type": self.chart_type.value,
            "title": self.title,
            "labels": self.labels,
            "datasets": self.datasets,
        }


@dataclass
class TableData:
    """Data for table generation."""
    headers: List[str] = field(default_factory=list)
    rows: List[List[Any]] = field(default_factory=list)
    style: TableStyle = TableStyle.PROFESSIONAL
    caption: str = ""
    column_widths: Optional[List[float]] = None
    alignment: List[str] = field(default_factory=list)
    highlight_header: bool = True
    alternate_rows: bool = True

    def to_dict(self) -> Dict[str, Any]:
        return {
            "headers": self.headers,
            "row_count": len(self.rows),
            "style": self.style.value,
        }


@dataclass
class ImageData:
    """Data for image insertion."""
    source: str  # Path, URL, or base64
    caption: str = ""
    alt_text: str = ""
    width: Optional[int] = None
    height: Optional[int] = None
    position: ImagePosition = ImagePosition.CENTER

    def to_dict(self) -> Dict[str, Any]:
        return {
            "caption": self.caption,
            "position": self.position.value,
        }


@dataclass
class Section:
    """Document section."""
    id: str
    title: str
    content: str = ""
    level: int = 1
    charts: List[ChartData] = field(default_factory=list)
    tables: List[TableData] = field(default_factory=list)
    images: List[ImageData] = field(default_factory=list)
    subsections: List['Section'] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    page_break_before: bool = False  # Bug 5 Fix: Add missing field

    def to_dict(self) -> Dict[str, Any]:
        return {
            "id": self.id,
            "title": self.title,
            "level": self.level,
            "content_length": len(self.content),
            "charts": len(self.charts),
            "tables": len(self.tables),
            "page_break_before": self.page_break_before,  # Include new field
        }


@dataclass
class DocumentContent:
    """Complete document content."""
    title: str
    subtitle: str = ""
    authors: List[str] = field(default_factory=list)
    date: Optional[str] = None
    abstract: str = ""
    sections: List[Section] = field(default_factory=list)
    bibliography: List[Dict[str, Any]] = field(default_factory=list)
    appendices: List[Section] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "section_count": len(self.sections),
            "has_bibliography": len(self.bibliography) > 0,
        }


@dataclass
class RenderConfig:
    """Configuration for rendering."""
    format: OutputFormat = OutputFormat.PDF
    style: DocumentStyle = field(default_factory=DocumentStyle)
    include_toc: bool = True
    include_cover: bool = True
    include_page_numbers: bool = True
    include_header: bool = False
    include_footer: bool = True
    compress: bool = False
    output_path: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "format": self.format.value,
            "include_toc": self.include_toc,
            "include_cover": self.include_cover,
        }


@dataclass
class RenderResult:
    """Result of rendering operation."""
    success: bool
    format: OutputFormat
    content: Optional[bytes] = None
    content_base64: Optional[str] = None
    path: Optional[str] = None
    page_count: int = 0
    file_size: int = 0
    time_ms: float = 0.0
    error: Optional[str] = None
    warnings: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "success": self.success,
            "format": self.format.value,
            "page_count": self.page_count,
            "file_size": self.file_size,
            "time_ms": self.time_ms,
            "error": self.error,
        }


# ============================================================================
# PPTX Specific Data Classes
# ============================================================================


@dataclass
class SlideLayout:
    """Slide layout configuration."""
    layout_type: str = "title_content"  # title, title_content, two_column, blank, etc.
    background_color: Optional[str] = None
    background_image: Optional[str] = None


@dataclass
class SlideContent:
    """Content for a single slide."""
    title: str
    content: str = ""
    bullet_points: List[str] = field(default_factory=list)
    speaker_notes: str = ""
    layout: SlideLayout = field(default_factory=SlideLayout)
    charts: List[ChartData] = field(default_factory=list)
    tables: List[TableData] = field(default_factory=list)
    images: List[ImageData] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "bullet_count": len(self.bullet_points),
            "has_notes": bool(self.speaker_notes),
        }


@dataclass
class PresentationContent:
    """Complete presentation content."""
    title: str
    subtitle: str = ""
    author: str = ""
    date: Optional[str] = None
    slides: List[SlideContent] = field(default_factory=list)
    theme: str = "professional"
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "slide_count": len(self.slides),
            "theme": self.theme,
        }


# ============================================================================
# Chart Builder
# ============================================================================


class ChartBuilder:
    """Builds charts for documents."""

    def __init__(self, style: Optional[ColorScheme] = None):
        self.style = style or ColorScheme()
        self._matplotlib_available = self._check_matplotlib()

    def _check_matplotlib(self) -> bool:
        """Check if matplotlib is available."""
        try:
            import matplotlib
            return True
        except ImportError:
            return False

    def build_chart(
        self,
        chart_data: ChartData,
        output_format: str = "png",
        dpi: int = 150,
    ) -> Optional[bytes]:
        """Build a chart and return as bytes."""
        if not self._matplotlib_available:
            logger.warning("Matplotlib not available, skipping chart")
            return None

        try:
            import matplotlib.pyplot as plt
            import matplotlib
            matplotlib.use('Agg')

            fig, ax = plt.subplots(figsize=(
                chart_data.width / 100,
                chart_data.height / 100
            ))

            # Build based on type
            if chart_data.chart_type == ChartType.BAR:
                self._build_bar(ax, chart_data)
            elif chart_data.chart_type == ChartType.LINE:
                self._build_line(ax, chart_data)
            elif chart_data.chart_type == ChartType.PIE:
                self._build_pie(ax, chart_data)
            elif chart_data.chart_type == ChartType.SCATTER:
                self._build_scatter(ax, chart_data)
            else:
                self._build_bar(ax, chart_data)  # Default

            # Add labels
            if chart_data.title:
                ax.set_title(chart_data.title)
            if chart_data.x_label:
                ax.set_xlabel(chart_data.x_label)
            if chart_data.y_label:
                ax.set_ylabel(chart_data.y_label)
            if chart_data.show_legend and chart_data.datasets:
                ax.legend()
            if chart_data.show_grid:
                ax.grid(True, alpha=0.3)

            # Save to bytes
            buf = io.BytesIO()
            plt.savefig(buf, format=output_format, dpi=dpi, bbox_inches='tight')
            plt.close(fig)
            buf.seek(0)
            return buf.read()

        except Exception as e:
            logger.error(f"Chart build failed: {e}")
            return None

    def _build_bar(self, ax, data: ChartData) -> None:
        """Build bar chart."""
        import numpy as np
        x = np.arange(len(data.labels))
        width = 0.8 / max(len(data.datasets), 1)
        
        for i, dataset in enumerate(data.datasets):
            offset = (i - len(data.datasets) / 2 + 0.5) * width
            color = self.style.chart_colors[i % len(self.style.chart_colors)]
            ax.bar(x + offset, dataset.get('data', []), width,
                   label=dataset.get('label', ''), color=color)
        
        ax.set_xticks(x)
        ax.set_xticklabels(data.labels)

    def _build_line(self, ax, data: ChartData) -> None:
        """Build line chart."""
        for i, dataset in enumerate(data.datasets):
            color = self.style.chart_colors[i % len(self.style.chart_colors)]
            ax.plot(data.labels, dataset.get('data', []),
                    label=dataset.get('label', ''), color=color, marker='o')

    def _build_pie(self, ax, data: ChartData) -> None:
        """Build pie chart."""
        if data.datasets:
            values = data.datasets[0].get('data', [])
            colors = self.style.chart_colors[:len(values)]
            ax.pie(values, labels=data.labels, colors=colors, autopct='%1.1f%%')
            ax.axis('equal')

    def _build_scatter(self, ax, data: ChartData) -> None:
        """Build scatter chart."""
        for i, dataset in enumerate(data.datasets):
            color = self.style.chart_colors[i % len(self.style.chart_colors)]
            x_data = dataset.get('x', list(range(len(dataset.get('data', [])))))
            y_data = dataset.get('data', [])
            ax.scatter(x_data, y_data, label=dataset.get('label', ''), color=color)


# ============================================================================
# Table Builder
# ============================================================================


class TableBuilder:
    """Builds tables for documents."""

    def __init__(self, style: Optional[ColorScheme] = None):
        self.style = style or ColorScheme()

    def build_html_table(self, table_data: TableData) -> str:
        """Build HTML table."""
        style_class = self._get_style_class(table_data.style)
        
        html = f'<table class="{style_class}">'
        
        # Caption
        if table_data.caption:
            html += f'<caption>{table_data.caption}</caption>'
        
        # Header
        if table_data.headers:
            html += '<thead><tr>'
            for header in table_data.headers:
                html += f'<th>{header}</th>'
            html += '</tr></thead>'
        
        # Body
        html += '<tbody>'
        for i, row in enumerate(table_data.rows):
            row_class = 'even' if i % 2 == 0 else 'odd'
            html += f'<tr class="{row_class}">'
            for j, cell in enumerate(row):
                align = table_data.alignment[j] if j < len(table_data.alignment) else 'left'
                html += f'<td style="text-align:{align}">{cell}</td>'
            html += '</tr>'
        html += '</tbody></table>'
        
        return html

    def _get_style_class(self, style: TableStyle) -> str:
        """Get CSS class for table style."""
        return f"table-{style.value}"

    def get_table_css(self) -> str:
        """Get CSS for table styles."""
        return """
        .table-professional { border-collapse: collapse; width: 100%; }
        .table-professional th { background: #f8f9fa; border-bottom: 2px solid #dee2e6; }
        .table-professional td { border-bottom: 1px solid #dee2e6; padding: 8px; }
        .table-minimal { border-collapse: collapse; }
        .table-minimal td, .table-minimal th { padding: 8px; }
        .table-bordered { border: 1px solid #dee2e6; }
        .table-bordered td, .table-bordered th { border: 1px solid #dee2e6; padding: 8px; }
        .table-striped tr.odd { background: #f8f9fa; }
        """


# ============================================================================
# Base Renderer
# ============================================================================


class BaseRenderer(ABC):
    """Base class for document renderers."""

    def __init__(self, style: Optional[DocumentStyle] = None):
        self.style = style or DocumentStyle()
        self.chart_builder = ChartBuilder(self.style.colors)
        self.table_builder = TableBuilder(self.style.colors)

    @abstractmethod
    async def render(
        self,
        content: DocumentContent,
        config: RenderConfig,
    ) -> RenderResult:
        """Render document."""
        pass


# ============================================================================
# PDF Renderer
# ============================================================================


class PDFRenderer(BaseRenderer):
    """Renders documents to PDF."""

    def __init__(self, style: Optional[DocumentStyle] = None):
        super().__init__(style)
        self._reportlab_available = self._check_reportlab()
        self._weasyprint_available = self._check_weasyprint()

    def _check_reportlab(self) -> bool:
        try:
            import reportlab
            return True
        except ImportError:
            return False

    def _check_weasyprint(self) -> bool:
        try:
            import weasyprint
            return True
        except ImportError:
            return False

    async def render(
        self,
        content: DocumentContent,
        config: RenderConfig,
    ) -> RenderResult:
        """Render to PDF."""
        import time
        start = time.perf_counter()
        
        try:
            if self._reportlab_available:
                pdf_bytes = await self._render_with_reportlab(content, config)
            elif self._weasyprint_available:
                pdf_bytes = await self._render_with_weasyprint(content, config)
            else:
                return RenderResult(
                    success=False,
                    format=OutputFormat.PDF,
                    error="No PDF library available (install reportlab or weasyprint)",
                )

            elapsed = (time.perf_counter() - start) * 1000
            
            return RenderResult(
                success=True,
                format=OutputFormat.PDF,
                content=pdf_bytes,
                content_base64=base64.b64encode(pdf_bytes).decode() if pdf_bytes else None,
                file_size=len(pdf_bytes) if pdf_bytes else 0,
                time_ms=elapsed,
            )
        except Exception as e:
            return RenderResult(
                success=False,
                format=OutputFormat.PDF,
                error=str(e),
                time_ms=(time.perf_counter() - start) * 1000,
            )

    async def _render_with_reportlab(
        self,
        content: DocumentContent,
        config: RenderConfig,
    ) -> bytes:
        """Render using ReportLab."""
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.units import cm
        
        buffer = io.BytesIO()
        
        page_size = A4 if self.style.page_size == PageSize.A4 else letter
        doc = SimpleDocTemplate(
            buffer,
            pagesize=page_size,
            leftMargin=self.style.margins.left * cm,
            rightMargin=self.style.margins.right * cm,
            topMargin=self.style.margins.top * cm,
            bottomMargin=self.style.margins.bottom * cm,
        )
        
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
        )
        story.append(Paragraph(content.title, title_style))
        
        if content.subtitle:
            story.append(Paragraph(content.subtitle, styles['Heading2']))
        
        story.append(Spacer(1, 20))
        
        # Sections
        for section in content.sections:
            story.append(Paragraph(section.title, styles['Heading1']))
            if section.content:
                for para in section.content.split('\n\n'):
                    if para.strip():
                        story.append(Paragraph(para, styles['Normal']))
                        story.append(Spacer(1, 10))
        
        doc.build(story)
        buffer.seek(0)
        return buffer.read()

    async def _render_with_weasyprint(
        self,
        content: DocumentContent,
        config: RenderConfig,
    ) -> bytes:
        """Render using WeasyPrint."""
        html = self._content_to_html(content)
        from weasyprint import HTML
        return HTML(string=html).write_pdf()

    def _content_to_html(self, content: DocumentContent) -> str:
        """Convert content to HTML."""
        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{content.title}</title>
    <style>
        body {{ font-family: {self.style.body_font.family}; font-size: {self.style.body_font.size}pt; }}
        h1 {{ color: {self.style.colors.heading}; }}
        {self.table_builder.get_table_css()}
    </style>
</head>
<body>
    <h1>{content.title}</h1>
"""
        if content.subtitle:
            html += f"<h2>{content.subtitle}</h2>"
        
        for section in content.sections:
            html += f"<h{section.level + 1}>{section.title}</h{section.level + 1}>"
            if section.content:
                html += f"<p>{section.content}</p>"
        
        html += "</body></html>"
        return html


# ============================================================================
# DOCX Renderer
# ============================================================================


class DOCXRenderer(BaseRenderer):
    """Renders documents to DOCX."""

    def __init__(self, style: Optional[DocumentStyle] = None):
        super().__init__(style)
        self._docx_available = self._check_docx()

    def _check_docx(self) -> bool:
        try:
            import docx
            return True
        except ImportError:
            return False

    async def render(
        self,
        content: DocumentContent,
        config: RenderConfig,
    ) -> RenderResult:
        """Render to DOCX."""
        import time
        start = time.perf_counter()
        
        if not self._docx_available:
            return RenderResult(
                success=False,
                format=OutputFormat.DOCX,
                error="python-docx not available",
            )
        
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            
            doc = Document()
            
            # Title
            doc.add_heading(content.title, 0)
            if content.subtitle:
                doc.add_paragraph(content.subtitle)
            
            # Sections
            for section in content.sections:
                doc.add_heading(section.title, level=section.level)
                if section.content:
                    doc.add_paragraph(section.content)
            
            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            docx_bytes = buffer.read()
            
            elapsed = (time.perf_counter() - start) * 1000
            
            return RenderResult(
                success=True,
                format=OutputFormat.DOCX,
                content=docx_bytes,
                content_base64=base64.b64encode(docx_bytes).decode(),
                file_size=len(docx_bytes),
                time_ms=elapsed,
            )
        except Exception as e:
            return RenderResult(
                success=False,
                format=OutputFormat.DOCX,
                error=str(e),
                time_ms=(time.perf_counter() - start) * 1000,
            )


# ============================================================================
# PPTX Renderer
# ============================================================================


class PPTXRenderer(BaseRenderer):
    """Renders presentations to PPTX."""

    def __init__(self, style: Optional[DocumentStyle] = None):
        super().__init__(style)
        self._pptx_available = self._check_pptx()

    def _check_pptx(self) -> bool:
        try:
            import pptx
            return True
        except ImportError:
            return False

    async def render(
        self,
        content: Union[DocumentContent, PresentationContent],
        config: RenderConfig,
    ) -> RenderResult:
        """Render to PPTX."""
        import time
        start = time.perf_counter()
        
        if not self._pptx_available:
            return RenderResult(
                success=False,
                format=OutputFormat.PPTX,
                error="python-pptx not available",
            )
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            if isinstance(content, PresentationContent):
                title.text = content.title
                subtitle_shape.text = content.subtitle or ""
                
                # Content slides
                for slide_content in content.slides:
                    self._add_content_slide(prs, slide_content)
            else:
                title.text = content.title
                subtitle_shape.text = content.subtitle or ""
                
                # Convert sections to slides
                for section in content.sections:
                    self._add_section_slide(prs, section)
            
            # Save
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            pptx_bytes = buffer.read()
            
            elapsed = (time.perf_counter() - start) * 1000
            
            return RenderResult(
                success=True,
                format=OutputFormat.PPTX,
                content=pptx_bytes,
                content_base64=base64.b64encode(pptx_bytes).decode(),
                file_size=len(pptx_bytes),
                page_count=len(prs.slides),
                time_ms=elapsed,
            )
        except Exception as e:
            return RenderResult(
                success=False,
                format=OutputFormat.PPTX,
                error=str(e),
                time_ms=(time.perf_counter() - start) * 1000,
            )

    def _add_content_slide(self, prs, slide_content: SlideContent) -> None:
        """Add a content slide."""
        from pptx.util import Pt
        
        layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(layout)
        
        title = slide.shapes.title
        title.text = slide_content.title
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        if slide_content.bullet_points:
            for i, point in enumerate(slide_content.bullet_points):
                if i == 0:
                    tf.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        elif slide_content.content:
            tf.text = slide_content.content

    def _add_section_slide(self, prs, section: Section) -> None:
        """Add slide from section."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        
        title = slide.shapes.title
        title.text = section.title
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        if section.content:
            # Split into bullet points
            points = section.content.split('\n')
            for i, point in enumerate(points[:5]):  # Max 5 points
                if point.strip():
                    if i == 0:
                        tf.text = point.strip()
                    else:
                        p = tf.add_paragraph()
                        p.text = point.strip()


# ============================================================================
# Additional Config Classes
# ============================================================================


@dataclass
class TableCell:
    """A table cell."""
    value: Any = ""
    colspan: int = 1
    rowspan: int = 1
    style: Optional[FontConfig] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "value": str(self.value),
            "colspan": self.colspan,
            "rowspan": self.rowspan,
        }


@dataclass
class PdfOptions:
    """PDF-specific rendering options."""
    page_size: PageSize = PageSize.A4
    orientation: Orientation = Orientation.PORTRAIT
    include_toc: bool = True
    include_cover: bool = True
    include_page_numbers: bool = True
    include_header: bool = False
    include_footer: bool = True
    include_bookmarks: bool = True
    compress: bool = False
    encrypt: bool = False
    watermark: Optional[str] = None
    watermark_text: Optional[str] = None
    password: Optional[str] = None
    font_family: str = "Helvetica"
    font_size: int = 11

    def to_dict(self) -> Dict[str, Any]:
        return {
            "page_size": self.page_size.value,
            "orientation": self.orientation.value,
            "include_toc": self.include_toc,
            "include_cover": self.include_cover,
            "compress": self.compress,
            "encrypt": self.encrypt,
            "font_family": self.font_family,
            "font_size": self.font_size,
        }


@dataclass
class DocxOptions:
    """DOCX-specific rendering options."""
    template_path: Optional[str] = None
    include_toc: bool = True
    track_changes: bool = False
    compatibility_mode: bool = False

    def to_dict(self) -> Dict[str, Any]:
        return {
            "include_toc": self.include_toc,
            "has_template": self.template_path is not None,
        }


@dataclass
class PptxOptions:
    """PPTX-specific rendering options."""
    template_path: Optional[str] = None
    include_title_slide: bool = True
    include_toc_slide: bool = False
    include_speaker_notes: bool = True
    transition: str = "fade"

    def to_dict(self) -> Dict[str, Any]:
        return {
            "include_title_slide": self.include_title_slide,
            "include_speaker_notes": self.include_speaker_notes,
        }


@dataclass
class XlsxOptions:
    """XLSX-specific rendering options."""
    sheet_name: str = "Sheet1"
    auto_filter: bool = False
    freeze_header: bool = True
    column_widths: Optional[List[int]] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "sheet_name": self.sheet_name,
            "freeze_header": self.freeze_header,
        }


@dataclass
class ChartResult:
    """Result of chart generation."""
    success: bool
    chart_id: str = ""
    content: Optional[bytes] = None
    format: str = "png"
    width: int = 0
    height: int = 0
    error: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "success": self.success,
            "chart_id": self.chart_id,
            "format": self.format,
            "error": self.error,
        }


@dataclass
class MultiRenderResult:
    """Result of multi-format rendering."""
    success: bool
    results: Dict[str, RenderResult] = field(default_factory=dict)
    total_time_ms: float = 0.0
    errors: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "success": self.success,
            "formats_rendered": list(self.results.keys()),
            "total_time_ms": self.total_time_ms,
            "errors": self.errors,
        }


@dataclass
class TemplateInfo:
    """Information about a template."""
    id: str
    name: str
    format: OutputFormat
    description: str = ""
    path: str = ""
    preview_path: Optional[str] = None
    tags: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "id": self.id,
            "name": self.name,
            "format": self.format.value,
            "description": self.description,
        }


# ============================================================================
# Aliases for backward compatibility
# ============================================================================

StyleConfig = DocumentStyle
SectionContent = Section
PdfRenderer = PDFRenderer


# ============================================================================
# Exports
# ============================================================================


__all__ = [
    # Enums
    "OutputFormat",
    "PageSize",
    "Orientation",
    "ChartType",
    "TableStyle",
    "ImagePosition",
    "RenderStatus",
    # Config classes
    "Margins",
    "FontConfig",
    "ColorScheme",
    "HeaderFooter",
    "DocumentStyle",
    "StyleConfig",  # Alias
    # Content classes
    "ChartData",
    "TableData",
    "TableCell",
    "ImageData",
    "Section",
    "SectionContent",  # Alias
    "DocumentContent",
    "RenderConfig",
    "RenderResult",
    # PPTX classes
    "SlideLayout",
    "SlideContent",
    "PresentationContent",
    # Option classes
    "PdfOptions",
    "DocxOptions",
    "PptxOptions",
    "XlsxOptions",
    # Result classes
    "ChartResult",
    "MultiRenderResult",
    "TemplateInfo",
    # Builders
    "ChartBuilder",
    "TableBuilder",
    # Renderers
    "BaseRenderer",
    "PDFRenderer",
    "PdfRenderer",  # Alias
    "DOCXRenderer",
    "PPTXRenderer",
]

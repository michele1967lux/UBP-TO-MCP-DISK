"""
PowerPoint Renderer for Multi-Format Artifact Engine (v2.5.1 - FEAT-ARTIFACT-002)

Provides rendering for PPTX (PowerPoint) format.

Features:
- Title and content slides
- Section dividers
- Bullet points with proper formatting
- Two-column layouts
- Configurable slide dimensions (4:3, 16:9)
- Professional styling

Architecture:
- AI generates slide content as JSON (slide_chunking mode)
- Python handles layout, formatting, and file generation

Usage:
    from renderers.pptx_renderer import PptxRenderer
    from renderers.base import RenderContext

    context = RenderContext(
        title="My Presentation",
        content="...",
        extracted_data=[
            {"slide_number": 1, "type": "title", "title": "Welcome"},
            {"slide_number": 2, "type": "content", "title": "Overview", "bullets": ["A", "B"]}
        ]
    )

    renderer = PptxRenderer()
    result = renderer.render(context)
"""

import io
import logging
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional
from dataclasses import dataclass
from enum import Enum

from .base import (
    FormatRenderer,
    OutputFormat,
    RenderContext,
    RenderResult,
    DataExtractor,
)

logger = logging.getLogger(__name__)


class SlideType(str, Enum):
    """Types of slides that can be generated."""

    TITLE = "title"           # Title slide (first slide)
    CONTENT = "content"       # Standard content slide with bullets
    SECTION = "section"       # Section divider/header
    TWO_COLUMN = "two_column" # Two-column layout
    BLANK = "blank"           # Blank slide for custom content
    SUMMARY = "summary"       # Summary/conclusion slide


class SlideLayout(int, Enum):
    """PowerPoint slide layout indices."""

    TITLE_SLIDE = 0           # Title and subtitle
    TITLE_AND_CONTENT = 1     # Title with body content
    SECTION_HEADER = 2        # Section header
    TWO_CONTENT = 3           # Two content areas
    COMPARISON = 4            # Comparison layout
    TITLE_ONLY = 5            # Title only
    BLANK = 6                 # Blank slide


@dataclass
class SlideContent:
    """Content for a single slide."""

    slide_number: int
    slide_type: SlideType
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = None
    left_content: List[str] = None
    right_content: List[str] = None
    notes: str = ""

    def __post_init__(self):
        if self.bullets is None:
            self.bullets = []
        if self.left_content is None:
            self.left_content = []
        if self.right_content is None:
            self.right_content = []


class PptxRenderer(FormatRenderer):
    """
    Renderer for PowerPoint PPTX format.

    Creates professional presentations from structured slide content.
    Supports multiple slide types and layouts with configurable styling.
    """

    def __init__(self, settings: Optional[Dict[str, Any]] = None):
        """Initialize PPTX renderer."""
        super().__init__(settings)

        # Slide dimensions (in EMUs - English Metric Units)
        # 1 inch = 914400 EMUs
        width_inches = self._get_setting("pptx_default_width_inches", 13.333)
        height_inches = self._get_setting("pptx_default_height_inches", 7.5)

        self._width_emu = int(width_inches * 914400)
        self._height_emu = int(height_inches * 914400)

        # Font settings
        self._default_font = self._get_setting("pptx_default_font", "Calibri")
        self._title_font_size = self._get_setting("pptx_title_font_size", 44)
        self._body_font_size = self._get_setting("pptx_body_font_size", 18)

        # Limits
        self._max_slides = self._get_setting("pptx_max_slides", 50)
        self._max_bullets_per_slide = 6

    @property
    def output_format(self) -> OutputFormat:
        """Return PPTX output format."""
        return OutputFormat.PPTX

    def render(self, context: RenderContext) -> RenderResult:
        """
        Render presentation to PPTX format.

        Args:
            context: RenderContext with extracted_data or content

        Returns:
            RenderResult with PPTX bytes
        """
        import time

        start_time = time.time()
        self._log_render_start(context)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            logger.error("[PPTX] python-pptx not installed")
            return RenderResult.failure(
                format=self.output_format,
                error="python-pptx is required for PowerPoint export",
                filename="",
            )

        try:
            # Get slide content
            slides = self._get_slides(context)
            if not slides:
                return RenderResult.failure(
                    format=self.output_format,
                    error="No slide content available for presentation",
                    filename=self.generate_filename(context.title, context.blueprint_id),
                )

            # Limit slides
            if len(slides) > self._max_slides:
                logger.warning(
                    f"[PPTX] Truncating slides from {len(slides)} to {self._max_slides}"
                )
                slides = slides[:self._max_slides]

            # Create presentation
            prs = Presentation()
            prs.slide_width = self._width_emu
            prs.slide_height = self._height_emu

            # Generate slides
            for slide_content in slides:
                self._add_slide(prs, slide_content)

            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            pptx_bytes = buffer.getvalue()

            # Create result
            duration_ms = int((time.time() - start_time) * 1000)
            filename = self.generate_filename(context.title, context.blueprint_id)

            result = RenderResult(
                content=pptx_bytes,
                format=self.output_format,
                filename=filename,
                render_time_ms=duration_ms,
                metadata={
                    "slides": len(slides),
                    "width_inches": self._width_emu / 914400,
                    "height_inches": self._height_emu / 914400,
                },
            )

            self._log_render_complete(result, duration_ms)
            return result

        except Exception as e:
            logger.error(f"[PPTX] Render error: {e}", exc_info=True)
            return RenderResult.failure(
                format=self.output_format,
                error=str(e),
                filename=self.generate_filename(context.title, context.blueprint_id),
            )

    def _get_slides(self, context: RenderContext) -> List[SlideContent]:
        """Extract slide content from context."""
        slides = []

        # First try extracted_data
        data = context.extracted_data
        if not data:
            data = DataExtractor.extract_json(context.content)

        if not data:
            # Generate basic slides from markdown content
            return self._slides_from_markdown(context.content, context.title)

        if isinstance(data, list):
            for item in data:
                slide = self._parse_slide_data(item)
                if slide:
                    slides.append(slide)
        elif isinstance(data, dict):
            # Single slide
            slide = self._parse_slide_data(data)
            if slide:
                slides.append(slide)

        return slides

    def _parse_slide_data(self, data: Dict) -> Optional[SlideContent]:
        """Parse a dictionary into SlideContent."""
        if not data:
            return None

        slide_type_str = data.get("type", "content")
        try:
            slide_type = SlideType(slide_type_str.lower())
        except ValueError:
            slide_type = SlideType.CONTENT

        return SlideContent(
            slide_number=data.get("slide_number", 0),
            slide_type=slide_type,
            title=data.get("title", ""),
            subtitle=data.get("subtitle", ""),
            bullets=data.get("bullets", []),
            left_content=data.get("left", data.get("left_content", [])),
            right_content=data.get("right", data.get("right_content", [])),
            notes=data.get("notes", ""),
        )

    def _slides_from_markdown(self, content: str, title: str) -> List[SlideContent]:
        """
        Generate slides from markdown content.

        Falls back to this when no structured slide data is provided.
        """
        slides = []
        import re

        # Add title slide
        slides.append(SlideContent(
            slide_number=1,
            slide_type=SlideType.TITLE,
            title=title,
            subtitle=f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
        ))

        # Parse markdown sections
        lines = content.split('\n')
        current_slide = None
        current_bullets = []

        for line in lines:
            # Check for heading (new slide)
            heading_match = re.match(r'^(#{1,3})\s+(.+)$', line)
            if heading_match:
                # Save previous slide
                if current_slide:
                    current_slide.bullets = current_bullets[:self._max_bullets_per_slide]
                    slides.append(current_slide)

                level = len(heading_match.group(1))
                text = heading_match.group(2).strip()

                if level == 1:
                    # Section header
                    current_slide = SlideContent(
                        slide_number=len(slides) + 1,
                        slide_type=SlideType.SECTION,
                        title=text,
                    )
                    current_bullets = []
                else:
                    # Content slide
                    current_slide = SlideContent(
                        slide_number=len(slides) + 1,
                        slide_type=SlideType.CONTENT,
                        title=text,
                    )
                    current_bullets = []
                continue

            # Check for bullet point
            bullet_match = re.match(r'^[\s]*[-*]\s+(.+)$', line)
            if bullet_match and current_slide:
                bullet_text = bullet_match.group(1).strip()
                # Limit bullet length
                if len(bullet_text) > 80:
                    bullet_text = bullet_text[:77] + "..."
                current_bullets.append(bullet_text)

        # Save last slide
        if current_slide:
            current_slide.bullets = current_bullets[:self._max_bullets_per_slide]
            slides.append(current_slide)

        return slides

    def _add_slide(self, prs, slide_content: SlideContent) -> None:
        """Add a slide to the presentation."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        # Determine layout
        if slide_content.slide_type == SlideType.TITLE:
            layout_idx = SlideLayout.TITLE_SLIDE.value
        elif slide_content.slide_type == SlideType.SECTION:
            layout_idx = SlideLayout.SECTION_HEADER.value
        elif slide_content.slide_type == SlideType.TWO_COLUMN:
            layout_idx = SlideLayout.TWO_CONTENT.value
        elif slide_content.slide_type == SlideType.BLANK:
            layout_idx = SlideLayout.BLANK.value
        else:
            layout_idx = SlideLayout.TITLE_AND_CONTENT.value

        # Get layout (handle missing layouts gracefully)
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title and slide_content.title:
            slide.shapes.title.text = slide_content.title

        # Handle different slide types
        if slide_content.slide_type == SlideType.TITLE:
            self._format_title_slide(slide, slide_content)
        elif slide_content.slide_type == SlideType.SECTION:
            self._format_section_slide(slide, slide_content)
        elif slide_content.slide_type == SlideType.TWO_COLUMN:
            self._format_two_column_slide(slide, slide_content)
        else:
            self._format_content_slide(slide, slide_content)

        # Add notes if present
        if slide_content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes

    def _format_title_slide(self, slide, content: SlideContent) -> None:
        """Format a title slide."""
        from pptx.util import Pt

        # Find subtitle placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = content.subtitle
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(20)
                break

    def _format_section_slide(self, slide, content: SlideContent) -> None:
        """Format a section header slide."""
        from pptx.util import Pt

        # Section slides typically just have a title
        if slide.shapes.title:
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True

    def _format_content_slide(self, slide, content: SlideContent) -> None:
        """Format a standard content slide with bullets."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        # Find content placeholder
        content_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_shape = shape
                break

        if content_shape and content.bullets:
            tf = content_shape.text_frame
            tf.clear()

            for i, bullet in enumerate(content.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = bullet
                p.font.size = Pt(self._body_font_size)
                p.font.name = self._default_font
                p.level = 0

    def _format_two_column_slide(self, slide, content: SlideContent) -> None:
        """Format a two-column slide."""
        from pptx.util import Pt

        # Find left and right content placeholders
        left_shape = None
        right_shape = None

        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 1:
                left_shape = shape
            elif idx == 2:
                right_shape = shape

        # Populate left column
        if left_shape and content.left_content:
            tf = left_shape.text_frame
            tf.clear()
            for i, item in enumerate(content.left_content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(self._body_font_size - 2)

        # Populate right column
        if right_shape and content.right_content:
            tf = right_shape.text_frame
            tf.clear()
            for i, item in enumerate(content.right_content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(self._body_font_size - 2)

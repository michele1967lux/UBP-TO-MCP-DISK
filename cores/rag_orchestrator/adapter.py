"""
RAG Orchestrator Adapter - UBP Bridge Layer

This module provides the UBP framework integration for RAG orchestration.
Acts as a bridge between UBP's module system and technical RAG implementations.

Features:
- Dependency Injection (rag_qdrant, inference_ollama_grok)
- Access Control Lists (ACL) via Redis
- Per-user/client RAG configuration
- Complete RAG pipeline orchestration
- Security context enforcement

MCP-COMPAT (ARCH-008): Added OperationContext support for dual REST/MCP compatibility.
"""

from pathlib import Path
from types import SimpleNamespace
from typing import Dict, Any, List, Optional, Union
import logging
import json
import uuid
import re
import time
import asyncio
import hashlib
from datetime import datetime
import redis.asyncio as aioredis
from fastapi import HTTPException

from ubp_enterprise_hybrid.modules.cores._shared import (
    BaseHybridModule,
    PLATFORM_ADMIN_CLIENT_ID,
    SYSTEM_USER_ID,
    SYSTEM_COLLECTIONS,
    EntityType,
    AccessLevel,
    PERSONAL_KB_PREFIX,
    CLIENT_KB_PREFIX,
    is_platform_admin_client,
    normalize_client_id,
    is_personal_kb,
    ErrorMessages,
    # v1.10.0: Role-Based Configuration Engine
    ProviderMapper,
    ProviderConfigurationError,
)
# MCP-COMPAT (ARCH-008): Import OperationContext for dual path support
from ubp_enterprise_hybrid.modules.cores._shared.operation_context import (
    OperationContext,
    is_operation_context_like,
)

from .providers import (
    RAGPipeline,
    DocumentChunker,
    ACLManager,
    ConfigManager,
    ConversationManager,
    KeywordManager,
)
from .semantic_router import (
    SemanticRouter,
    RouteType,
    RouterResult,
    RouterConfig,
    SmartFallbackStrategy,
)
from .capabilities import CapabilityManager

# v2.3/v2.6: Interactive Analyst - Report Session Manager
try:
    from .agents import (
        ReportSessionManager,
        ReportState,
        Researcher,
        SourcePreference,
    )
    # v2.6: Import enrichment config and debug structures
    from .agents.report_session import (
        SectionEnrichmentConfig,
        WorkerDebugEvent,
        DebugEventBatch,
    )

    REPORT_AGENTS_AVAILABLE = True
except ImportError as e:
    ReportSessionManager = None  # type: ignore
    ReportState = None  # type: ignore
    Researcher = None  # type: ignore
    SourcePreference = None  # type: ignore
    SectionEnrichmentConfig = None  # type: ignore
    WorkerDebugEvent = None  # type: ignore
    DebugEventBatch = None  # type: ignore
    REPORT_AGENTS_AVAILABLE = False
    logging.getLogger(__name__).warning(f"Report agents not available: {e}")

# Batch Ingestion Job model
from .agents.ingestion_job import IngestJob, IngestJobState, IngestFileEntry, FileStatus

# Import settings manager for efficient settings access (v3.7.1+)
try:
    from ubp_enterprise_hybrid.backend.app.api.admin_settings_routes import settings_manager

    _settings = None  # Lazy-load via settings_manager.get_settings()
    _router_settings = None
except ImportError:
    _settings = None
    _router_settings = None
    logger = logging.getLogger(__name__)
    logger.warning("Could not import settings_manager, using defaults")

# Configure structured logging
logger = logging.getLogger(__name__)

# FIX-PROP-002 v4.3.0: Mapping override_config enrichment keys → pipeline_options keys
# Used to propagate user/architect overrides into the 4-layer resolution (Layer 4: User Request)
ENRICHMENT_PROPAGATION_KEYS = {
    "query_expansion_enabled": "query_expansion",
    "hyde_enabled": "hyde",
    "investigative_enabled": "investigative",
    "reranking_enabled": "rerank",
    "rerank_enabled": "rerank",
    "fusion_enabled": "fusion",
    "dedup_enabled": "dedup",
    "compression_enabled": "compression",
}


def _build_router_config_from_settings() -> RouterConfig:
    """
    Build RouterConfig from Pydantic Settings (12-Factor compliant).

    Maps UBP_ROUTER__* environment variables to RouterConfig dataclass.
    Falls back to defaults if settings not available.
    """
    global _settings, _router_settings
    # Lazy-load settings on first call (BUG-001 fix: were never populated)
    if _settings is None:
        try:
            _settings = settings_manager.get_settings()
            _router_settings = _settings.router if _settings else None
        except Exception as e:
            logger.warning(f"Could not load settings: {e}")
    if _router_settings is None:
        logger.info("Using default RouterConfig (settings not available)")
        return RouterConfig()

    # Map fallback strategy string to enum
    strategy_map = {
        "conservative": SmartFallbackStrategy.CONSERVATIVE,
        "smart": SmartFallbackStrategy.SMART,
        "aggressive": SmartFallbackStrategy.AGGRESSIVE,
    }
    fallback_strategy = strategy_map.get(
        _router_settings.fallback_strategy.lower(), SmartFallbackStrategy.SMART
    )

    # Map default route string to enum
    route_map = {
        "CHAT": RouteType.CHAT,
        "RAG": RouteType.RAG,
        "WEB": RouteType.WEB,
    }
    default_route = route_map.get(
        _router_settings.default_route_for_unknown.upper(), RouteType.WEB
    )

    config = RouterConfig(
        confidence_threshold=_router_settings.confidence_threshold,
        fallback_strategy=fallback_strategy,
        default_route_for_unknown=default_route,
        enable_auto_retry=_router_settings.enable_auto_retry,
        max_retry_attempts=_router_settings.max_retry_attempts,
        empty_rag_triggers_web=_router_settings.empty_rag_triggers_web,
        commercial_boost_threshold=_router_settings.commercial_boost_threshold,
    )

    logger.info(
        f"RouterConfig built from settings: "
        f"confidence={config.confidence_threshold}, "
        f"fallback={config.fallback_strategy.value}, "
        f"auto_retry={config.enable_auto_retry}"
    )
    return config


# =============================================================================
# SECURITY: Sensitive Data Filter (VULN-008 Mitigation)
# =============================================================================
class SensitiveDataFilter:
    """
    Filter to redact sensitive data from RAG responses.

    Mitigates VULN-008 (Data Leakage) by sanitizing potential secrets,
    API keys, passwords, and authentication tokens from LLM outputs.
    """

    PATTERNS = [
        r"(password|passwd|pwd)\s*[:=]\s*\S+",
        r"(api[_-]?key|access[_-]?token|secret[_-]?key)\s*[:=]\s*\S+",
        r"Authorization:\s*Bearer\s+\S+",
        r"sk-[a-zA-Z0-9]{20,}",  # OpenAI/Generic API keys
        r"(aws[_-]?access[_-]?key[_-]?id|aws[_-]?secret)\s*[:=]\s*\S+",
        r"[a-zA-Z0-9]{32,}(?=.*[A-Z])(?=.*[a-z])(?=.*[0-9])",  # Long mixed-case tokens
    ]

    @staticmethod
    def sanitize(text: str) -> str:
        """
        Remove sensitive patterns from text.

        Args:
            text: Input text potentially containing sensitive data

        Returns:
            Sanitized text with sensitive patterns replaced by [REDACTED]
        """
        if not text:
            return text
        sanitized = text
        for pattern in SensitiveDataFilter.PATTERNS:
            sanitized = re.sub(pattern, "[REDACTED]", sanitized, flags=re.IGNORECASE)
        return sanitized


# =============================================================================
# METADATA ENRICHMENT: Auto-tag extraction & kb_id injection
# =============================================================================

# Stop-words excluded from auto-tag extraction (numbers, noise tokens)
_TAG_STOPWORDS = frozenset({
    "md", "txt", "pdf", "json", "yaml", "csv", "html", "xml", "docx", "xlsx",
    "the", "and", "for", "with", "from", "this", "that", "manual", "readme",
})

# Fields stripped from API-facing chunk responses (retrieve_enriched, retrieve_with_acl).
# Embedding vectors and query metadata are internal to the retrieval engine.
# Internal callers needing raw vectors should use _retrieve() directly.
_INTERNAL_ONLY_KEYS = frozenset({"embedding", "vector", "query_source"})


def _extract_tags_from_filename(filename: str) -> List[str]:
    """Extract semantic tags from a filename.

    Splits on separators, removes numeric-only tokens, extensions and common
    stop-words to produce meaningful keyword tags.

    Examples:
        "MANUAL_07_INFERENCE_PROVIDERS.md" → ["inference", "providers"]
        "report_cardiovascular_2026.pdf"   → ["report", "cardiovascular"]
    """
    if not filename:
        return []
    # Remove extension
    name = filename.rsplit(".", 1)[0] if "." in filename else filename
    parts = re.split(r"[_\-\s]+", name.lower())
    return [p for p in parts if len(p) > 2 and not p.isdigit() and p not in _TAG_STOPWORDS]


def _enrich_ingest_metadata(
    metadata: Dict[str, Any],
    collection_id: str,
    *,
    explicit_tags: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Enrich metadata before Qdrant upsert.

    Injects:
      - ``kb_id``   = collection_id  (so NL filter on kb_id works)
      - ``collection`` = collection_id (for downstream pipeline mapping)
      - ``tags``    = explicit_tags OR auto-extracted from filename

    This is the single injection point for chunk-level metadata enrichment.
    A future centralised tag-extraction module can hook here.
    """
    metadata["kb_id"] = collection_id
    metadata["collection"] = collection_id

    if explicit_tags is not None:
        metadata["tags"] = explicit_tags
    elif "tags" not in metadata:
        metadata["tags"] = _extract_tags_from_filename(metadata.get("filename", ""))
    return metadata


class RagOrchestratorAdapter(BaseHybridModule):
    """
    UBP Adapter for RAG Orchestration.

    This class handles:
    - UBP lifecycle integration (initialize, shutdown, health_check)
    - Dependency injection (rag_qdrant, inference_ollama_grok, Redis)
    - ACL management via ACLManager provider
    - RAG configuration management via ConfigManager provider
    - Full RAG pipeline orchestration
    - Security context (ctx) enforcement
    """

    def __init__(self, module_path: Path, **kwargs):
        """Initialize the RAG orchestrator adapter."""
        super().__init__(module_path, **kwargs)

        # Injected dependencies (will be set during initialize)
        self.qdrant_module = None
        self.llm_module = None
        self.redis_client = None
        self.memory_module = None  # ROADMAP v1.5.0 - FEAT-MEM-001

        # Provider instances
        self.rag_pipeline: Optional[RAGPipeline] = None
        self.acl_manager: Optional[ACLManager] = None
        self.config_manager: Optional[ConfigManager] = None
        self.conversation_manager: Optional[ConversationManager] = None  # FEAT-MEM-001
        self.semantic_router: Optional[SemanticRouter] = None  # FEAT-ROUTER-001

        # v3.5.0: Dedicated Architect RAG pipeline with isolated provider
        # This pipeline BYPASSES standard provider resolution to ensure
        # the Architect always uses its configured providers
        self.architect_pipeline: Optional[RAGPipeline] = None
        self.architect_llm_module = None  # Dedicated LLM module for Architect

        # Optional modules (resolved via DI, may be None)
        self.web_search_module = None  # FEAT-ROUTER-001 - Web search routing
        self.enrichment_module = None  # v1.7.1 - GPU Reranking

        # v2.3: Interactive Analyst - Report Generation
        self.report_session_manager = None  # ReportSessionManager instance
        self.researcher = None  # Researcher agent instance

        # Capability aggregation
        self.capability_manager: Optional[CapabilityManager] = None

        # Statistics
        self.total_queries = 0
        self.total_documents_ingested = 0

        # v4.3.0: Conversation context cache for parallel execution optimization
        # Cache key: f"{user_id}:{conversation_id}", value: (context_str, timestamp)
        self._conversation_cache: Dict[str, tuple[str, float]] = {}
        self._conversation_cache_ttl = 300  # 5 minutes TTL
        self._conversation_cache_max_entries = 100  # LRU eviction threshold

    async def initialize(self) -> Dict[str, Any]:
        """Initialize RAG Orchestrator with dependency injection."""
        logger.info(
            f"Initializing {self.manifest.name} module",
            extra={"mod_name": self.manifest.name},
        )

        # Inject dependencies via DI container
        if not self.di_container:
            raise RuntimeError("DI container not available")

        # === Infrastructure Services (Typed Resolution) ===
        try:
            self.redis_client = await self.di_container.resolve(aioredis.Redis)
            if not self.redis_client:
                raise RuntimeError("Redis client not found in DI container")
            logger.info("✅ Redis client resolved from DI container")
        except Exception as e:
            logger.error(f"Failed to resolve Redis client: {e}")
            raise

        # === Core Modules (String-based Resolution) ===
        try:
            self.qdrant_module = await self.di_container.resolve("rag_qdrant")
            if not self.qdrant_module:
                raise RuntimeError("rag_qdrant module not found in DI container")
            logger.info("✅ rag_qdrant module resolved")
        except Exception as e:
            logger.error(f"Failed to resolve rag_qdrant: {e}")
            raise

        # === LLM Module (v1.10.0: Role-Based Configuration Engine) ===
        # === High Availability Fallback (v1.10.1) ===
        # Resolution: UBP_ROLE_RAG_PROVIDER -> ProviderMapper -> Module
        # Fallback: UBP_ROLE_RAG_FALLBACK_PROVIDER -> Secondary Module
        # Supports: grok, ollama, vllm, openai, anthropic
        try:
            # Use ProviderMapper.resolve_chain() for HA support (v1.10.1)
            provider_chain = ProviderMapper.resolve_chain("rag")

            if not provider_chain:
                raise ProviderConfigurationError(
                    "No valid providers configured for role 'rag'. "
                    "Check UBP_ROLE_RAG_PROVIDER and UBP_PROVIDER_*__ENABLED in .env"
                )

            # Store the chain: [(module_name, provider_name), ...]
            self.llm_modules_chain = []

            for idx, (module_name, provider_name) in enumerate(provider_chain):
                role_type = "Primary" if idx == 0 else "Fallback"
                logger.info(
                    f"[{role_type}] Resolving '{provider_name}' -> module '{module_name}'"
                )

                # Resolve module from DI container
                module = await self.di_container.resolve(module_name)
                if not module:
                    logger.warning(
                        f"[{role_type}] {module_name} not found in DI container - skipping"
                    )
                    continue

                # Configure provider if supported
                if hasattr(module, "set_default_provider"):
                    module.set_default_provider(provider_name)
                    logger.info(
                        f"[{role_type}] ✅ {module_name} configured with provider '{provider_name}'"
                    )

                self.llm_modules_chain.append((module, provider_name))

            if not self.llm_modules_chain:
                raise RuntimeError(
                    "No LLM modules could be resolved from provider chain"
                )

            # Primary module for backward compatibility
            self.llm_module, primary_provider = self.llm_modules_chain[0]

            # Log HA status
            if len(self.llm_modules_chain) > 1:
                fallback_module, fallback_provider = self.llm_modules_chain[1]
                logger.info(
                    f"✅ RAG LLM with HA: Primary='{primary_provider}', "
                    f"Fallback='{fallback_provider}'"
                )
            else:
                logger.info(
                    f"✅ RAG LLM resolved: '{primary_provider}' (no fallback configured)"
                )

        except ProviderConfigurationError as e:
            logger.error(f"Role-Based Configuration Error: {e}")
            raise RuntimeError(f"RAG LLM configuration error: {e}")
        except Exception as e:
            logger.error(f"Failed to resolve RAG LLM module: {e}")
            raise

        # === Conversation Memory Module (Optional - FEAT-MEM-001) ===
        try:
            self.memory_module = await self.di_container.resolve(
                "rag_conversation_memory"
            )
            if self.memory_module:
                logger.info("✅ rag_conversation_memory module resolved")
            else:
                logger.warning(
                    "rag_conversation_memory module not available - sessions disabled"
                )
        except Exception as e:
            logger.warning(
                f"Could not resolve rag_conversation_memory: {e} - sessions disabled"
            )
            self.memory_module = None

        # === Initialize Providers (Business Logic Layer) ===
        self.acl_manager = ACLManager(self.redis_client, self.config["acl"])
        logger.info("✅ ACLManager provider initialized")

        self.config_manager = ConfigManager(
            self.redis_client, self.config["rag_config_storage"]
        )
        logger.info("✅ ConfigManager provider initialized")

        # Initialize ConversationManager (FEAT-MEM-001 - Task #15)
        self.conversation_manager = ConversationManager(
            self.redis_client, self.config.get("conversation_storage", {})
        )
        logger.info("✅ ConversationManager provider initialized")

        # Initialize KeywordManager (FEAT-DKI-001 - Dynamic Knowledge Injection v2.0)
        # Extracts keywords from documents during ingestion for intelligent routing
        # v2.0: Configurable via UBP_DKI__* env vars, supports batch and summarize
        from ubp_enterprise_hybrid.backend.app.core.config import DKISettings

        dki_settings = DKISettings()

        # Resolve DKI LLM provider (dedicated or RAG primary)
        dki_llm_module = self.llm_module  # default: RAG primary
        if dki_settings.provider:
            try:
                from ubp_enterprise_hybrid.modules.cores._shared import ProviderMapper as _DKIMapper

                mapping = _DKIMapper.PROVIDER_MAP.get(dki_settings.provider.lower())
                if mapping:
                    module_name, provider_name = mapping
                    dki_llm_module = await self.di_container.resolve(module_name)
                    logger.info(
                        f"DKI: Using dedicated provider {dki_settings.provider} ({module_name})"
                    )
                else:
                    logger.warning(
                        f"DKI: Unknown provider '{dki_settings.provider}', falling back to RAG primary"
                    )
            except Exception as e:
                logger.warning(
                    f"DKI: Failed to resolve provider '{dki_settings.provider}': {e}, "
                    "falling back to RAG primary"
                )

        self.keyword_manager = KeywordManager(
            redis_client=self.redis_client,
            llm_module=dki_llm_module,
            config=dki_settings,
        )
        logger.info(
            f"✅ KeywordManager provider initialized (DKI v2.0, "
            f"enabled={dki_settings.enabled}, batch_size={dki_settings.batch_size}, "
            f"summarize={dki_settings.summarize_enabled})"
        )

        # === Semantic Router (FEAT-ROUTER-001 - Intelligent Query Routing) ===
        # Build RouterConfig from Pydantic Settings (12-Factor compliant)
        router_config = _build_router_config_from_settings()
        self.semantic_router = SemanticRouter(
            llm_module=self.llm_module,
            default_route=RouteType.RAG,
            confidence_threshold=router_config.confidence_threshold,
            router_config=router_config,
        )
        logger.info(
            f"✅ SemanticRouter initialized with 12-Factor config "
            f"(confidence={router_config.confidence_threshold})"
        )

        # === Web Search Module (Optional - FEAT-ROUTER-001) ===
        try:
            self.web_search_module = await self.di_container.resolve("web_search")
            if self.web_search_module:
                logger.info("✅ web_search module resolved - WEB routing enabled")
            else:
                logger.warning(
                    "web_search module not available - WEB routing will fallback to RAG"
                )
        except Exception as e:
            logger.warning(
                f"Could not resolve web_search: {e} - WEB routing will fallback to RAG"
            )
            self.web_search_module = None

        # === Enrichment Pipeline Module (Optional - v1.7.1 GPU Reranking) ===
        # LAZY RESOLUTION: enrichment_pipeline initializes AFTER rag_orchestrator
        # So we create a resolver callable for lazy resolution on first use
        self.enrichment_module = None  # Will be resolved lazily
        self._enrichment_resolution_attempted = False

        async def enrichment_resolver():
            """Lazy resolver for enrichment_pipeline module."""
            if self._enrichment_resolution_attempted:
                return self.enrichment_module
            self._enrichment_resolution_attempted = True
            try:
                self.enrichment_module = await self.di_container.resolve(
                    "enrichment_pipeline"
                )
                return self.enrichment_module
            except Exception as e:
                logger.warning(f"Could not resolve enrichment_pipeline: {e}")
                return None
        
        # === RAG Reranker Module (v3.7.0 - Dedicated Reranking) ===
        # Resolve rag_reranker module for dedicated reranking operations
        self.rag_reranker_module = None
        self._reranker_resolution_attempted = False
        
        async def reranker_resolver():
            """Lazy resolver for rag_reranker module."""
            if self._reranker_resolution_attempted:
                return self.rag_reranker_module
            self._reranker_resolution_attempted = True
            try:
                self.rag_reranker_module = await self.di_container.resolve("rag_reranker")
                logger.info("✅ rag_reranker module resolved successfully")
                return self.rag_reranker_module
            except Exception as e:
                logger.warning(f"Could not resolve rag_reranker: {e}")
                return None

        # === Adaptive Budget Manager (v4.0.0 - Dynamic Token Budget) ===
        # Manages prompt/token allocation to prevent "Prompt exceeds maximum length"
        # by dynamically adjusting context size based on LLM limits
        self.adaptive_memory_module = None
        self._adaptive_memory_resolution_attempted = False

        async def adaptive_memory_resolver():
            """Lazy resolver for adaptive_budget_manager module."""
            if self._adaptive_memory_resolution_attempted:
                return self.adaptive_memory_module
            self._adaptive_memory_resolution_attempted = True
            try:
                self.adaptive_memory_module = await self.di_container.resolve(
                    "adaptive_budget_manager"
                )
                logger.info("adaptive_budget_manager module resolved successfully")
                return self.adaptive_memory_module
            except Exception as e:
                logger.warning(f"Could not resolve adaptive_budget_manager: {e}")
                return None

        # === RAG Pipeline Chain (v1.10.1: HA Fallback Support) ===
        # Create a RAGPipeline for each provider in the chain
        self.rag_pipelines_chain = []
        for llm_module, provider_name in self.llm_modules_chain:
            pipeline = RAGPipeline(
                self.qdrant_module,
                llm_module,
                enrichment_resolver,  # Pass resolver callable for lazy resolution
                reranker_resolver,    # v3.7.0: Add reranker resolver
                adaptive_memory_resolver,  # v4.0.0: Dynamic token budget management
            )
            self.rag_pipelines_chain.append((pipeline, provider_name))
            logger.info(f"✅ RAG pipeline created for provider '{provider_name}'")

        # Primary pipeline for backward compatibility
        self.rag_pipeline = self.rag_pipelines_chain[0][0]
        logger.info(
            f"✅ RAG pipeline chain initialized: {len(self.rag_pipelines_chain)} provider(s)"
        )

        # === v2.3/v2.4/v2.5: Interactive Analyst - Report Session Manager ===
        if REPORT_AGENTS_AVAILABLE and self.redis_client:
            try:
                # Initialize Researcher first (needed by ReportSessionManager)
                self.researcher = Researcher(
                    rag_module=self.qdrant_module,
                    web_module=self.web_search_module,
                    enrichment_module=None,  # Will be resolved lazily
                )

                # v2.5: Try to resolve ArtifactManager from DI container
                artifact_manager = None
                try:
                    artifact_manager = await self.di_container.resolve(
                        "system_artifact_manager"
                    )
                    logger.info("✅ ArtifactManager resolved from DI container")
                except Exception as am_err:
                    logger.warning(
                        f"ArtifactManager not available: {am_err} - Artifact export disabled"
                    )

                # v2.6: Try to resolve enrichment_pipeline for per-section enrichment
                enrichment_for_reports = None
                try:
                    enrichment_for_reports = await self.di_container.resolve(
                        "enrichment_pipeline"
                    )
                    logger.info("✅ EnrichmentPipeline resolved for report generation")
                except Exception as ep_err:
                    logger.warning(
                        f"EnrichmentPipeline not available for reports: {ep_err} - "
                        "Per-section enrichment disabled"
                    )

                # v6.0.1: Resolve correct LLM module for worker based on WORKER_PROVIDER
                worker_llm_module = None
                try:
                    import os
                    worker_provider = os.getenv("UBP_REPORT__WORKER_PROVIDER", "")
                    # v6.8.x: empty default → resolve via ProviderMapper (no grok hardcode)
                    if not worker_provider:
                        try:
                            chain = ProviderMapper.resolve_chain("enrichment")
                            if chain:
                                worker_provider = chain[0][1]
                        except Exception:
                            pass
                    if worker_provider in ProviderMapper.VALID_PROVIDERS:
                        worker_module_name = ProviderMapper.PROVIDER_MAP[worker_provider][0]
                        worker_llm_module = await self.di_container.resolve(worker_module_name)
                        logger.info(
                            f"[REPORT] Worker LLM resolved: provider='{worker_provider}' "
                            f"-> module='{worker_module_name}'"
                        )
                except Exception as wlm_err:
                    logger.warning(f"Could not resolve worker LLM module: {wlm_err} - using main LLM")

                # v6.0.1: Resolve correct LLM module for planner based on PLANNER_PROVIDER
                planner_llm_module = None
                try:
                    import os
                    planner_provider = os.getenv("UBP_REPORT__PLANNER_PROVIDER", "")
                    # v6.8.x: empty default → resolve via ProviderMapper (no grok hardcode)
                    if not planner_provider:
                        try:
                            chain = ProviderMapper.resolve_chain("enrichment")
                            if chain:
                                planner_provider = chain[0][1]
                        except Exception:
                            pass
                    if planner_provider in ProviderMapper.VALID_PROVIDERS:
                        planner_module_name = ProviderMapper.PROVIDER_MAP[planner_provider][0]
                        planner_llm_module = await self.di_container.resolve(planner_module_name)
                        logger.info(
                            f"[REPORT] Planner LLM resolved: provider='{planner_provider}' "
                            f"-> module='{planner_module_name}'"
                        )
                except Exception as plm_err:
                    logger.warning(f"Could not resolve planner LLM module: {plm_err} - using main LLM")

                # Initialize ReportSessionManager with all dependencies
                self.report_session_manager = ReportSessionManager(
                    redis_client=self.redis_client,
                    templates_path=None,  # Uses default path
                    llm_module=self.llm_module,  # v2.4: Enable DynamicPlanner
                    researcher=self.researcher,  # v2.4: Enable SwarmExecutor
                    artifact_manager=artifact_manager,  # v2.5: Enable artifact export
                    enrichment_module=enrichment_for_reports,  # v2.6: Enable per-section enrichment
                    worker_llm_module=worker_llm_module,  # v5.0.3 RPT-001: Separate worker LLM
                    planner_llm_module=planner_llm_module,  # v6.0.1: Separate planner LLM
                )
                await self.report_session_manager.initialize()

                logger.info(
                    f"✅ Interactive Analyst initialized: "
                    f"{len(self.report_session_manager._templates)} templates loaded"
                    + (
                        f", artifact export enabled"
                        if artifact_manager
                        else ", artifact export disabled"
                    )
                )
            except Exception as e:
                logger.warning(
                    f"Could not initialize Interactive Analyst: {e} - Report features disabled"
                )
                self.report_session_manager = None
                self.researcher = None
        else:
            logger.info(
                "Interactive Analyst not available (REPORT_AGENTS_AVAILABLE=False)"
            )

        logger.info(f"✅ {self.manifest.name} initialized successfully")

        # v2.0: Subscribe to memory.topic_shifted for context cleanup
        if hasattr(self, "subscriber") and self.subscriber:
            try:
                await self.subscriber.subscribe(
                    "memory.topic_shifted", self._on_topic_shifted
                )
                logger.info(
                    "✅ Subscribed to memory.topic_shifted for context management"
                )
            except Exception as e:
                logger.warning(f"Could not subscribe to memory.topic_shifted: {e}")

        if not self.capability_manager:
            self.capability_manager = CapabilityManager(adapter=self)

        # v3.5.0: Initialize dedicated Architect RAG pipeline
        await self._initialize_architect_pipeline()

        return {
            "status": "initialized",
            "module": self.manifest.name,
            "dependencies": {
                "rag_qdrant": self.qdrant_module is not None,
                "inference_ollama_grok": self.llm_module is not None,
                "redis": self.redis_client is not None,
                "rag_conversation_memory": self.memory_module is not None,
                "conversation_manager": self.conversation_manager is not None,
                "semantic_router": self.semantic_router is not None,
                "web_search": self.web_search_module is not None,
                "enrichment_pipeline": self.enrichment_module is not None,
                "report_session_manager": self.report_session_manager
                is not None,  # v2.3
                "researcher": self.researcher is not None,  # v2.3
                # v3.5.0: Architect dedicated pipeline
                "architect_pipeline": self.architect_pipeline is not None,
                "architect_llm_module": self.architect_llm_module is not None,
            },
        }

    async def _initialize_architect_pipeline(self) -> None:
        """
        Initialize dedicated Architect RAG pipeline with isolated provider (v3.5.0).
        
        The Architect pipeline BYPASSES standard provider resolution to ensure:
        1. Always uses the model specified in UBP_ARCHITECT__MODEL
        2. Uses dedicated enrichment provider from UBP_ARCHITECT__ENRICHMENT_*
        3. No interference with system-wide provider settings or overrides
        4. Separate LLM module instance for complete isolation
        
        This method resolves the Architect's dedicated providers and creates
        a separate RAG pipeline instance that the ask_architect method will use.
        """
        try:
            from ubp_enterprise_hybrid.backend.app.api.admin_settings_routes import settings_manager
            settings = settings_manager.get_settings()
            architect_settings = settings.architect
        except ImportError as e:
            logger.warning(
                f"Could not import settings_manager for Architect pipeline initialization: {e}",
                extra={"import_error": str(e)}
            )
            return

        if not architect_settings.enabled:
            logger.info("Architect Agent disabled, skipping dedicated pipeline initialization")
            return

        logger.info(
            "Initializing dedicated Architect RAG pipeline",
            extra={
                "provider": architect_settings.provider,
                "enrichment_provider": architect_settings.enrichment_provider,
            }
        )

        try:
            # v6.0.1: Provider-only resolution — model resolved by inference module
            from ubp_enterprise_hybrid.modules.cores._shared import ProviderMapper

            # v6.8.x: empty default → resolve via ProviderMapper (no grok hardcode)
            configured_provider = architect_settings.provider or ""
            if not configured_provider:
                try:
                    chain = ProviderMapper.resolve_chain("enrichment")
                    if chain:
                        configured_provider = chain[0][1]
                except Exception:
                    pass

            # v6.2.1 BUG-002: candidate list — configured first, then fallback chain
            candidate_providers = [configured_provider.lower()]
            for fb in settings.system.get_fallback_chain():
                fb_lower = fb.lower()
                if fb_lower not in candidate_providers and fb_lower in ProviderMapper.PROVIDER_MAP:
                    candidate_providers.append(fb_lower)

            logger.info(
                f"[ARCHITECT] Provider candidates: {candidate_providers} "
                f"(configured='{configured_provider}')"
            )

            architect_llm = None
            main_provider = None
            main_module_name = None

            for candidate in candidate_providers:
                if candidate not in ProviderMapper.PROVIDER_MAP:
                    logger.warning(f"[ARCHITECT] FALLBACK: '{candidate}' not in PROVIDER_MAP, skipping")
                    continue
                if not settings.providers.is_enabled(candidate):
                    logger.warning(f"[ARCHITECT] FALLBACK: '{candidate}' is DISABLED, skipping")
                    continue

                candidate_module = ProviderMapper.PROVIDER_MAP[candidate][0]
                resolved = await self.di_container.resolve(candidate_module)
                if not resolved:
                    logger.warning(
                        f"[ARCHITECT] FALLBACK: module '{candidate_module}' "
                        f"(provider '{candidate}') could not be resolved, trying next"
                    )
                    continue

                if hasattr(resolved, "set_default_provider"):
                    resolved.set_default_provider(candidate)

                architect_llm = resolved
                main_provider = candidate
                main_module_name = candidate_module

                if candidate != configured_provider.lower():
                    logger.warning(
                        f"[ARCHITECT] FALLBACK ACTIVATED: configured '{configured_provider}' "
                        f"unavailable, using '{candidate}' (module='{candidate_module}')"
                    )
                else:
                    logger.info(
                        f"[ARCHITECT] LLM module configured: {main_module_name} "
                        f"with provider '{candidate}'"
                    )
                break

            if not architect_llm:
                logger.error(
                    f"[ARCHITECT] FALLBACK EXHAUSTED: no provider resolved from "
                    f"{candidate_providers}. Architect will use standard RAG pipeline."
                )
                return

            self.architect_llm_module = architect_llm
            self._architect_provider = main_provider

            # Enrichment resolver (shared with system, uses UBP_ROLE_ENRICHMENT_PROVIDER)
            # Enrichment handles HyDE, Query Expansion, Reranking - separate from answer LLM
            # FIX-HYDE-001: Must attempt lazy resolution, not just return cached value
            async def architect_enrichment_resolver():
                """Resolver for Architect enrichment pipeline (shared with system).
                
                FIX-HYDE-001: Previously returned None if enrichment_module wasn't already
                cached, which caused HyDE and Investigative to be silently skipped.
                Now uses lazy resolution pattern like the main enrichment_resolver.
                """
                if self.enrichment_module:
                    return self.enrichment_module
                # FIX-HYDE-001: Attempt lazy resolution if not already cached
                if not self._enrichment_resolution_attempted:
                    self._enrichment_resolution_attempted = True
                    try:
                        self.enrichment_module = await self.di_container.resolve(
                            "enrichment_pipeline"
                        )
                        if self.enrichment_module:
                            logger.info(
                                "[ARCHITECT] enrichment_pipeline resolved lazily - "
                                "HyDE/Investigative now available"
                            )
                        return self.enrichment_module
                    except Exception as e:
                        logger.warning(f"[ARCHITECT] Could not resolve enrichment_pipeline: {e}")
                return None

            # Adaptive budget manager resolver (shared with system)
            async def architect_adaptive_memory_resolver():
                """Resolver for Architect adaptive budget manager (shared with system)."""
                if self.adaptive_memory_module:
                    return self.adaptive_memory_module
                if not self._adaptive_memory_resolution_attempted:
                    self._adaptive_memory_resolution_attempted = True
                    try:
                        self.adaptive_memory_module = await self.di_container.resolve(
                            "adaptive_budget_manager"
                        )
                        return self.adaptive_memory_module
                    except Exception as e:
                        logger.warning(f"Could not resolve adaptive_budget_manager for architect: {e}")
                return None

            # v6.8.0: Tool-calling LLM resolver for dual-LLM orchestration.
            # If the tool provider (e.g. vllm_remote) maps to a different inference module
            # than the main LLM (e.g. inference_ollama_grok), resolve it separately so
            # _generate_with_tools() can route tool calls to the correct module.
            tool_provider_name = None
            try:
                tool_search = settings.tool.search
                if tool_search.enabled and tool_search.enabled_architect and tool_search.provider_architect:
                    tool_provider_name = tool_search.provider_architect.lower()
            except Exception:
                pass

            tool_llm_resolver = None
            if (tool_provider_name
                    and tool_provider_name in ProviderMapper.PROVIDER_MAP
                    and ProviderMapper.PROVIDER_MAP[tool_provider_name][0] != main_module_name):
                # Tool provider maps to a DIFFERENT module than main LLM — need dual-LLM
                tool_module_name = ProviderMapper.PROVIDER_MAP[tool_provider_name][0]
                _di = self.di_container

                async def architect_tool_llm_resolver(_mod=tool_module_name, _prov=tool_provider_name):
                    """Resolve dedicated tool-calling LLM for Architect dual-LLM."""
                    try:
                        resolved = await _di.resolve(_mod)
                        if resolved:
                            if hasattr(resolved, "set_default_provider"):
                                resolved.set_default_provider(_prov)
                            logger.info(
                                f"[ARCHITECT] Tool LLM resolved: module={_mod}, "
                                f"provider={_prov} (dual-LLM with main={main_module_name})"
                            )
                        return resolved
                    except Exception as e:
                        logger.warning(f"[ARCHITECT] Could not resolve tool LLM {_mod}: {e}")
                        return None

                tool_llm_resolver = architect_tool_llm_resolver
                logger.info(
                    f"[ARCHITECT] Dual-LLM configured: main={main_module_name}/{main_provider}, "
                    f"tool={tool_module_name}/{tool_provider_name}"
                )

            # Create dedicated RAG pipeline for Architect
            # v3.6.0: architect_llm is now resolved from MAIN model provider (grok),
            # NOT from enrichment role (vllm). This ensures answer generation uses
            # the correct LLM with appropriate prompt length limits.
            self.architect_pipeline = RAGPipeline(
                self.qdrant_module,  # Same Qdrant (no isolation needed for vector store)
                architect_llm,  # Main LLM for answer generation (e.g. inference_ollama_grok)
                architect_enrichment_resolver,  # Enrichment pipeline (shared, uses vllm)
                None,  # reranker_resolver: architect uses enrichment reranking
                architect_adaptive_memory_resolver,  # v4.0.0: Dynamic token budget management
                tool_llm_resolver=tool_llm_resolver,  # v6.8.0: Dual-LLM for tool calls
            )

            logger.info(
                f"[ARCHITECT] Pipeline initialized: main_llm={main_module_name} "
                f"(provider={main_provider}), enrichment=shared",
                extra={
                    "pipeline_type": "isolated",
                    "main_module": main_module_name,
                    "main_provider": main_provider,
                }
            )

        except Exception as e:
            logger.error(
                f"Failed to initialize Architect dedicated pipeline: {e}",
                extra={"error_type": type(e).__name__},
            )
            # Fallback: Architect will use standard RAG pipeline
            logger.warning("Architect will use standard RAG pipeline as fallback")

    async def shutdown(self) -> None:
        """Shutdown RAG Orchestrator and cleanup."""
        logger.info(f"Shutting down {self.manifest.name} module")

        # Redis client is managed by DI container, just clear reference
        self.redis_client = None
        self.acl_manager = None
        self.config_manager = None
        self.conversation_manager = None
        self.rag_pipeline = None
        self.rag_pipelines_chain = []  # v1.10.1: HA Fallback
        self.llm_modules_chain = []  # v1.10.1: HA Fallback
        self.memory_module = None
        self.semantic_router = None  # FEAT-ROUTER-001
        self.web_search_module = None  # FEAT-ROUTER-001
        self.enrichment_module = None  # v1.7.1 GPU Reranking
        self.report_session_manager = None  # v2.3 Interactive Analyst
        self.researcher = None  # v2.3 Interactive Analyst
        self.capability_manager = None
        # v3.5.0: Architect dedicated pipeline cleanup
        self.architect_pipeline = None
        self.architect_llm_module = None

        logger.info(f"✅ {self.manifest.name} shutdown successfully")

    # =========================================================================
    # v2.0: STRUCTURED MEMORY - Topic Shift Handler
    # =========================================================================

    async def _on_topic_shifted(self, payload: Dict[str, Any]) -> None:
        """
        Handle topic shift events from rag_conversation_memory.

        When a topic shift is detected, we can:
        1. Log the event for analytics
        2. Optionally clear session-specific RAG caches (if any)
        3. Signal router to potentially adjust strategy

        Args:
            payload: Event payload with session_id, old_topic, new_topic
        """
        session_id = payload.get("session_id")
        old_topic = payload.get("old_topic", "unknown")
        new_topic = payload.get("new_topic", "unknown")

        logger.info(
            f"Topic shift detected in session {session_id}: "
            f"'{old_topic}' → '{new_topic}'"
        )

        # Future: Clear any session-specific RAG context cache
        # This is where we would invalidate cached chunks/context
        # for the session if we implement session-level caching

        # Future: Could signal router to adjust strategy for new topic
        # e.g., if moving from "technical" to "conversational" topic

    # =========================================================================
    # v1.10.1: High Availability - Pipeline with Fallback
    # =========================================================================

    async def _load_conversation_context(
        self,
        user_id: str,
        conversation_id: str,
        ctx=None,
        route_label: str = "RAG",
    ) -> str:
        """
        v4.1.0: Centralized conversation context loader.
        v4.3.0: Added in-process caching for parallel execution optimization.

        Strategy: summary-first with raw fallback.
        1. Check in-process cache (TTL 5min, LRU eviction)
        2. Try structured memory summary (compressed, ~500 chars)
        3. If unavailable, fallback to raw messages (max_turns=2)

        Args:
            user_id: User identifier
            conversation_id: Conversation/session identifier
            ctx: Security context for access verification
            route_label: Label for logging (RAG, CHAT, WEB, ARCHITECT)

        Returns:
            Formatted conversation context string (may be empty)
        """
        # v4.3.0: Check cache first
        cache_key = f"{user_id}:{conversation_id}"
        current_time = time.time()

        if cache_key in self._conversation_cache:
            cached_context, cached_time = self._conversation_cache[cache_key]
            if current_time - cached_time < self._conversation_cache_ttl:
                logger.debug(
                    f"[{route_label}] Cache hit for conversation context",
                    extra={"cache_key": cache_key, "context_len": len(cached_context)},
                )
                return cached_context
            else:
                # Expired entry, remove it
                del self._conversation_cache[cache_key]

        # 1. Try structured memory summary (replaces raw messages)
        if conversation_id and self.memory_module:
            try:
                memory_result = await self.memory_module.get_structured_context(
                    session_id=conversation_id, ctx=ctx
                )
                if memory_result.get("has_structured_context"):
                    structured_ctx = memory_result.get("system_message", "")
                    if structured_ctx:
                        logger.info(
                            f"[{route_label}] Using structured memory context ({len(structured_ctx)} chars)",
                            extra={
                                "conversation_id": conversation_id,
                                "context_source": "structured_memory",
                            },
                        )
                        # v4.3.0: Cache before returning
                        self._cache_conversation_context(cache_key, structured_ctx, current_time)
                        return structured_ctx
            except Exception as e:
                logger.warning(f"[{route_label}] Could not load structured memory: {e}")

        # 2. Fallback: raw messages (max_turns=2 to keep context small)
        if self.conversation_manager:
            try:
                context_result = await self.conversation_manager.get_context_for_llm(
                    user_id=user_id,
                    conversation_id=conversation_id,
                    max_turns=2,
                )
                if not context_result.get("error"):
                    raw_ctx = context_result.get("context", "")
                    if raw_ctx:
                        logger.info(
                            f"[{route_label}] Fallback to raw context (max_turns=2, {len(raw_ctx)} chars)",
                            extra={
                                "conversation_id": conversation_id,
                                "context_source": "raw_fallback",
                            },
                        )
                        # v4.3.0: Cache before returning
                        self._cache_conversation_context(cache_key, raw_ctx, current_time)
                        return raw_ctx
            except Exception as e:
                logger.warning(f"[{route_label}] Could not load raw conversation context: {e}")

        return ""

    def _cache_conversation_context(
        self, cache_key: str, context: str, timestamp: float
    ) -> None:
        """
        v4.3.0: Cache conversation context with LRU eviction.

        Args:
            cache_key: Cache key (user_id:conversation_id)
            context: Context string to cache
            timestamp: Current timestamp for TTL tracking
        """
        # LRU eviction: remove oldest entries if at capacity
        if len(self._conversation_cache) >= self._conversation_cache_max_entries:
            # Find and remove oldest entry
            oldest_key = None
            oldest_time = float("inf")
            for key, (_, cached_time) in self._conversation_cache.items():
                if cached_time < oldest_time:
                    oldest_time = cached_time
                    oldest_key = key
            if oldest_key:
                del self._conversation_cache[oldest_key]
                logger.debug(f"[CACHE] LRU eviction: removed {oldest_key}")

        # Store in cache
        self._conversation_cache[cache_key] = (context, timestamp)
        logger.debug(
            f"[CACHE] Stored conversation context",
            extra={"cache_key": cache_key, "context_len": len(context)},
        )

    def _invalidate_conversation_cache(
        self, user_id: str, conversation_id: str
    ) -> None:
        """
        Invalidate in-process conversation context cache after saving new messages.

        Must be called after add_message(role='assistant') so the next query
        reads fresh compressed context from Redis instead of stale in-process cache.
        """
        cache_key = f"{user_id}:{conversation_id}"
        if self._conversation_cache.pop(cache_key, None) is not None:
            logger.debug(f"[CACHE] Invalidated conversation context for {conversation_id}")

    async def _chat_with_fallback(
        self,
        query: str,
        collections: List[str],
        config: Dict[str, Any],
        return_debug: bool = False,
        web_context: Optional[str] = None,
        conversation_context: Optional[str] = None,
        pipeline_config: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """
        Execute RAG chat with automatic fallback on provider failure.

        Tries each provider in the chain sequentially. If primary fails,
        automatically retries with fallback provider(s).

        v1.10.2: Provider-Aware Context Limits
        Each provider has a different context window (Grok=128k, vLLM=8k).
        When falling back, the prompt is regenerated with the smaller context limit
        to prevent "Context Length Exceeded" errors on the fallback provider.

        Args:
            query: User question
            collections: List of collection IDs to query
            config: RAG configuration
            return_debug: Include debug information in response
            web_context: Optional web search results
            conversation_context: Optional conversation history
            pipeline_config: Permission-aware enrichment settings

        Returns:
            Chat result from the first successful provider

        Raises:
            Exception: If all providers fail
        """
        errors = []

        for idx, (pipeline, provider_name) in enumerate(self.rag_pipelines_chain):
            role_type = "Primary" if idx == 0 else f"Fallback-{idx}"

            try:
                # v1.10.2: Get provider-specific context window and pass to pipeline
                # The RAGPipeline._calculate_document_budget() will compute the exact
                # character budget considering query, system_prompt, web_context, etc.
                provider_config = config.copy()
                try:
                    # Get context window from the pipeline's LLM module
                    if hasattr(pipeline.llm, "get_context_window"):
                        context_window_tokens = pipeline.llm.get_context_window()
                        # Pass tokens directly - RAGPipeline will calculate budget
                        provider_config["context_limit_tokens"] = context_window_tokens
                        logger.info(
                            f"[HA] Provider '{provider_name}' context limit: {context_window_tokens} tokens",
                            extra={
                                "provider": provider_name,
                                "context_limit_tokens": context_window_tokens,
                            },
                        )
                except Exception as e:
                    logger.warning(
                        f"[HA] Could not get context window for {provider_name}: {e}"
                    )

                logger.info(
                    f"[HA] Trying {role_type} provider: {provider_name}",
                    extra={"provider": provider_name, "role": role_type},
                )

                result = await pipeline.chat(
                    query=query,
                    collections=collections,
                    config=provider_config,  # Use provider-specific config
                    return_debug=return_debug,
                    web_context=web_context,
                    conversation_context=conversation_context,
                    pipeline_config=pipeline_config,
                )

                # Success - add metadata about which provider was used
                result["_ha_provider"] = provider_name
                result["_ha_role"] = role_type
                result["_ha_attempts"] = idx + 1

                if idx > 0:
                    logger.warning(
                        f"[HA] Success with {role_type} provider '{provider_name}' "
                        f"after {idx} failed attempt(s)",
                        extra={
                            "provider": provider_name,
                            "attempts": idx + 1,
                            "errors": [str(e) for e in errors],
                        },
                    )
                else:
                    logger.debug(
                        f"[HA] Success with Primary provider '{provider_name}'"
                    )

                return result

            except Exception as e:
                error_msg = (
                    f"[{role_type}] {provider_name}: {type(e).__name__}: {str(e)}"
                )
                errors.append(error_msg)
                logger.warning(
                    f"[HA] {role_type} provider '{provider_name}' failed: {e}",
                    extra={
                        "provider": provider_name,
                        "role": role_type,
                        "error": str(e),
                        "error_type": type(e).__name__,
                    },
                )

                # If there are more providers, continue to next
                if idx < len(self.rag_pipelines_chain) - 1:
                    logger.info(
                        f"[HA] Attempting fallback to next provider in chain..."
                    )
                    continue

        # All providers failed
        error_summary = "; ".join(errors)
        logger.error(
            f"[HA] All providers failed ({len(errors)} attempts): {error_summary}",
            extra={"errors": errors, "attempts": len(errors)},
        )
        raise RuntimeError(
            f"All RAG providers failed ({len(errors)} attempts). Errors: {error_summary}"
        )

    async def health_check(self) -> Dict[str, Any]:
        """Perform health check on module and dependencies."""
        health = {
            "module": self.manifest.name,
            "status": "healthy",
            "dependencies": {
                "rag_qdrant": {"status": "unknown"},
                "inference_ollama_grok": {"status": "unknown"},
                "redis": {"status": "unknown"},
                "rag_conversation_memory": {"status": "not_configured"},
            },
        }

        # Check qdrant module
        if self.qdrant_module:
            try:
                qdrant_health = await self.qdrant_module.health_check()
                health["dependencies"]["rag_qdrant"] = qdrant_health
            except Exception as e:
                health["dependencies"]["rag_qdrant"] = {
                    "status": "unhealthy",
                    "error": str(e),
                }

        # Check LLM module
        if self.llm_module:
            try:
                llm_health = await self.llm_module.health_check()
                health["dependencies"]["inference_ollama_grok"] = llm_health
            except Exception as e:
                health["dependencies"]["inference_ollama_grok"] = {
                    "status": "unhealthy",
                    "error": str(e),
                }

        # Check Redis
        if self.redis_client:
            try:
                await self.redis_client.ping()
                health["dependencies"]["redis"] = {"status": "healthy"}
            except Exception as e:
                health["dependencies"]["redis"] = {
                    "status": "unhealthy",
                    "error": str(e),
                }

        # Check conversation memory module (optional - FEAT-MEM-001)
        if self.memory_module:
            try:
                memory_health = await self.memory_module.health_check()
                health["dependencies"]["rag_conversation_memory"] = memory_health
            except Exception as e:
                health["dependencies"]["rag_conversation_memory"] = {
                    "status": "unhealthy",
                    "error": str(e),
                }

        # Determine overall status (memory is optional, so exclude from degraded check)
        core_deps = ["rag_qdrant", "inference_ollama_grok", "redis"]
        dep_statuses = [
            health["dependencies"][dep].get("status", "unknown") for dep in core_deps
        ]
        if "unhealthy" in dep_statuses:
            health["status"] = "degraded"
        elif all(s == "healthy" for s in dep_statuses):
            health["status"] = "healthy"
        else:
            health["status"] = "degraded"

        return health

    # ===== Security Context Helpers =====

    def _build_legacy_ctx_compat_view(self, ctx: Any) -> Any:
        """Build a legacy-compatible ctx view from a flat OperationContext-like input."""
        user_id = getattr(ctx, "user_id", None)
        if not user_id:
            raise ValueError("Security context must contain user_id")

        client_id = getattr(ctx, "client_id", None)
        if not client_id:
            raise ValueError("Security context must contain client_id")

        roles = getattr(ctx, "roles", [])
        if not isinstance(roles, (list, set, tuple)):
            roles = []

        user = SimpleNamespace(
            user_id=str(user_id),
            client_id=str(client_id),
            roles=list(roles),
        )
        return SimpleNamespace(
            user=user,
            user_id=user.user_id,
            client_id=user.client_id,
            roles=list(roles),
            source=getattr(ctx, "source", None),
            session_id=getattr(ctx, "session_id", None),
        )

    def _require_ctx(self, ctx: Any) -> Any:
        """
        Validate and return security context.

        Args:
            ctx: Security context to validate

        Returns:
            The validated ctx (for chaining)

        Raises:
            ValueError: If ctx is None or missing user info
        """
        if ctx and hasattr(ctx, "user") and ctx.user:
            if not hasattr(ctx.user, "user_id"):
                raise ValueError("Security context must contain user_id")
            return ctx

        if is_operation_context_like(ctx):
            # Preserve ctx.user.* for legacy ACL/retrieve callers while accepting
            # the authoritative flat OperationContext contract from MCP.
            return self._build_legacy_ctx_compat_view(ctx)

        raise ValueError("Security context required for this operation")

    def _is_admin(self, ctx: Any) -> bool:
        """
        Check if the current user is an administrator.

        Args:
            ctx: Security context with user info

        Returns:
            True if user has admin role, False otherwise
        """
        if not ctx or not hasattr(ctx, "user") or not ctx.user:
            return False
        if not hasattr(ctx.user, "roles"):
            return False

        roles = ctx.user.roles
        # Ensure roles is iterable
        if not isinstance(roles, (list, set, tuple)):
            return False

        return "admin" in roles

    # ===== MCP-COMPAT: OperationContext helpers (ARCH-008) =====

    def _build_context_from_di(self) -> OperationContext:
        """
        Build OperationContext from DI container — backward compatibility for REST path.
        
        MCP-COMPAT: When ctx is not provided (REST path), this method constructs
        an OperationContext from the DI container state.
        
        Returns:
            OperationContext with default values
        """
        # In rag_orchestrator, security context is typically passed via ctx parameter.
        # This provides a minimal fallback for internal calls.
        return OperationContext(
            client_id="default",
            user_id=None,
            session_id=None,
            source="rest",
        )

    def _normalize_ctx(self, ctx: Any) -> OperationContext:
        """
        Normalize any context format to OperationContext.
        
        MCP-COMPAT: Handles both legacy security context (ctx.user.user_id) 
        and new OperationContext format for backward compatibility.
        
        Args:
            ctx: Either OperationContext, legacy security context, or None
            
        Returns:
            OperationContext instance
        """
        if ctx is None:
            return self._build_context_from_di()
        
        # Already an OperationContext
        if isinstance(ctx, OperationContext):
            return ctx
        
        # Legacy security context format (ctx.user.user_id, ctx.user.roles)
        if hasattr(ctx, "user") and ctx.user:
            user_id = getattr(ctx.user, "user_id", None)
            roles = getattr(ctx.user, "roles", [])
            client_id = getattr(ctx.user, "client_id", "default")
            if not isinstance(roles, (list, tuple)):
                roles = []
            return OperationContext(
                client_id=str(client_id) if client_id else "default",
                user_id=str(user_id) if user_id else None,
                roles=list(roles),
                source="rest",
            )
        
        # Fallback
        return self._build_context_from_di()

    def _is_report_approval_query(self, query: str) -> bool:
        """
        v2.6 FIX: Check if query is a report approval command.

        This prevents approval queries from being sent to the enrichment
        pipeline when the conversation_id doesn't match the report session.

        Uses word boundary matching to avoid false positives like:
        - "sicurezza" should NOT match "si"
        - "appropriato" should NOT match "approva"

        Args:
            query: User input query

        Returns:
            True if query matches an approval pattern
        """
        import re

        query_lower = query.lower().strip()

        # Approval patterns (Italian and English)
        approval_patterns = [
            # Italian
            r"\bsì\b", r"\bsi\b", r"\bok\b", r"\bprocedi\b", r"\bvai\b",
            r"\bconferma\b", r"\bapprova\b", r"\bapprovato\b", r"\bva bene\b",
            r"\besegui\b", r"\binizia\b", r"\bavvia\b", r"\bparti\b",
            r"\bfallo\b", r"\bperfetto\b", r"\bottimo\b",
            # English
            r"\byes\b", r"\bproceed\b", r"\bgo\b", r"\bconfirm\b",
            r"\bstart\b", r"\bexecute\b", r"\bdo it\b", r"\bapproved\b",
            r"\bapprove\b", r"\bperfect\b", r"\bgreat\b", r"\bfine\b",
        ]

        for pattern in approval_patterns:
            if re.search(pattern, query_lower):
                return True

        return False

    def _format_web_results(self, results: List[Dict[str, Any]]) -> str:
        """
        Format web search results into context string for LLM.

        FEAT-ROUTER-001 - Web routing support.

        Args:
            results: List of web search results from web_search module

        Returns:
            Formatted string with web results for LLM context
        """
        if not results:
            return ""

        formatted_parts = ["=== Web Search Results ===\n"]

        for i, result in enumerate(results[:5], 1):  # Limit to 5 results
            title = result.get("title", "No title")
            url = result.get("href", result.get("url", ""))
            snippet = result.get(
                "body", result.get("snippet", result.get("description", ""))
            )

            formatted_parts.append(f"[{i}] {title}")
            if url:
                formatted_parts.append(f"    URL: {url}")
            if snippet:
                # Truncate long snippets
                snippet_clean = snippet[:500] + "..." if len(snippet) > 500 else snippet
                formatted_parts.append(f"    {snippet_clean}")
            formatted_parts.append("")  # Empty line between results

        formatted_parts.append("=== End Web Results ===")
        return "\n".join(formatted_parts)

    def _is_report_request(self, query: str) -> bool:
        """
        Detect if query is requesting a structured report.

        v2.3: Interactive Analyst - Report pattern detection.

        Patterns detected:
        - "Fammi un report su..."
        - "Genera un'analisi di..."
        - "Crea un documento su..."
        - "Audit di sicurezza..."
        - "Confronto tra X e Y..."

        Args:
            query: User query text

        Returns:
            True if query matches report request patterns
        """
        if not query:
            return False

        query_lower = query.lower().strip()

        # Report request patterns (Italian and English)
        report_patterns = [
            # Italian patterns
            r"^(fammi|fai|genera|crea|scrivi)\s+(un\s+)?(report|analisi|documento|sintesi|riassunto)",
            r"^(voglio|vorrei|mi serve)\s+(un\s+)?(report|analisi|documento)",
            r"^report\s+(su|di|per|sulla|del|della)",
            r"^analisi\s+(su|di|per|tecnica|approfondita)",
            r"^audit\s+(di|della|del)\s+sicurezza",
            r"^confronto\s+(tra|fra|di)",
            # English patterns
            r"^(generate|create|write|make)\s+(a\s+)?(report|analysis|document|summary)",
            r"^(i\s+want|i\s+need)\s+(a\s+)?(report|analysis)",
            r"^report\s+(on|about|for)",
            r"^(security|technical)\s+audit",
            r"^comparison\s+(between|of)",
        ]

        for pattern in report_patterns:
            if re.search(pattern, query_lower, re.IGNORECASE):
                return True

        return False

    async def _handle_web_fallback_to_chat(
        self,
        query: str,
        user_id: str,
        conversation_id: str,
        route_result: Optional[RouterResult],
        fallback_reason: str,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        Handle fallback from WEB route to CHAT mode when web search is unavailable.

        FEAT-ROUTER-001 - Graceful degradation for WEB routing.

        Args:
            query: Original user query
            user_id: User ID for conversation tracking
            conversation_id: Conversation ID
            route_result: Router result for metadata
            fallback_reason: Reason for fallback (for logging/debugging)

        Returns:
            Response dict with CHAT mode answer and fallback notification
        """
        logger.info(
            f"WEB->CHAT fallback: {fallback_reason}",
            extra={
                "user_id": user_id,
                "conversation_id": conversation_id,
                "fallback_reason": fallback_reason,
            },
        )

        # Load conversation context for continuity
        conversation_context = ""
        if self.conversation_manager:
            try:
                context_result = await self.conversation_manager.get_context_for_llm(
                    user_id=user_id,
                    conversation_id=conversation_id,
                    max_turns=self.config.get("default_rag_config", {}).get(
                        "context_max_turns", 10
                    ),
                )
                if not context_result.get("error"):
                    conversation_context = context_result.get("context", "")
            except Exception as e:
                logger.warning(f"Could not load conversation context: {e}")

        # Generate response with LLM only (no web context)
        logger.info("[RAG] op=web_fallback_chat mode=pure_llm")
        chat_result = await self.rag_pipeline.chat(
            query=query,
            collections=[],  # Empty = Pure LLM mode
            config=self.config.get("default_rag_config", {}),
            return_debug=True,
            conversation_context=conversation_context,
        )

        # Save to conversation history
        if self.conversation_manager:
            try:
                await self.conversation_manager.add_message(
                    user_id=user_id,
                    conversation_id=conversation_id,
                    role="user",
                    content=query,
                    metadata={"route": "web", "fallback": "chat"},
                )
                await self.conversation_manager.add_message(
                    user_id=user_id,
                    conversation_id=conversation_id,
                    role="assistant",
                    content=chat_result.get("answer", ""),
                    metadata={
                        "route": "chat",
                        "fallback_from": "web",
                        "fallback_reason": fallback_reason,
                    },
                )
            except Exception as e:
                logger.warning(f"Could not save chat messages: {e}")

        # v4.1.0: Save to structured memory for eager compression
        if conversation_id and self.memory_module:
            try:
                await self.memory_module.add_message(
                    session_id=conversation_id, role="user",
                    content=query, ctx=ctx,
                )
                await self.memory_module.add_message(
                    session_id=conversation_id, role="assistant",
                    content=chat_result.get("answer", ""),
                    metadata={"route": "chat_fallback"}, ctx=ctx,
                )
                self._invalidate_conversation_cache(user_id, conversation_id)
            except Exception as e:
                logger.warning(f"Could not save to structured memory: {e}")

        return {
            "answer": chat_result.get("answer", ""),
            "sources": [],
            "config_used": {},
            "permissions_checked": [],
            "conversation_id": conversation_id,
            "mode": "chat",
            "mode_reason": "web_fallback",
            "fallback_reason": fallback_reason,
            "router": route_result.to_dict() if route_result else None,
        }

    # =========================================================================
    # SMART AUTO-RETRY SYSTEM (FEAT-ROUTER-003)
    # =========================================================================

    # Phrases that indicate empty/uncertain RAG responses (Italian + English)
    EMPTY_RESPONSE_PATTERNS = [
        # Italian
        r"non ho informazioni",
        r"non ho trovato",
        r"non sono in grado",
        r"non ho dati",
        r"non dispongo di",
        r"non posso rispondere",
        r"non ho abbastanza informazioni",
        r"non ho conoscenze",
        r"non ho accesso a",
        r"non trovo informazioni",
        r"non ci sono documenti",
        r"nessun documento",
        r"nessuna informazione",
        r"mi dispiace.*non",
        # English
        r"i don'?t have information",
        r"i couldn'?t find",
        r"i'?m not able",
        r"i don'?t have data",
        r"i don'?t have enough",
        r"i cannot answer",
        r"no information available",
        r"no documents found",
        r"i'?m sorry.*not",
        r"i apologize.*cannot",
        # Generic
        r"^(?:mi dispiace|sorry|scusa)\.?$",
    ]

    # Compiled patterns for efficiency
    _empty_response_compiled: Optional[List[re.Pattern]] = None

    def _get_empty_patterns(self) -> List[re.Pattern]:
        """Get compiled empty response patterns (lazy initialization)."""
        if self._empty_response_compiled is None:
            self._empty_response_compiled = [
                re.compile(p, re.IGNORECASE) for p in self.EMPTY_RESPONSE_PATTERNS
            ]
        return self._empty_response_compiled

    def _is_empty_response(self, answer: str, sources_count: int = 0) -> bool:
        """
        Detect if a RAG response indicates lack of information.

        This method identifies "I don't know" type responses that should
        trigger auto-retry with a different route.

        Args:
            answer: The LLM response text
            sources_count: Number of sources found (0 is a strong signal)

        Returns:
            True if response indicates empty/uncertain result
        """
        if not answer:
            return True

        # Strong signal: RAG route but 0 sources found
        if sources_count == 0:
            # Check if answer seems like a deflection
            answer_lower = answer.lower().strip()
            if len(answer_lower) < 200:  # Short answers more likely to be deflections
                for pattern in self._get_empty_patterns():
                    if pattern.search(answer_lower):
                        logger.debug(
                            f"Empty response detected via pattern: {pattern.pattern}"
                        )
                        return True

        # Very short answers are suspicious for information queries
        if len(answer.strip()) < 50 and sources_count == 0:
            return True

        return False

    async def _llm_general_fallback(
        self,
        query: str,
        user_id: str,
        conversation_id: str,
        conversation_context: str = "",
    ) -> Dict[str, Any]:
        """
        Ultimate fallback: Use LLM general knowledge when all routes fail.

        This method is the last resort when:
        - RAG returned empty (no documents)
        - WEB search failed or unavailable

        Args:
            query: User query
            user_id: User ID
            conversation_id: Conversation ID
            conversation_context: Previous conversation for continuity

        Returns:
            Response dict with LLM general knowledge answer
        """
        logger.info(
            "LLM general fallback triggered",
            extra={
                "user_id": user_id,
                "conversation_id": conversation_id,
            },
        )

        # Generate response with LLM only, with explicit instruction
        fallback_prompt = f"""L'utente ha chiesto: "{query}"

Non ho trovato informazioni specifiche nei documenti interni né tramite ricerca web.
Rispondi usando le tue conoscenze generali, ma specifica chiaramente che questa è una risposta basata su conoscenze generali e non su documenti specifici dell'azienda.

Se non sei sicuro, suggerisci all'utente di verificare con fonti ufficiali o di riformulare la domanda."""

        try:
            if self.rag_pipeline:
                logger.info("[RAG] op=llm_general_fallback mode=pure_llm")
                result = await self.rag_pipeline.chat(
                    query=fallback_prompt,
                    collections=[],  # Pure LLM mode
                    config=self.config.get("default_rag_config", {}),
                    return_debug=True,
                    conversation_context=conversation_context,
                )
                answer = result.get("answer", "")
            else:
                answer = "Mi dispiace, non sono riuscito a trovare informazioni pertinenti alla tua richiesta."

        except Exception as e:
            logger.error(f"LLM general fallback error: {e}")
            answer = "Mi dispiace, si è verificato un errore durante l'elaborazione della richiesta."

        return {
            "answer": answer,
            "sources": [],
            "config_used": {},
            "permissions_checked": [],
            "conversation_id": conversation_id,
            "mode": "llm_general",
            "mode_reason": "all_routes_failed",
        }

    async def _auto_retry_with_fallback(
        self,
        query: str,
        original_result: Dict[str, Any],
        original_route: RouteType,
        route_result: RouterResult,
        user_id: str,
        client_id: Optional[str],
        conversation_id: str,
        rag_config: Dict[str, Any],
        conversation_context: str = "",
        pipeline_config: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """
        Attempt auto-retry with fallback route after empty primary result.

        Smart Auto-Retry Flow:
        1. RAG returned empty → try WEB
        2. WEB returned empty → try RAG (rare)
        3. Both failed → LLM general fallback

        Args:
            query: Original query
            original_result: Result from primary route (contains empty answer)
            original_route: The route that failed (RAG/WEB)
            route_result: Router classification result
            user_id: User ID
            client_id: Client ID (optional)
            conversation_id: Conversation ID
            rag_config: RAG configuration
            conversation_context: Conversation context for continuity
            pipeline_config: Permission-aware enrichment settings

        Returns:
            Response dict from fallback route or LLM general
        """
        if not self.semantic_router:
            return original_result

        # Check if empty response
        sources_count = len(original_result.get("sources", []))
        if not self._is_empty_response(
            original_result.get("answer", ""), sources_count
        ):
            return original_result  # Not empty, keep original

        # Get smart fallback
        fallback_result = self.semantic_router.get_smart_fallback(
            route_result, rag_sources_count=sources_count
        )

        if not fallback_result:
            return original_result  # No fallback available

        logger.info(
            f"Auto-retry: {original_route.value} → {fallback_result.route.value}",
            extra={
                "user_id": user_id,
                "conversation_id": conversation_id,
                "original_route": original_route.value,
                "fallback_route": fallback_result.route.value,
            },
        )

        # Execute fallback route
        try:
            if fallback_result.route == RouteType.WEB:
                # Try WEB search
                if self.web_search_module:
                    web_result = await self.web_search_module.search(
                        query=query,
                        max_results=self.config.get("router", {}).get(
                            "web_max_results", 5
                        ),
                    )
                    web_results_list = web_result.get("results", [])

                    if web_results_list:
                        web_context = self._format_web_results(web_results_list)

                        # Generate response with web context
                        if self.rag_pipeline:
                            logger.info("[RAG] op=web_enriched_chat mode=web_context")
                            web_chat_result = await self.rag_pipeline.chat(
                                query=query,
                                collections=[],
                                config=rag_config,
                                return_debug=True,
                                web_context=web_context,
                                conversation_context=conversation_context,
                                pipeline_config=pipeline_config,
                            )

                            return {
                                "answer": web_chat_result.get("answer", ""),
                                "sources": [
                                    {
                                        "type": "web",
                                        "title": r.get("title", ""),
                                        "url": r.get("href", r.get("url", "")),
                                    }
                                    for r in web_results_list[:5]
                                ],
                                "config_used": rag_config,
                                "permissions_checked": [],
                                "conversation_id": conversation_id,
                                "mode": "web",
                                "mode_reason": "auto_retry_from_rag",
                                "router": fallback_result.to_dict(),
                            }

            elif fallback_result.route == RouteType.RAG:
                # Try RAG (rare case: WEB failed)
                # Get accessible collections
                if self.acl_manager:
                    accessible = await self.acl_manager.get_accessible_collections(
                        user_id, client_id
                    )
                    if accessible and self.rag_pipeline:
                        logger.info("[RAG] op=rag_retry mode=vector collections=%d", len(accessible))
                        rag_retry_result = await self.rag_pipeline.chat(
                            query=query,
                            collections=accessible,
                            config=rag_config,
                            return_debug=True,
                            conversation_context=conversation_context,
                            pipeline_config=pipeline_config,
                        )

                        if rag_retry_result.get("sources"):
                            return {
                                **rag_retry_result,
                                "conversation_id": conversation_id,
                                "mode": "rag",
                                "mode_reason": "auto_retry_from_web",
                                "router": fallback_result.to_dict(),
                            }

        except Exception as e:
            logger.warning(f"Auto-retry fallback failed: {e}")

        # All fallbacks failed - use LLM general knowledge
        return await self._llm_general_fallback(
            query=query,
            user_id=user_id,
            conversation_id=conversation_id,
            conversation_context=conversation_context,
        )

    def _require_admin(self, ctx: Any, operation: str) -> Any:
        """
        Require admin privileges for an operation.

        Args:
            ctx: Security context
            operation: Name of the operation for error logging

        Returns:
            The validated ctx

        Raises:
            PermissionError: If user is not admin
            ValueError: If ctx is invalid
        """
        ctx = self._require_ctx(ctx)
        if not self._is_admin(ctx):
            logger.warning(
                f"Unauthorized {operation} attempt by user {ctx.user.user_id}",
                extra={"user_id": ctx.user.user_id, "operation": operation},
            )
            raise PermissionError(f"Only administrators can perform: {operation}")
        return ctx

    async def _check_client_kb_creation_permission(self, ctx: Any) -> Dict[str, Any]:
        """
        Check if the caller's client is authorized to create knowledge bases.

        Enterprise v2.0 - GAP-002 Implementation.

        This method checks the client's kb_config.can_create_universal_kb setting.
        If True, the client is allowed to create KBs with specific naming conventions.

        Args:
            ctx: Security context with user info (must have client_id)

        Returns:
            Dict with:
                - allowed: bool - Whether creation is permitted
                - client_id: str - The client ID
                - max_kbs: int - Maximum KBs this client can create
                - current_count: int - Current KB count for this client
                - kb_prefix: str - Required prefix for KB names
                - error: str (optional) - Error message if not allowed
        """
        client_id = getattr(ctx.user, "client_id", None)

        # Platform admin (no client_id) should use _require_admin instead
        if is_platform_admin_client(client_id):
            return {
                "allowed": False,
                "client_id": PLATFORM_ADMIN_CLIENT_ID,
                "error": "Platform admins should use admin endpoints for KB creation",
            }

        # Resolve admin_clients module to fetch client config
        try:
            admin_clients = await self.di_container.resolve("admin_clients")
            if not admin_clients:
                logger.error("[GAP-002] admin_clients module not available")
                return {
                    "allowed": False,
                    "client_id": client_id,
                    "error": "Client management service unavailable",
                }
        except Exception as e:
            logger.error(f"[GAP-002] Failed to resolve admin_clients: {e}")
            return {
                "allowed": False,
                "client_id": client_id,
                "error": f"Failed to verify client permissions: {e}",
            }

        # Fetch client configuration
        try:
            # Use internal get to avoid permission loops
            client_key = f"ubp:admin:client:{client_id}"
            client_data = await self.redis_client.get(client_key)

            if not client_data:
                logger.warning(f"[GAP-002] Client not found: {client_id}")
                return {
                    "allowed": False,
                    "client_id": client_id,
                    "error": ErrorMessages.CLIENT_NOT_FOUND.format(client_id=client_id),
                }

            client = json.loads(client_data)
            kb_config = client.get("kb_config", {})

            # Check can_create_universal_kb flag
            can_create = kb_config.get("can_create_universal_kb", False)
            max_kbs = kb_config.get("max_universal_kbs", 5)

            if not can_create:
                logger.info(
                    f"[GAP-002] Client {client_id} not authorized to create KBs",
                    extra={"client_id": client_id, "can_create": False},
                )
                return {
                    "allowed": False,
                    "client_id": client_id,
                    "error": ErrorMessages.CLIENT_KB_CREATE_DENIED,
                }

            # Count existing KBs created by this client
            # Convention: client KBs are prefixed with client_{client_id[:8]}_
            kb_prefix = f"{CLIENT_KB_PREFIX}{client_id[:8]}_"
            current_count = 0

            # Scan Redis for KB metadata with this client's prefix
            async for key in self.redis_client.scan_iter(
                match=f"rag:kb:{kb_prefix}*:metadata"
            ):
                current_count += 1

            if current_count >= max_kbs:
                logger.warning(
                    f"[GAP-002] Client {client_id} has reached max KB limit ({max_kbs})",
                    extra={
                        "client_id": client_id,
                        "current_count": current_count,
                        "max_kbs": max_kbs,
                    },
                )
                return {
                    "allowed": False,
                    "client_id": client_id,
                    "current_count": current_count,
                    "max_kbs": max_kbs,
                    "error": f"Maximum KB limit reached ({current_count}/{max_kbs})",
                }

            logger.info(
                f"[GAP-002] Client {client_id} authorized to create KB ({current_count}/{max_kbs})",
                extra={
                    "client_id": client_id,
                    "can_create": True,
                    "current_count": current_count,
                    "max_kbs": max_kbs,
                },
            )

            return {
                "allowed": True,
                "client_id": client_id,
                "max_kbs": max_kbs,
                "current_count": current_count,
                "kb_prefix": kb_prefix,
            }

        except json.JSONDecodeError as e:
            logger.error(f"[GAP-002] Invalid client data for {client_id}: {e}")
            return {
                "allowed": False,
                "client_id": client_id,
                "error": "Invalid client configuration",
            }
        except Exception as e:
            logger.error(f"[GAP-002] Error checking client KB permissions: {e}")
            return {
                "allowed": False,
                "client_id": client_id,
                "error": f"Error checking permissions: {e}",
            }

    def _require_initialized(self) -> None:
        """
        Ensure module dependencies are initialized.

        Raises:
            RuntimeError: If required dependencies are not initialized
        """
        if not self.qdrant_module:
            raise RuntimeError("qdrant_module not initialized")
        if not self.redis_client:
            raise RuntimeError("redis_client not initialized")
        if not self.acl_manager:
            raise RuntimeError("acl_manager not initialized")
        if not self.config_manager:
            raise RuntimeError("config_manager not initialized")
        if not self.rag_pipeline:
            raise RuntimeError("rag_pipeline not initialized")

    async def get_collection_metadata_internal(
        self, collection_id: str
    ) -> Optional[Dict[str, Any]]:
        """Internal helper to fetch KB metadata without security checks."""
        if not collection_id:
            return None
        self._require_initialized()

        metadata_key = f"rag:kb:{collection_id}:metadata"
        try:
            raw = await self.redis_client.get(metadata_key)
            if not raw:
                return None
            if isinstance(raw, bytes):
                raw = raw.decode("utf-8")
            return json.loads(raw)
        except json.JSONDecodeError as exc:
            logger.warning(
                "Malformed KB metadata",
                extra={"collection_id": collection_id, "error": str(exc)},
            )
            return None
        except Exception as exc:  # pragma: no cover - defensive logging
            logger.debug(
                "Unable to read KB metadata",
                extra={"collection_id": collection_id, "error": str(exc)},
            )
            return None

    async def update_collection_metadata_internal(
        self, collection_id: str, metadata: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Internal helper to upsert KB metadata in Redis."""
        if not collection_id:
            raise ValueError("collection_id is required")
        if metadata is None:
            metadata = {}
        self._require_initialized()

        metadata_key = f"rag:kb:{collection_id}:metadata"
        try:
            existing = await self.get_collection_metadata_internal(collection_id) or {}
            existing.update(metadata)
            await self.redis_client.set(metadata_key, json.dumps(existing))
            return {
                "status": "success",
                "collection_id": collection_id,
                "metadata": existing,
            }
        except Exception as exc:
            logger.error(
                "Failed to update KB metadata",
                extra={"collection_id": collection_id, "error": str(exc)},
            )
            return {
                "status": "error",
                "collection_id": collection_id,
                "error": str(exc),
            }

    # ===== Knowledge Base Management =====

    async def create_knowledge_base(
        self,
        name: str,
        description: str = "",
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Create a new knowledge base (Qdrant collection).

        Enterprise v2.0 - GAP-002: Supports both admin and authorized client creation.

        Security Model:
        - Platform Admin (role=admin): Can create any KB, no naming restrictions
        - Authorized Client (can_create_universal_kb=true): Can create KBs with
          prefix client_{client_id[:8]}_, limited by max_universal_kbs

        Args:
            name: Collection name (clients must use their assigned prefix)
            description: Optional description
            ctx: Security context with user info
            **kwargs: Additional config (embedding_model, kb_type, etc.)

        Returns:
            Status with collection_id, or error details
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        # Determine caller type and authorization
        is_admin = self._is_admin(ctx)
        client_id = getattr(ctx.user, "client_id", None)
        user_id = ctx.user.user_id

        # Track authorization source for audit
        auth_source = "unknown"
        owner_client_id = None  # The client that owns this KB (None for admin-created)

        if is_admin:
            # Platform admin - full access, no restrictions
            auth_source = "platform_admin"
            logger.info(
                f"[GAP-002] Admin KB creation: {name}",
                extra={
                    "user_id": user_id,
                    "kb_name": name,
                    "auth_source": auth_source,
                },
            )
        elif client_id and not is_platform_admin_client(client_id):
            # Client user - check can_create_universal_kb permission
            permission_check = await self._check_client_kb_creation_permission(ctx)

            if not permission_check.get("allowed"):
                error_msg = permission_check.get(
                    "error", ErrorMessages.CLIENT_KB_CREATE_DENIED
                )
                logger.warning(
                    f"[GAP-002] Client KB creation denied: {error_msg}",
                    extra={
                        "client_id": client_id,
                        "user_id": user_id,
                        "kb_name": name,
                    },
                )
                return {
                    "collection_id": "",
                    "status": "error",
                    "message": error_msg,
                    "error_code": "KB_CREATE_DENIED",
                }

            # Validate KB name follows client naming convention
            required_prefix = permission_check.get(
                "kb_prefix", f"{CLIENT_KB_PREFIX}{client_id[:8]}_"
            )
            if not name.startswith(required_prefix):
                error_msg = (
                    f"Client KBs must be prefixed with '{required_prefix}'. "
                    f"Example: {required_prefix}my_knowledge_base"
                )
                logger.warning(
                    f"[GAP-002] Invalid KB name for client: {name}",
                    extra={
                        "client_id": client_id,
                        "required_prefix": required_prefix,
                        "provided_name": name,
                    },
                )
                return {
                    "collection_id": "",
                    "status": "error",
                    "message": error_msg,
                    "error_code": "INVALID_KB_NAME",
                    "required_prefix": required_prefix,
                }

            auth_source = "client_authorized"
            owner_client_id = client_id
            logger.info(
                f"[GAP-002] Client KB creation authorized: {name}",
                extra={
                    "client_id": client_id,
                    "user_id": user_id,
                    "kb_name": name,
                    "auth_source": auth_source,
                    "kb_count": f"{permission_check.get('current_count', 0)}/{permission_check.get('max_kbs', 5)}",
                },
            )
        else:
            # No valid authorization
            logger.warning(
                f"[GAP-002] Unauthorized KB creation attempt",
                extra={
                    "user_id": user_id,
                    "client_id": client_id,
                    "kb_name": name,
                },
            )
            raise PermissionError(ErrorMessages.ADMIN_REQUIRED)

        # Validate name
        if not name:
            return {
                "collection_id": "",
                "status": "error",
                "message": "Collection name is required",
                "error_code": "MISSING_NAME",
            }

        # Validate name format (alphanumeric, underscores, hyphens)
        if not re.match(r"^[a-zA-Z][a-zA-Z0-9_-]*$", name):
            return {
                "collection_id": "",
                "status": "error",
                "message": "KB name must start with a letter and contain only alphanumeric characters, underscores, and hyphens",
                "error_code": "INVALID_NAME_FORMAT",
            }

        try:
            # Extract config from kwargs
            # FIX-DIM-v4.1.2: Do NOT use hardcoded default for embedding_model
            # If not explicitly provided, pass None to let rag_qdrant read from ENV/config
            embedding_model = kwargs.get("embedding_model")  # None → rag_qdrant uses config
            vector_size = kwargs.get(
                "vector_size"
            )  # Let backend derive from model if not set
            distance = kwargs.get("distance", "Cosine")
            chunk_size = kwargs.get("chunk_size", 512)
            chunk_overlap = kwargs.get("chunk_overlap", 50)

            # Determine KB type based on creator
            if owner_client_id:
                kb_type = "client"  # Client-owned KB
            elif is_personal_kb(name):
                kb_type = "personal"
            else:
                kb_type = kwargs.get("kb_type", "universal")

            # Create collection via qdrant module
            result = await self.qdrant_module.create_collection_internal(
                collection_name=name,
                vector_size=vector_size,
                distance=distance,
                description=description,
                embedding_model=embedding_model,
                kb_type=kb_type,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
            )

            if result.get("status") == "created" or result.get("status") == "exists":
                logger.info(f"[GAP-002] Knowledge base '{name}' created successfully")

                # Store full config in Redis metadata (including ownership)
                if self.redis_client:
                    metadata_key = f"rag:kb:{name}:metadata"
                    kb_metadata = {
                        "name": name,
                        "description": description,
                        "embedding_model": embedding_model,
                        "kb_type": kb_type,
                        "chunk_size": chunk_size,
                        "chunk_overlap": chunk_overlap,
                        "created_by_user_id": user_id,
                        "created_by_client_id": owner_client_id
                        or PLATFORM_ADMIN_CLIENT_ID,
                        "auth_source": auth_source,
                        "created_at": json.dumps({}),  # Will be set by Redis
                    }
                    await self.redis_client.set(metadata_key, json.dumps(kb_metadata))

                # AUTO-GRANT: Give creator 'write' permission
                if self.acl_manager:
                    try:
                        # Grant to creator user
                        await self.acl_manager.set_permission(
                            entity_type=EntityType.USER.value,
                            entity_id=user_id,
                            collection_id=name,
                            access_level=AccessLevel.WRITE.value,
                        )
                        logger.info(
                            f"[GAP-002] ACL: Granted 'write' to creator {user_id} on KB '{name}'"
                        )

                        # If client-created, also grant to the client entity
                        # This allows all users of the client to access it
                        if owner_client_id:
                            await self.acl_manager.set_permission(
                                entity_type=EntityType.CLIENT.value,
                                entity_id=owner_client_id,
                                collection_id=name,
                                access_level=AccessLevel.WRITE.value,
                            )
                            logger.info(
                                f"[GAP-002] ACL: Granted 'write' to client {owner_client_id} on KB '{name}'"
                            )

                    except Exception as acl_error:
                        logger.warning(
                            f"[GAP-002] Could not auto-grant ACL permission: {acl_error}"
                        )

                return {
                    "collection_id": name,
                    "status": "created",
                    "message": f"Knowledge base '{name}' created successfully",
                    "kb_type": kb_type,
                    "owner_client_id": owner_client_id,
                    "auth_source": auth_source,
                }
            else:
                error_msg = result.get("error", "Unknown error")
                logger.error(f"[GAP-002] Failed to create knowledge base: {error_msg}")
                return {
                    "collection_id": name,
                    "status": "error",
                    "message": f"Failed to create knowledge base: {error_msg}",
                }

        except Exception as e:
            logger.error(f"[GAP-002] Error creating knowledge base: {e}")
            return {
                "collection_id": name,
                "status": "error",
                "message": f"Error creating knowledge base: {str(e)}",
            }

    async def ingest_document(
        self,
        collection_id: str,
        text: str,
        metadata: Optional[Dict[str, Any]] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Ingest a text document into a knowledge base.

        Security: Admin only operation.

        Args:
            collection_id: Target collection
            text: Document text
            metadata: Optional metadata dict
            ctx: Security context

        Returns:
            Status with document_id and chunks_count
        """
        ctx = self._require_admin(ctx, "ingest_document")
        self._require_initialized()

        metadata = metadata or {}

        if not collection_id or not text:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "collection_id and text are required",
            }

        # TASK #82: Inject uploader_id from security context for document lifecycle tracking
        if ctx and hasattr(ctx, "user") and ctx.user:
            metadata["uploader_id"] = getattr(ctx.user, "user_id", "system")
        else:
            metadata["uploader_id"] = "system"

        logger.info(
            "Ingesting document",
            extra={
                "collection": collection_id,
                "text_length": len(text),
                "admin_user_id": ctx.user.user_id,
                "uploader_id": metadata.get("uploader_id"),
            },
        )

        try:
            # --- DEDUP CHECK ---
            content_hash = hashlib.md5(text.encode()).hexdigest()
            existing = await self.qdrant_module.check_duplicate_internal(
                collection=collection_id,
                content_hash=content_hash,
                filename=metadata.get("filename"),
            )
            if existing:
                logger.info(
                    f"DEDUP: Skipping duplicate document "
                    f"[collection={collection_id}, file={metadata.get('filename')}, "
                    f"existing_doc={existing.get('doc_id')}, "
                    f"existing_file={existing.get('filename')}, hash={content_hash}]"
                )
                return {
                    "document_id": existing.get("doc_id", ""),
                    "chunks_count": existing.get("chunk_count", 0),
                    "status": "duplicate",
                    "message": f"Document already ingested as '{existing.get('filename')}' "
                               f"(uploaded {existing.get('upload_timestamp', 'unknown')})",
                }

            # Generate document ID
            document_id = str(uuid.uuid4())

            # Enrich metadata with kb_id, collection and auto-tags
            tags = kwargs.get("tags")
            _enrich_ingest_metadata(metadata, collection_id, explicit_tags=tags)

            # Delegate indexing to rag_qdrant (embeddings + chunking live there)
            # Use internal method - admin auth already verified by _require_admin above
            result = await self.qdrant_module.add_document_internal(
                doc_id=document_id,
                text=text,
                metadata=metadata,
                collection=collection_id,
            )

            # rag_qdrant returns simplified format: {status: indexed|failed, ...}
            if result.get("status") != "indexed":
                error_msg = result.get("error") or "Failed to ingest document"
                logger.error(
                    "❌ Document ingestion FAILED",
                    extra={
                        "document_id": document_id,
                        "collection": collection_id,
                        "doc_filename": metadata.get("filename", "unknown"),
                        "text_length": len(text),
                        "uploader_id": metadata.get("uploader_id"),
                        "error": error_msg,
                        "qdrant_result": result,
                    },
                )
                return {
                    "document_id": document_id,
                    "chunks_count": 0,
                    "status": "error",
                    "message": error_msg,
                }

            self.total_documents_ingested += 1
            chunks_count = int(result.get("chunks_count") or 0)

            # FEAT-DKI-001 v2.0: Configurable keyword extraction
            keyword_result = {"status": "disabled", "keywords_count": 0}
            if self.keyword_manager.enabled:
                try:
                    keyword_result = await self.keyword_manager.extract_and_store_keywords(
                        collection_name=collection_id,
                        text=text,
                    )
                    logger.info(
                        f"DKI: Keywords extracted for {collection_id}",
                        extra={"keyword_result": keyword_result},
                    )
                except Exception as kw_err:
                    logger.warning(
                        f"DKI: Keyword extraction failed (non-blocking): {kw_err}"
                    )

            logger.info(
                "✅ Document ingested",
                extra={
                    "document_id": document_id,
                    "collection": collection_id,
                    "chunks": chunks_count,
                    "keywords_extracted": keyword_result.get("keywords_extracted", 0),
                },
            )

            return {
                "document_id": document_id,
                "chunks_count": chunks_count,
                "status": "success",
                "message": f"Document ingested with {chunks_count} chunks",
                "keywords": keyword_result,  # Include keyword extraction result
            }

        except Exception as e:
            import traceback
            error_traceback = traceback.format_exc()
            logger.error(
                f"❌ Document ingestion EXCEPTION: {type(e).__name__}: {e} "
                f"[collection={collection_id}, file={metadata.get('filename', 'unknown')}, "
                f"text_len={len(text) if text else 0}]\n{error_traceback}"
            )
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": f"Error ingesting document: {str(e)}",
            }

    async def ingest_file(
        self,
        collection_name: str,
        file_content: str,
        filename: str,
        file_type: str = "text/plain",
        chunking_config: Optional[Dict[str, Any]] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Ingest a file into a knowledge base with server-side text extraction.

        Handles text extraction from various file types:
        - TXT, MD, JSON: Direct text (UTF-8)
        - PDF: PyPDF2 extraction (base64)
        - Excel (.xlsx, .xls): openpyxl tabular (base64)
        - DOCX: python-docx extraction (base64)
        - PPTX: python-pptx slide extraction (base64)
        - CSV: Tabular text extraction (base64)
        - HTML: BeautifulSoup extraction (base64)
        - XML: lxml text extraction (base64)
        - YAML (.yaml, .yml): PyYAML parsing (base64)

        Security: Admin only operation.

        Args:
            collection_name: Target collection
            file_content: File content (base64 for PDF, plaintext for TXT/MD)
            filename: Original filename
            file_type: MIME type (default text/plain)
            chunking_config: Optional override for chunk_size, chunk_overlap
            ctx: Security context

        Returns:
            Status with document_id, chunks_count, extraction_method
        """
        ctx = self._require_admin(ctx, "ingest_file")
        self._require_initialized()

        if not collection_name:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "collection_name is required",
            }
        if not file_content:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "file_content is required",
            }
        if not filename:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "filename is required",
            }

        # Extract text based on file type
        text = ""
        extraction_method = "direct"

        try:
            # Normalize file type detection
            is_pdf = file_type == "application/pdf" or filename.lower().endswith(".pdf")
            is_json = file_type == "application/json" or filename.lower().endswith(
                ".json"
            )
            # FIX-EXCEL-001: Add Excel file detection
            excel_mimes = [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ]
            is_excel = (
                file_type in excel_mimes
                or filename.lower().endswith(".xlsx")
                or filename.lower().endswith(".xls")
            )

            # FIX-DOCX-001: Add DOCX file detection
            is_docx = (
                file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                or filename.lower().endswith(".docx")
            )

            # FIX-CSV-001: Add CSV file detection
            is_csv = (
                "csv" in file_type.lower()
                or filename.lower().endswith(".csv")
            )

            # FIX-HTML-001: Add HTML file detection
            is_html = (
                "html" in file_type.lower()
                or filename.lower().endswith(".html")
                or filename.lower().endswith(".htm")
            )

            # FIX-PPTX-001: Add PowerPoint file detection
            is_pptx = (
                file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                or filename.lower().endswith(".pptx")
            )

            # FIX-XML-001: Add XML file detection
            is_xml = (
                "xml" in file_type.lower()
                or filename.lower().endswith(".xml")
            )

            # FIX-YAML-001: Add YAML file detection
            is_yaml = (
                "yaml" in file_type.lower()
                or filename.lower().endswith(".yaml")
                or filename.lower().endswith(".yml")
            )

            if is_pdf:
                # PDF extraction: decode base64 and extract text
                # TODO v1.7: Implement OCR for image-based PDFs using pytesseract/tesseract
                # Current limitation: PyPDF2 only extracts embedded text, not text in images
                # OCR would require: pip install pytesseract pdf2image
                # and system dependency: tesseract-ocr
                import base64
                import io

                try:
                    from PyPDF2 import PdfReader

                    pdf_bytes = base64.b64decode(file_content)
                    pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
                    text_parts = []
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    text = "\n\n".join(text_parts)
                    extraction_method = "pypdf2"
                    logger.info(
                        f"PDF extracted: {len(pdf_reader.pages)} pages, {len(text)} chars"
                    )
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "PDF extraction requires PyPDF2. Install with: pip install PyPDF2",
                    }
                except Exception as pdf_err:
                    logger.error(f"PDF extraction failed: {pdf_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"PDF extraction failed: {str(pdf_err)}",
                    }

            elif is_json:
                # JSON: parse and re-serialize formatted
                import json as json_module

                try:
                    data = json_module.loads(file_content)
                    text = json_module.dumps(data, indent=2, ensure_ascii=False)
                    extraction_method = "json"
                except json_module.JSONDecodeError as je:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"Invalid JSON: {str(je)}",
                    }

            elif is_excel:
                # FIX-EXCEL-001: Excel extraction using openpyxl
                import base64
                import io

                try:
                    from openpyxl import load_workbook

                    excel_bytes = base64.b64decode(file_content)
                    workbook = load_workbook(
                        io.BytesIO(excel_bytes), read_only=True, data_only=True
                    )
                    text_parts = []

                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        sheet_lines = [f"=== Sheet: {sheet_name} ==="]

                        for row in sheet.iter_rows(values_only=True):
                            # Filter out completely empty rows
                            if any(cell is not None for cell in row):
                                row_values = [
                                    str(cell) if cell is not None else ""
                                    for cell in row
                                ]
                                sheet_lines.append(" | ".join(row_values))

                        if len(sheet_lines) > 1:  # Has data beyond header
                            text_parts.append("\n".join(sheet_lines))

                    workbook.close()
                    text = "\n\n".join(text_parts)
                    extraction_method = "openpyxl"
                    logger.info(
                        f"Excel extracted: {len(workbook.sheetnames)} sheets, {len(text)} chars"
                    )
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "Excel extraction requires openpyxl. Install with: pip install openpyxl",
                    }
                except Exception as excel_err:
                    logger.error(f"Excel extraction failed: {excel_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"Excel extraction failed: {str(excel_err)}",
                    }

            elif is_docx:
                # FIX-DOCX-001: DOCX extraction using python-docx
                import base64
                import io

                try:
                    from docx import Document

                    docx_bytes = base64.b64decode(file_content)
                    doc = Document(io.BytesIO(docx_bytes))
                    text_parts = []

                    for para in doc.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)

                    # Also extract text from tables
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                text_parts.append(row_text)

                    text = "\n\n".join(text_parts)
                    extraction_method = "python-docx"
                    logger.info(f"DOCX extracted: {len(text_parts)} paragraphs, {len(text)} chars")
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "DOCX extraction requires python-docx. Install with: pip install python-docx",
                    }
                except Exception as docx_err:
                    logger.error(f"DOCX extraction failed: {docx_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"DOCX extraction failed: {str(docx_err)}",
                    }

            elif is_csv:
                # FIX-CSV-001: CSV extraction to tabular text
                import base64
                import csv as csv_module
                import io

                try:
                    csv_bytes = base64.b64decode(file_content)
                    # Try UTF-8 first, then fallback to latin-1
                    try:
                        csv_content = csv_bytes.decode("utf-8")
                    except UnicodeDecodeError:
                        csv_content = csv_bytes.decode("latin-1")

                    reader = csv_module.reader(io.StringIO(csv_content))
                    lines = []

                    for row in reader:
                        if any(cell.strip() for cell in row):
                            lines.append(" | ".join(cell.strip() for cell in row))

                    text = "\n".join(lines)
                    extraction_method = "csv"
                    logger.info(f"CSV extracted: {len(lines)} rows, {len(text)} chars")
                except Exception as csv_err:
                    logger.error(f"CSV extraction failed: {csv_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"CSV extraction failed: {str(csv_err)}",
                    }

            elif is_html:
                # FIX-HTML-001: HTML extraction using BeautifulSoup
                import base64
                import re

                try:
                    from bs4 import BeautifulSoup

                    html_bytes = base64.b64decode(file_content)
                    # Try UTF-8 first, then fallback
                    try:
                        html_content = html_bytes.decode("utf-8")
                    except UnicodeDecodeError:
                        html_content = html_bytes.decode("latin-1")

                    soup = BeautifulSoup(html_content, "html.parser")

                    # Remove script and style elements
                    for element in soup(["script", "style", "head", "meta", "link"]):
                        element.decompose()

                    # Get text with reasonable spacing
                    text = soup.get_text(separator="\n", strip=True)
                    # Clean up multiple newlines
                    text = re.sub(r"\n{3,}", "\n\n", text)
                    extraction_method = "beautifulsoup"
                    logger.info(f"HTML extracted: {len(text)} chars")
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "HTML extraction requires beautifulsoup4. Install with: pip install beautifulsoup4",
                    }
                except Exception as html_err:
                    logger.error(f"HTML extraction failed: {html_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"HTML extraction failed: {str(html_err)}",
                    }

            elif is_pptx:
                # FIX-PPTX-001: PowerPoint extraction using python-pptx
                import base64
                import io

                try:
                    from pptx import Presentation

                    pptx_bytes = base64.b64decode(file_content)
                    prs = Presentation(io.BytesIO(pptx_bytes))
                    text_parts = []

                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_texts = [f"=== Slide {slide_num} ==="]
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_texts.append(shape.text.strip())
                        if len(slide_texts) > 1:
                            text_parts.append("\n".join(slide_texts))

                    text = "\n\n".join(text_parts)
                    extraction_method = "python-pptx"
                    logger.info(f"PPTX extracted: {len(prs.slides)} slides, {len(text)} chars")
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "PPTX extraction requires python-pptx. Install with: pip install python-pptx",
                    }
                except Exception as pptx_err:
                    logger.error(f"PPTX extraction failed: {pptx_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"PPTX extraction failed: {str(pptx_err)}",
                    }

            elif is_xml:
                # FIX-XML-001: XML extraction using lxml
                import base64

                try:
                    from lxml import etree

                    xml_bytes = base64.b64decode(file_content)
                    # Try UTF-8 first, then fallback
                    try:
                        xml_content = xml_bytes.decode("utf-8")
                    except UnicodeDecodeError:
                        xml_content = xml_bytes.decode("latin-1")

                    root = etree.fromstring(xml_content.encode())

                    # Extract all text content from XML
                    texts = []
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
                        if elem.tail and elem.tail.strip():
                            texts.append(elem.tail.strip())

                    text = "\n".join(texts)
                    extraction_method = "lxml"
                    logger.info(f"XML extracted: {len(texts)} text nodes, {len(text)} chars")
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "XML extraction requires lxml. Install with: pip install lxml",
                    }
                except Exception as xml_err:
                    logger.error(f"XML extraction failed: {xml_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"XML extraction failed: {str(xml_err)}",
                    }

            elif is_yaml:
                # FIX-YAML-001: YAML extraction using PyYAML
                import base64

                try:
                    import yaml

                    yaml_bytes = base64.b64decode(file_content)
                    # Try UTF-8 first, then fallback
                    try:
                        yaml_content = yaml_bytes.decode("utf-8")
                    except UnicodeDecodeError:
                        yaml_content = yaml_bytes.decode("latin-1")

                    data = yaml.safe_load(yaml_content)
                    text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
                    extraction_method = "pyyaml"
                    logger.info(f"YAML extracted: {len(text)} chars")
                except ImportError:
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": "YAML extraction requires PyYAML. Install with: pip install PyYAML",
                    }
                except Exception as yaml_err:
                    logger.error(f"YAML extraction failed: {yaml_err}")
                    return {
                        "document_id": "",
                        "chunks_count": 0,
                        "status": "error",
                        "message": f"YAML extraction failed: {str(yaml_err)}",
                    }

            else:
                # Plain text files (TXT, MD, etc.)
                # BUG-005 FIX: Handle base64-encoded text files from frontend
                import base64

                # Check if content looks like base64 (no whitespace at start, valid base64 chars)
                # Frontend typically sends base64 for all files read via FileReader.readAsDataURL
                content_to_process = file_content

                # Try to detect and decode base64
                # Base64 pattern: alphanumeric + /+ and = padding, no newlines in middle
                if file_content and len(file_content) > 0:
                    # Strip potential data URL prefix (e.g., "data:text/plain;base64,")
                    if ";base64," in file_content:
                        content_to_process = file_content.split(";base64,", 1)[1]

                    # Heuristic: if content looks like base64 (valid chars, length multiple of 4)
                    # and doesn't look like natural text, try to decode
                    is_likely_base64 = False
                    stripped = content_to_process.strip()

                    # Check if it's valid base64 format
                    import re

                    base64_pattern = re.compile(r"^[A-Za-z0-9+/]*={0,2}$")
                    if (
                        len(stripped) > 20
                        and len(stripped) % 4 == 0
                        and base64_pattern.match(
                            stripped.replace("\n", "").replace("\r", "")
                        )
                    ):
                        is_likely_base64 = True

                    if is_likely_base64:
                        try:
                            decoded_bytes = base64.b64decode(content_to_process)
                            # Try to decode as UTF-8 text
                            text = decoded_bytes.decode("utf-8")
                            extraction_method = "plaintext_base64_decoded"
                            logger.info(
                                f"Base64 text file decoded: {len(file_content)} -> {len(text)} chars"
                            )
                        except (ValueError, UnicodeDecodeError) as decode_err:
                            # Not valid base64 or not UTF-8, use as-is
                            logger.debug(
                                f"Base64 decode failed, using raw content: {decode_err}"
                            )
                            text = file_content
                            extraction_method = "plaintext"
                    else:
                        # Not base64, use as plaintext
                        text = file_content
                        extraction_method = "plaintext"
                else:
                    text = file_content
                    extraction_method = "plaintext"

        except Exception as e:
            logger.error(f"Error extracting text from file: {e}")
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": f"Failed to extract text from file: {str(e)}",
            }

        if not text or not text.strip():
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "No text could be extracted from file",
            }

        # Build metadata
        metadata = {
            "filename": filename,
            "file_type": file_type,
            "extraction_method": extraction_method,
            "source": "file_upload",
            "original_size": len(file_content),
            "extracted_size": len(text),
        }

        # Add chunking config if provided
        if chunking_config:
            metadata["chunking_config"] = chunking_config

        logger.info(
            f"Ingesting file: {filename} ({file_type}) -> {collection_name}",
            extra={
                "extraction_method": extraction_method,
                "text_length": len(text),
                "admin_user_id": ctx.user.user_id,
            },
        )

        # Delegate to ingest_document (reuse existing logic)
        return await self.ingest_document(
            collection_id=collection_name,
            text=text,
            metadata=metadata,
            ctx=ctx,
            tags=kwargs.get("tags"),
        )

    async def ingest_document_authorized(
        self,
        collection_id: str,
        text: str,
        metadata: Optional[Dict[str, Any]] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Ingest a text document into a knowledge base with ACL-based authorization.

        Security: Requires write permission on the target collection.
        Fixes GAP-INGEST-002: Allows users/clients to ingest to their authorized KBs.

        Args:
            collection_id: Target collection
            text: Document text
            metadata: Optional metadata dict
            ctx: Security context (required)

        Returns:
            Status with document_id and chunks_count
        """
        # Require authentication context
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        metadata = metadata or {}

        if not collection_id or not text:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "collection_id and text are required",
            }

        # Extract user/client identity from context
        user_id = getattr(ctx.user, "user_id", None)
        client_id = getattr(ctx.user, "client_id", None)

        if not user_id:
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": "User identity required for authorization",
            }

        # Check if user/client has write permission on the collection
        if not self.acl_manager:
            raise RuntimeError("ACL manager not initialized")

        has_write = await self.acl_manager.check_write_access(
            user_id, client_id, collection_id
        )

        if not has_write:
            logger.warning(
                "Unauthorized ingest attempt",
                extra={
                    "user_id": user_id,
                    "client_id": client_id,
                    "collection": collection_id,
                },
            )
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": f"Write permission required for collection {collection_id}",
            }

        # Track uploader for audit
        metadata["uploader_id"] = user_id
        metadata["client_id"] = client_id if client_id else "none"

        logger.info(
            "Ingesting document (authorized)",
            extra={
                "collection": collection_id,
                "text_length": len(text),
                "user_id": user_id,
                "client_id": client_id,
            },
        )

        try:
            # Generate document ID
            document_id = str(uuid.uuid4())

            # Enrich metadata with kb_id, collection and auto-tags
            tags = kwargs.get("tags")
            _enrich_ingest_metadata(metadata, collection_id, explicit_tags=tags)

            # Delegate indexing to rag_qdrant
            # Use internal method - ACL write permission already verified above
            result = await self.qdrant_module.add_document_internal(
                doc_id=document_id,
                text=text,
                metadata=metadata,
                collection=collection_id,
            )

            if result.get("status") != "indexed":
                return {
                    "document_id": document_id,
                    "chunks_count": 0,
                    "status": "error",
                    "message": result.get("error") or "Failed to ingest document",
                }

            self.total_documents_ingested += 1
            chunks_count = int(result.get("chunks_count") or 0)

            # FEAT-DKI-001 v2.0: Configurable keyword extraction
            keyword_result = {"status": "disabled", "keywords_count": 0}
            if self.keyword_manager.enabled:
                try:
                    keyword_result = await self.keyword_manager.extract_and_store_keywords(
                        collection_name=collection_id,
                        text=text,
                    )
                    logger.info(
                        f"DKI: Keywords extracted for {collection_id} (authorized)",
                        extra={"keyword_result": keyword_result},
                    )
                except Exception as kw_err:
                    logger.warning(
                        f"DKI: Keyword extraction failed (non-blocking): {kw_err}"
                    )

            logger.info(
                "✅ Document ingested (authorized)",
                extra={
                    "document_id": document_id,
                    "collection": collection_id,
                    "chunks": chunks_count,
                    "user_id": user_id,
                    "keywords_extracted": keyword_result.get("keywords_extracted", 0),
                },
            )

            return {
                "document_id": document_id,
                "chunks_count": chunks_count,
                "status": "success",
                "message": f"Document ingested with {chunks_count} chunks",
                "keywords": keyword_result,
            }

        except Exception as e:
            logger.error(f"Error ingesting document (authorized): {e}")
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": f"Error ingesting document: {str(e)}",
            }

    # ===== ACL Management (Delegated to ACLManager) =====

    async def set_permission(
        self,
        entity_type: str,
        entity_id: str,
        collection_id: str,
        access_level: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Set access permission for an entity on a collection.

        Security: Admin only operation.

        Args:
            entity_type: "user" or "client"
            entity_id: user_id or client_id
            collection_id: Target collection
            access_level: "read", "write", or "none"
            ctx: Security context

        Returns:
            Status dict
        """
        ctx = self._require_admin(ctx, "set_permission")
        self._require_initialized()

        logger.info(
            "Setting permission",
            extra={
                "entity_type": entity_type,
                "entity_id": entity_id,
                "collection": collection_id,
                "access": access_level,
                "admin_user_id": ctx.user.user_id,
            },
        )

        # Delegate to ACL manager (already validated as not None)
        return await self.acl_manager.set_permission(
            entity_type=entity_type,
            entity_id=entity_id,
            collection_id=collection_id,
            access_level=access_level,
        )

    async def set_permission_internal(
        self,
        entity_type: str,
        entity_id: str,
        collection_id: str,
        access_level: str,
    ) -> Dict[str, Any]:
        """Set permissions without admin context for trusted internal callers.

        Caller modules (e.g., admin_clients) MUST perform their own security checks
        before invoking this helper.
        """
        self._require_initialized()
        if not self.acl_manager:
            raise RuntimeError("ACL manager not initialized")

        logger.debug(
            "[INTERNAL] Setting permission",
            extra={
                "entity_type": entity_type,
                "entity_id": entity_id,
                "collection": collection_id,
                "access": access_level,
            },
        )

        return await self.acl_manager.set_permission(
            entity_type=entity_type,
            entity_id=entity_id,
            collection_id=collection_id,
            access_level=access_level,
        )

    async def get_permissions(
        self,
        collection_id: Optional[str] = None,
        entity_type: Optional[str] = None,
        entity_id: Optional[str] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Get permissions (optionally filtered).

        Security: Admin can see all, users can see only their own.

        Args:
            collection_id: Optional filter
            entity_type: Optional filter
            entity_id: Optional filter
            ctx: Security context

        Returns:
            Dict with permissions array and count
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        # Security: Non-admins can only see their own permissions
        if not self._is_admin(ctx):
            entity_type = "user"
            entity_id = ctx.user.user_id
            logger.info(
                "Non-admin user requesting permissions - filtered to own user",
                extra={"user_id": ctx.user.user_id},
            )

        # Delegate to ACL manager
        return await self.acl_manager.get_permissions(
            collection_id=collection_id, entity_type=entity_type, entity_id=entity_id
        )

    # ===== RAG Configuration Management (Delegated to ConfigManager) =====

    async def set_rag_config(
        self,
        entity_type: str,
        entity_id: Optional[str] = None,
        config_json: Optional[Dict[str, Any]] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Set RAG configuration for an entity.

        Security: Admin only operation.

        Args:
            entity_type: "user", "client", or "default"
            entity_id: user_id or client_id (null for default)
            config_json: RAG configuration dict
            ctx: Security context

        Returns:
            Status dict
        """
        ctx = self._require_admin(ctx, "set_rag_config")
        self._require_initialized()

        logger.info(
            "Setting RAG config",
            extra={
                "entity_type": entity_type,
                "entity_id": entity_id,
                "admin_user_id": ctx.user.user_id,
            },
        )

        # Delegate to config manager
        return await self.config_manager.set_rag_config(
            entity_type=entity_type,
            entity_id=entity_id,
            config_json=config_json or {},
        )

    async def get_user_capabilities(self, ctx: Any) -> Dict[str, Any]:
        """Expose facade-friendly capability snapshot."""

        if not self.capability_manager:
            self.capability_manager = CapabilityManager(adapter=self)
        client_id = getattr(ctx.user, "client_id", None)
        return await self.capability_manager.get_full_user_context(
            user_id=ctx.user.user_id,
            client_id=client_id,
            ctx=ctx,
        )

    async def _build_enrichment_pipeline_config(
        self,
        ctx: Any,
        pipeline_options: Optional[Dict[str, Any]] = None,
        target_collections: Optional[List[str]] = None,
    ) -> Optional[Dict[str, Any]]:
        """
        Compute permission-aware enrichment pipeline configuration with 4-layer merge.

        Priority order (highest to lowest):
        1. User Request (pipeline_options) - explicit override wins
        2. Collection Metadata (redis) - domain-specific config
        3. Client Policy (client_config) - tenant permissions
        4. System Defaults (.env) - fallback

        Args:
            ctx: Security context
            pipeline_options: User-provided options (highest priority)
            target_collections: List of collections being queried (for metadata lookup)
        """

        pipeline_options = pipeline_options or {}
        target_collections = target_collections or []

        try:
            capabilities = await self.get_user_capabilities(ctx)
        except Exception as exc:  # pragma: no cover - diagnostic only
            logger.warning(
                "Could not fetch capabilities for enrichment pipeline: %s", exc
            )
            capabilities = {}

        features = capabilities.get("features", {}) or {}
        client_scope = capabilities.get("client", {}) or {}
        client_features = client_scope.get("features", {}) or {}

        enrichment_caps = features.get("enrichment", {}) or {}
        system_steps = enrichment_caps.get("steps", {}) or {}

        client_enrichment = client_features.get("enrichment")
        client_steps: Dict[str, Any] = {}
        client_enrichment_enabled = True

        def _safe_bool(value: Any, default: bool) -> bool:
            return value if isinstance(value, bool) else default

        def _get_flag_nullable(
            source_dict: Dict[str, Any], keys: list
        ) -> Optional[bool]:
            """Helper to search for a feature with multiple possible names (aliasing).
            Returns None if not found (for priority chain)."""
            for k in keys:
                if k in source_dict:
                    val = source_dict[k]
                    if isinstance(val, bool):
                        return val
            return None

        def _get_flag(source_dict: Dict[str, Any], keys: list, default: bool) -> bool:
            """Helper with default value."""
            result = _get_flag_nullable(source_dict, keys)
            return result if result is not None else default

        client_enrichment_enabled = _safe_bool(
            client_features.get("enrichment_enabled"), True
        )

        if isinstance(client_enrichment, dict):
            client_enrichment_enabled = _safe_bool(
                client_enrichment.get("enabled"), client_enrichment_enabled
            )
            potential_steps = client_enrichment.get("steps")
            if isinstance(potential_steps, dict):
                client_steps = potential_steps
        elif isinstance(client_enrichment, bool):
            client_enrichment_enabled = _safe_bool(
                client_enrichment, client_enrichment_enabled
            )

        system_available = _safe_bool(enrichment_caps.get("available"), True)
        use_pipeline_from_user = pipeline_options.get("use_enrichment")

        # Default to True only if user didn't specify
        # Log when using default vs user preference
        if use_pipeline_from_user is None:
            use_pipeline = True
            logger.debug(
                "[RAG_CONFIG] use_enrichment not specified by user, using default: True",
                extra={"source": "system_default"},
            )
        else:
            use_pipeline = use_pipeline_from_user
            logger.info(
                f"[RAG_CONFIG] use_enrichment explicitly set by user: {use_pipeline}",
                extra={"source": "user_request", "value": use_pipeline},
            )

        pipeline_enabled = (
            system_available and client_enrichment_enabled and use_pipeline
        )

        # Log when pipeline is disabled and why
        if not pipeline_enabled:
            reasons = []
            if not system_available:
                reasons.append("system_unavailable")
            if not client_enrichment_enabled:
                reasons.append("client_disabled")
            if not use_pipeline:
                reasons.append("user_disabled")
            logger.warning(
                f"[RAG_CONFIG] Enrichment pipeline DISABLED. Reasons: {', '.join(reasons)}",
                extra={
                    "pipeline_enabled": False,
                    "system_available": system_available,
                    "client_enrichment_enabled": client_enrichment_enabled,
                    "use_pipeline": use_pipeline,
                    "disabled_reasons": reasons,
                },
            )

        # === LAYER 3: Collection Metadata ===
        # Fetch enrichment config from collection metadata (if any collection requires a feature, enable it)
        meta_tasks = [
            self.get_collection_metadata_internal(coll_name)
            for coll_name in target_collections
        ]
        meta_results = await asyncio.gather(*meta_tasks, return_exceptions=True)
        collection_metadata_list: List[Dict[str, Any]] = []
        for coll_name, meta in zip(target_collections, meta_results):
            if isinstance(meta, Exception):
                logger.debug(
                    "Could not fetch metadata for collection %s: %s", coll_name, meta
                )
            elif meta:
                collection_metadata_list.append(meta)

        # Alias definitions for backward compatibility and flexible naming
        hyde_keys = ["hyde", "hyde_enabled", "use_hyde"]
        exp_keys = ["query_expansion", "query_expansion_enabled", "use_expansion"]
        investigative_keys = [
            "investigative",
            "investigative_enabled",
            "use_investigative",
        ]
        filter_keys = ["query_filters", "query_filters_enabled", "use_filters"]
        rerank_keys = ["rerank", "rerank_enabled", "use_rerank"]
        fusion_keys = ["fusion", "fusion_enabled", "use_fusion"]
        dedup_keys = ["dedup", "dedup_enabled", "deduplication_enabled", "use_dedup"]
        compression_keys = ["compression", "compression_enabled", "use_compression"]

        def _resolve_step(
            aliases: list, system_default: bool, step_name: str = ""
        ) -> bool:
            """
            Resolve step enabled status using 4-layer priority.

            Priority:
            1. User Request (pipeline_options) - explicit override wins
            2. Collection Metadata - if ANY collection requires it, enable
            3. Client Policy
            4. System Defaults
            """
            if not pipeline_enabled:
                return False

            # Layer 4 (highest): User explicit override
            user_val = _get_flag_nullable(pipeline_options, aliases)
            if user_val is not None:
                logger.debug(
                    f"[RAG_CONFIG] {step_name}: Using USER REQUEST layer: {user_val}",
                    extra={
                        "step": step_name,
                        "value": user_val,
                        "layer": "user_request",
                    },
                )
                return user_val

            # Layer 3: Collection metadata (OR logic - any collection requiring it = True)
            collection_val: Optional[bool] = None
            for coll_meta in collection_metadata_list:
                coll_enrichment = coll_meta.get("enrichment", {}) or {}
                coll_steps = coll_enrichment.get("steps", {}) or {}
                coll_flag = _get_flag_nullable(coll_steps, aliases)
                if coll_flag is True:
                    collection_val = True
                    break  # One True is enough
                elif coll_flag is False and collection_val is None:
                    collection_val = False
            if collection_val is not None:
                logger.debug(
                    f"[RAG_CONFIG] {step_name}: Using COLLECTION METADATA layer: {collection_val}",
                    extra={
                        "step": step_name,
                        "value": collection_val,
                        "layer": "collection_metadata",
                    },
                )
                return collection_val

            # Layer 2: Client policy
            client_val = _get_flag_nullable(client_steps, aliases)
            if client_val is not None:
                logger.debug(
                    f"[RAG_CONFIG] {step_name}: Using CLIENT POLICY layer: {client_val}",
                    extra={
                        "step": step_name,
                        "value": client_val,
                        "layer": "client_policy",
                    },
                )
                return client_val

            # Layer 1 (lowest): System defaults
            system_val = _get_flag(system_steps, aliases, system_default)
            logger.debug(
                f"[RAG_CONFIG] {step_name}: Using SYSTEM DEFAULT layer: {system_val}",
                extra={
                    "step": step_name,
                    "value": system_val,
                    "layer": "system_default",
                    "fallback": True,
                },
            )
            return system_val

        # Resolve each step using 4-layer priority
        # FIX-PROP-001 v4.2.0: Use system_steps from CapabilityManager (populated from EnrichmentSettings)
        # instead of hardcoded defaults. This ensures ENV variables propagate correctly.
        final_hyde = _resolve_step(hyde_keys, system_steps.get("hyde", False), "hyde")
        final_exp = _resolve_step(exp_keys, system_steps.get("query_expansion", False), "query_expansion")
        final_investigative = _resolve_step(investigative_keys, system_steps.get("investigative", False), "investigative")
        final_filter = _resolve_step(filter_keys, system_steps.get("query_filters", True), "query_filters")
        final_rerank = _resolve_step(rerank_keys, system_steps.get("rerank", True), "rerank")
        final_fusion = _resolve_step(fusion_keys, system_steps.get("fusion", False), "fusion")
        final_dedup = _resolve_step(dedup_keys, system_steps.get("dedup", False), "dedup")
        final_compression = _resolve_step(compression_keys, system_steps.get("compression", False), "compression")

        # LOG 4-LAYER RESOLUTION RESULT - Critical for debugging config propagation
        logger.info(
            f"[RAG_CONFIG] 4-layer enrichment resolution: "
            f"hyde={final_hyde}, expansion={final_exp}, investigative={final_investigative}, "
            f"filters={final_filter}, rerank={final_rerank}, fusion={final_fusion}, "
            f"dedup={final_dedup}, compression={final_compression}, pipeline_enabled={pipeline_enabled}",
            extra={
                "user_request_keys": list(pipeline_options.keys())
                if pipeline_options
                else [],
                "collections_checked": len(collection_metadata_list),
                "final_hyde": final_hyde,
                "final_expansion": final_exp,
                "final_investigative": final_investigative,
                "final_filters": final_filter,
                "final_rerank": final_rerank,
                "final_fusion": final_fusion,
                "final_dedup": final_dedup,
                "final_compression": final_compression,
                "pipeline_enabled": pipeline_enabled,
            },
        )

        pipeline_config = {
            "enabled": pipeline_enabled,
            "steps": [
                {"step": "hyde", "enabled": final_hyde},
                {"step": "query_expansion", "enabled": final_exp},
                {"step": "investigative", "enabled": final_investigative},
                {"step": "query_filters", "enabled": final_filter},
                {"step": "rerank", "enabled": final_rerank},
                {"step": "fusion", "enabled": final_fusion},
                {"step": "dedup", "enabled": final_dedup},
                {"step": "compression", "enabled": final_compression},
            ],
            "options": pipeline_options,
            "_layers_debug": {
                "collections_with_metadata": len(collection_metadata_list),
                "target_collections": target_collections,
                "user_request_params": list(pipeline_options.keys())
                if pipeline_options
                else [],
                "resolution": {
                    "hyde": final_hyde,
                    "query_expansion": final_exp,
                    "investigative": final_investigative,
                    "query_filters": final_filter,
                    "rerank": final_rerank,
                    "fusion": final_fusion,
                    "dedup": final_dedup,
                    "compression": final_compression,
                },
            },
        }

        return pipeline_config

    async def get_rag_config(
        self,
        entity_type: Optional[str] = None,
        entity_id: Optional[str] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Get RAG configuration for an entity (with fallback to default).

        Security: Users can get their own config, admins can get any.

        Args:
            entity_type: "user", "client", or "default"
            entity_id: user_id or client_id
            ctx: Security context

        Returns:
            Dict with config and source
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        entity_type = entity_type or "user"

        # Security: Non-admins can only get their own config
        if not self._is_admin(ctx):
            entity_type = "user"
            entity_id = ctx.user.user_id
            logger.info(
                "Non-admin user requesting config - using own user_id",
                extra={"user_id": ctx.user.user_id},
            )

        # Delegate to config manager
        return await self.config_manager.get_rag_config(
            entity_type=entity_type,
            entity_id=entity_id,
            default_config=self.config.get("default_rag_config", {}),
        )

    # ===== RAG Chat (The Magic Function) =====

    async def rag_chat(
        self,
        query: str,
        collections: Optional[List[str]] = None,
        override_config: Optional[Dict[str, Any]] = None,
        web_context: Optional[str] = None,
        conversation_id: Optional[str] = None,
        ctx: Any = None,
        # FIX-BUG-002 v1.8.3: Parameter alias for API consistency
        collection_ids: Optional[List[str]] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        The magic function - Execute RAG pipeline with ACL checks.

        Security: Uses ctx.user for identity, enforces ACL.

        Args:
            query: User question
            collections: Optional list of collections to query (preferred parameter name)
            override_config: Optional config overrides
            web_context: Optional pre-fetched web search results as formatted text
                         (ROADMAP v1.5.0 - FEAT-WEB-001)
            conversation_id: Optional conversation ID to continue
                            (ROADMAP v1.5.0 - FEAT-MEM-001 / Task #15)
                            If not provided, a new conversation is created.
            ctx: Security context (provides user_id and client_id)
            collection_ids: Alias for 'collections' (FIX-BUG-002 v1.8.3 - API consistency)

        Returns:
            Dict with answer, sources, config_used, permissions_checked, conversation_id, debug
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        # FIX-BUG-002 v1.8.3: Accept both 'collections' and 'collection_ids' parameter names
        # API callers may use either name - we normalize to 'collections' internally
        if collections is None and collection_ids is not None:
            collections = collection_ids
            logger.debug(
                "Parameter alias used: collection_ids -> collections",
                extra={"collection_ids": collection_ids},
            )

        # Extract user identity from validated ctx
        user_id = ctx.user.user_id
        client_id = getattr(ctx.user, "client_id", None)

        # Generate conversation_id if not provided (new conversation)
        if not conversation_id:
            conversation_id = str(uuid.uuid4())
            logger.info(f"Created new conversation: {conversation_id}")

        logger.info(
            "Executing RAG chat",
            extra={
                "user_id": user_id,
                "client_id": client_id,
                "conversation_id": conversation_id,
            },
        )

        try:
            # Validate query length (v1.10.2: configurable via UBP_RAG__MAX_QUERY_LENGTH)
            max_length = int(self.config["security"].get("max_query_length", 32000))
            if len(query) > max_length:
                return {
                    "answer": f"Query too long ({len(query)} chars). Maximum {max_length} characters allowed.",
                    "sources": [],
                    "config_used": {},
                    "permissions_checked": [],
                    "conversation_id": conversation_id,
                    "error": "query_too_long",
                }

            # ═══════════════════════════════════════════════════════════════════
            # MAGIC COMMAND: /raw - Pure LLM Bypass Mode
            # ═══════════════════════════════════════════════════════════════════
            # Allows testing pure LLM intelligence by bypassing:
            # - Semantic Router
            # - RAG Retrieval (Qdrant)
            # - Web Search
            # Usage: /raw Chi erano gli Estensi?
            # ═══════════════════════════════════════════════════════════════════
            if query.strip().lower().startswith("/raw"):
                logger.info(
                    "[MAGIC] /raw command detected - Bypassing RAG/Web pipeline",
                    extra={
                        "user_id": user_id,
                        "conversation_id": conversation_id,
                        "original_query": query[:100],
                    },
                )

                # 1. Clean the query (remove "/raw" prefix)
                clean_query = query.strip()[4:].strip()

                if not clean_query:
                    return {
                        "answer": "Usage: /raw <your question>\nExample: /raw Chi erano gli Estensi?",
                        "sources": [],
                        "config_used": {},
                        "permissions_checked": [],
                        "conversation_id": conversation_id,
                        "mode": "pure_llm_forced",
                        "router": {
                            "route": "pure_llm",
                            "confidence": 1.0,
                            "method": "magic_command",
                            "reasoning": "/raw bypass - no query provided",
                        },
                    }

                # 2. Get base config
                base_config = self.config.get("default_rag_config", {})

                # 3. Execute Pure LLM (no retrieval, no web)
                start_time = time.time()

                # Use rag_pipeline.chat with empty collections = Pure LLM mode
                logger.info("[RAG] op=pure_llm_execute mode=pure_llm")
                llm_result = await self.rag_pipeline.chat(
                    query=clean_query,
                    collections=[],  # Empty = Pure LLM mode, no retrieval
                    config={
                        **base_config,
                        "system_prompt": (
                            "You are a helpful AI assistant. "
                            "Answer based on your internal knowledge only. "
                            "Be comprehensive and informative."
                        ),
                    },
                    return_debug=True,
                    web_context=None,  # No web search
                    conversation_context=None,  # Fresh context (optional: could pass history)
                )

                latency_ms = (time.time() - start_time) * 1000

                logger.info(
                    "[MAGIC] /raw command completed",
                    extra={
                        "latency_ms": latency_ms,
                        "answer_length": len(llm_result.get("answer", "")),
                        "conversation_id": conversation_id,
                    },
                )

                return {
                    "answer": llm_result.get("answer", ""),
                    "sources": [],
                    "config_used": llm_result.get("config_used", base_config),
                    "debug": {"pure_llm": True, "latency_ms": latency_ms},
                    "permissions_checked": [],
                    "conversation_id": conversation_id,
                    "mode": "pure_llm_forced",
                    "router": {
                        "route": "pure_llm",
                        "confidence": 1.0,
                        "method": "magic_command",
                        "reasoning": "/raw bypass - direct LLM query",
                    },
                }
            # ═══════════════════════════════════════════════════════════════════

            # ═══════════════════════════════════════════════════════════════════
            # v2.3: INTERACTIVE ANALYST - Check for Active Report Session
            # v2.6: FIX - Enhanced approval detection to bypass enrichment
            # ═══════════════════════════════════════════════════════════════════
            if self.report_session_manager:
                try:
                    # v2.6 FIX: First check if query is a report approval command
                    # This prevents approval queries ("si", "yes", "approve") from
                    # being sent to enrichment when conversation_id doesn't match
                    is_approval = self._is_report_approval_query(query)

                    # Try to find session with provided conversation_id first
                    active_session = (
                        await self.report_session_manager.get_active_session(
                            user_id=user_id,
                            conversation_id=conversation_id,
                        )
                    )

                    # v2.6 FIX: If no session found but query is approval,
                    # search for ANY active session in AWAITING_APPROVAL state
                    if not active_session and is_approval:
                        logger.info(
                            "[REPORT] Approval query detected, searching for any awaiting session",
                            extra={"user_id": user_id, "query": query[:20]},
                        )
                        # Search without conversation_id filter
                        active_session = (
                            await self.report_session_manager.get_active_session(
                                user_id=user_id,
                                conversation_id=None,  # No filter
                            )
                        )
                        if active_session and active_session.state == ReportState.AWAITING_APPROVAL:
                            logger.info(
                                f"[REPORT] Found awaiting session: {active_session.session_id}",
                                extra={
                                    "session_conversation_id": active_session.conversation_id,
                                    "original_query": active_session.metadata.get("original_query", "")[:50],
                                },
                            )

                    if active_session and active_session.state not in [
                        ReportState.COMPLETED,
                        ReportState.CANCELLED,
                        ReportState.IDLE,
                    ]:
                        # User has an active report session - process input
                        logger.info(
                            f"[REPORT] Active session found: {active_session.state.value}",
                            extra={
                                "session_id": active_session.session_id,
                                "conversation_id": conversation_id,
                                "is_approval": is_approval,
                            },
                        )

                        session_result = (
                            await self.report_session_manager.process_input(
                                session_id=active_session.session_id,
                                user_input=query,
                            )
                        )

                        # v2.6: Use session's conversation_id in response
                        response_conversation_id = active_session.conversation_id or conversation_id

                        # Return session response
                        report_session_data = {
                            "session_id": active_session.session_id,
                            "state": session_result.get(
                                "state", active_session.state.value
                            ),
                            "action": session_result.get("action"),
                            "original_query": active_session.metadata.get("original_query"),
                        }
                        # Propagate sections and plan preview when plan is modified
                        if session_result.get("sections"):
                            report_session_data["sections"] = session_result["sections"]
                        if session_result.get("current_plan_preview"):
                            report_session_data["current_plan_preview"] = session_result["current_plan_preview"]

                        return {
                            "answer": session_result.get("message", ""),
                            "sources": [],
                            "config_used": {},
                            "permissions_checked": [],
                            "conversation_id": response_conversation_id,
                            "mode": "report_interactive",
                            "report_session": report_session_data,
                            "router": {
                                "route": "report",
                                "confidence": 1.0,
                                "method": "active_session",
                                "reasoning": "Continuing active report session",
                            },
                        }
                except Exception as e:
                    logger.warning(f"[REPORT] Session check failed: {e}")
            # ═══════════════════════════════════════════════════════════════════

            # === USER EXPLICIT ROUTING OVERRIDE ===
            # If user provides use_rag parameter, it overrides semantic router
            # This ensures user intent is respected and not overridden by keyword matching
            use_rag_override = kwargs.get("use_rag")

            # === FEAT-ROUTER-001: Semantic Query Routing ===
            # Classify query intent to determine optimal handling strategy
            route_result = None

            if use_rag_override is not None:
                # User has explicitly requested RAG or CHAT mode
                # This takes precedence over semantic router classification
                forced_route = RouteType.RAG if use_rag_override else RouteType.CHAT
                route_result = RouterResult(
                    route=forced_route,
                    confidence=1.0,
                    method="user_explicit_override",
                    reasoning=f"User explicitly set use_rag={use_rag_override}",
                )
                logger.info(
                    f"[ROUTING] User explicit override: use_rag={use_rag_override} -> {forced_route.value}",
                    extra={
                        "user_id": user_id,
                        "conversation_id": conversation_id,
                        "use_rag": use_rag_override,
                        "forced_route": forced_route.value,
                    },
                )
            elif self.semantic_router:
                try:
                    # FEAT-DKI-001 (v1.8.1): Load dynamic keywords for user's accessible KBs
                    user_keywords = None
                    if self.keyword_manager and self.acl_manager:
                        try:
                            accessible_kbs = (
                                await self.acl_manager.get_accessible_collections(
                                    user_id=user_id, client_id=client_id
                                )
                            )
                            if accessible_kbs:
                                user_keywords = (
                                    await self.keyword_manager.get_all_keywords_flat(
                                        accessible_kbs
                                    )
                                )
                                if user_keywords:
                                    logger.debug(
                                        f"DKI: Loaded {len(user_keywords)} keywords from {len(accessible_kbs)} KBs",
                                        extra={"kbs": accessible_kbs},
                                    )
                        except Exception as kw_err:
                            logger.warning(
                                f"DKI: Could not load user keywords: {kw_err}"
                            )

                    route_result = await self.semantic_router.classify(
                        query, user_keywords=user_keywords
                    )
                    logger.info(
                        f"Query routed",
                        extra={
                            "route": route_result.route.value,
                            "confidence": route_result.confidence,
                            "method": route_result.method,
                            "conversation_id": conversation_id,
                            "dki_keywords_available": bool(user_keywords),
                        },
                    )
                except Exception as e:
                    logger.warning(
                        f"Router classification failed, defaulting to RAG: {e}"
                    )
                    route_result = RouterResult(
                        route=RouteType.RAG,
                        confidence=0.5,
                        method="fallback",
                        reasoning="Router error, using default",
                    )

            # ═══════════════════════════════════════════════════════════════════
            # FEAT-MEM-003: Memory-Aware Query Rewriting (v5.0)
            # ═══════════════════════════════════════════════════════════════════
            # Pre-compute retrieval_query using cached hints.
            # Works for ALL routes and ALL users.
            # query originale preservata per memory saves e LLM prompt.
            rewrite_result = None
            retrieval_query = query  # default: no rewrite
            if conversation_id and self.memory_module and hasattr(self.memory_module, 'rewrite_query'):
                try:
                    rewrite_result = await self.memory_module.rewrite_query(
                        session_id=conversation_id, query=query, ctx=ctx
                    )
                    if rewrite_result.get("rewrite_type") not in ("none", "error"):
                        retrieval_query = rewrite_result["query"]
                        logger.info(
                            f"[REWRITER] {rewrite_result['rewrite_type']}: "
                            f"'{query[:40]}' -> '{retrieval_query[:60]}'"
                        )
                except Exception as e:
                    logger.warning(f"[REWRITER] Query rewrite failed: {e}")

            # ═══════════════════════════════════════════════════════════════════
            # v2.3: REPORT Route - Start Interactive Report Session (pre-CHAT)
            # ═══════════════════════════════════════════════════════════════════
            if self.report_session_manager and self._is_report_request(query):
                logger.info(
                    "[REPORT] New report request detected",
                    extra={"user_id": user_id, "conversation_id": conversation_id},
                )

                try:
                    # Get user's accessible collections for the report
                    accessible_kbs = []
                    if collections:
                        accessible_kbs = collections
                    elif self.acl_manager:
                        try:
                            accessible_kbs = (
                                await self.acl_manager.get_accessible_collections(
                                    user_id=user_id, client_id=client_id
                                )
                            )
                        except Exception:
                            pass

                    # Start a new report session
                    session = await self.report_session_manager.start_session(
                        user_id=user_id,
                        query=query,
                        conversation_id=conversation_id,
                        client_id=client_id,
                        collections=accessible_kbs,
                    )

                    # Get the proposal from the session
                    proposal = (
                        session.plan.get_proposal_text()
                        if session.plan
                        else "Session started"
                    )

                    return {
                        "answer": proposal,
                        "sources": [],
                        "config_used": {},
                        "permissions_checked": accessible_kbs,
                        "conversation_id": conversation_id,
                        "mode": "report_planning",
                        "report_session": {
                            "session_id": session.session_id,
                            "state": session.state.value,
                            "template_id": session.plan.template_id
                            if session.plan
                            else None,
                            "template_name": session.plan.template_name
                            if session.plan
                            else None,
                            "sections_count": len(session.plan.sections)
                            if session.plan
                            else 0,
                        },
                        "router": {
                            "route": "report",
                            "confidence": 0.9,
                            "method": "pattern_match",
                            "reasoning": "Report request detected, starting interactive session",
                        },
                    }
                except Exception as e:
                    logger.error(f"[REPORT] Failed to start session: {e}")

            # === CHAT Route: Pure conversation, no retrieval needed ===
            if route_result and route_result.route == RouteType.CHAT:
                logger.info(
                    "CHAT route: Pure LLM conversation",
                    extra={"user_id": user_id, "conversation_id": conversation_id},
                )

                # v4.1.0: Centralized summary-first context loading
                conversation_context = await self._load_conversation_context(
                    user_id=user_id,
                    conversation_id=conversation_id,
                    ctx=ctx,
                    route_label="CHAT",
                )

                # Generate response with LLM only (no retrieval)
                # v5.0: Use retrieval_query for pipeline, inject _original_user_query for LLM prompt
                chat_config = self.config.get("default_rag_config", {}).copy()
                # v6.x: Merge override_config into chat_config (system_prompt, temperature, etc.)
                if override_config:
                    for key in ("system_prompt", "temperature", "max_tokens", "top_p"):
                        if key in override_config:
                            chat_config[key] = override_config[key]
                if retrieval_query != query:
                    chat_config["_original_user_query"] = query
                    if rewrite_result and rewrite_result.get("metadata"):
                        chat_config["_rewrite_focus"] = rewrite_result["metadata"].get("current_focus", "")
                logger.info("[RAG] op=architect_chat mode=pure_llm")
                chat_result = await self.rag_pipeline.chat(
                    query=retrieval_query,
                    collections=[],  # Empty = Pure LLM mode
                    config=chat_config,
                    return_debug=True,
                    conversation_context=conversation_context,
                )

                # Save to conversation history
                if self.conversation_manager:
                    try:
                        await self.conversation_manager.add_message(
                            user_id=user_id,
                            conversation_id=conversation_id,
                            role="user",
                            content=query,
                            metadata={"route": "chat"},
                        )
                        await self.conversation_manager.add_message(
                            user_id=user_id,
                            conversation_id=conversation_id,
                            role="assistant",
                            content=chat_result.get("answer", ""),
                            metadata={"route": "chat"},
                        )
                    except Exception as e:
                        logger.warning(f"Could not save chat messages: {e}")

                # v4.1.0: Save to structured memory for eager compression
                if conversation_id and self.memory_module:
                    try:
                        await self.memory_module.add_message(
                            session_id=conversation_id, role="user",
                            content=query, ctx=ctx,
                        )
                        await self.memory_module.add_message(
                            session_id=conversation_id, role="assistant",
                            content=chat_result.get("answer", ""),
                            metadata={"route": "chat"}, ctx=ctx,
                        )
                        self._invalidate_conversation_cache(user_id, conversation_id)
                    except Exception as e:
                        logger.warning(f"Could not save to structured memory: {e}")

                chat_return = {
                    "answer": chat_result.get("answer", ""),
                    "sources": [],
                    "config_used": {},
                    "permissions_checked": [],
                    "conversation_id": conversation_id,
                    "mode": "chat",
                    "mode_reason": "semantic_router",
                    "router": route_result.to_dict() if route_result else None,
                }
                # v5.0: Add rewrite debug
                if rewrite_result:
                    chat_return.setdefault("debug", {})["query_rewrite"] = rewrite_result
                return chat_return

                    # Fall through to normal RAG handling
            # ═══════════════════════════════════════════════════════════════════

            # === WEB Route: Web search + LLM synthesis ===
            if route_result and route_result.route == RouteType.WEB:
                if self.web_search_module:
                    logger.info(
                        "WEB route: Fetching web results",
                        extra={"user_id": user_id, "conversation_id": conversation_id},
                    )

                    try:
                        # Perform web search
                        web_result = await self.web_search_module.search(
                            query=query,
                            max_results=self.config.get("router", {}).get(
                                "web_max_results", 5
                            ),
                        )

                        # Check if web search returned actual results or service unavailable
                        web_results_list = web_result.get("results", [])
                        service_unavailable = len(
                            web_results_list
                        ) == 1 and "Search Service Unavailable" in web_results_list[
                            0
                        ].get("title", "")

                        # Format web results as context
                        web_context_formatted = self._format_web_results(
                            web_results_list
                        )

                        # v4.1.0: Centralized summary-first context loading
                        conversation_context = await self._load_conversation_context(
                            user_id=user_id,
                            conversation_id=conversation_id,
                            ctx=ctx,
                            route_label="WEB",
                        )

                        # Generate response with web context
                        # v5.0: Use retrieval_query for pipeline, inject _original_user_query for LLM prompt
                        web_config = self.config.get("default_rag_config", {}).copy()
                        if retrieval_query != query:
                            web_config["_original_user_query"] = query
                            if rewrite_result and rewrite_result.get("metadata"):
                                web_config["_rewrite_focus"] = rewrite_result["metadata"].get("current_focus", "")
                        logger.info("[RAG] op=architect_web_chat mode=web_context")
                        web_chat_result = await self.rag_pipeline.chat(
                            query=retrieval_query,
                            collections=[],  # No vector search
                            config=web_config,
                            return_debug=True,
                            web_context=web_context_formatted,
                            conversation_context=conversation_context,
                        )

                        # Save to conversation history
                        if self.conversation_manager:
                            try:
                                await self.conversation_manager.add_message(
                                    user_id=user_id,
                                    conversation_id=conversation_id,
                                    role="user",
                                    content=query,
                                    metadata={"route": "web"},
                                )
                                await self.conversation_manager.add_message(
                                    user_id=user_id,
                                    conversation_id=conversation_id,
                                    role="assistant",
                                    content=web_chat_result.get("answer", ""),
                                    metadata={
                                        "route": "web",
                                        "web_sources": len(web_results_list),
                                        "service_unavailable": service_unavailable,
                                    },
                                )
                            except Exception as e:
                                logger.warning(f"Could not save web chat messages: {e}")

                        # v4.1.0: Save to structured memory for eager compression
                        if conversation_id and self.memory_module:
                            try:
                                await self.memory_module.add_message(
                                    session_id=conversation_id, role="user",
                                    content=query, ctx=ctx,
                                )
                                await self.memory_module.add_message(
                                    session_id=conversation_id, role="assistant",
                                    content=web_chat_result.get("answer", ""),
                                    metadata={"route": "web"}, ctx=ctx,
                                )
                                self._invalidate_conversation_cache(user_id, conversation_id)
                            except Exception as e:
                                logger.warning(f"Could not save to structured memory: {e}")

                        web_return = {
                            "answer": web_chat_result.get("answer", ""),
                            "sources": [
                                {
                                    "type": "web",
                                    "title": r.get("title", ""),
                                    "url": r.get("href", r.get("url", "")),
                                }
                                for r in web_results_list[:5]
                            ],
                            "web_results": web_results_list,  # FIX-WEB-001: Include raw web results for observability
                            "web_results_count": len(web_results_list),
                            "config_used": {},
                            "permissions_checked": [],
                            "conversation_id": conversation_id,
                            "mode": "web",
                            "mode_reason": "semantic_router",
                            "service_unavailable": service_unavailable,
                            "router": route_result.to_dict() if route_result else None,
                        }
                        # v5.0: Add rewrite debug
                        if rewrite_result:
                            web_return.setdefault("debug", {})["query_rewrite"] = rewrite_result
                        return web_return

                    except Exception as e:
                        logger.warning(f"Web search failed, falling back to CHAT: {e}")
                        # Fallback to CHAT mode with user notification
                        return await self._handle_web_fallback_to_chat(
                            query=query,
                            user_id=user_id,
                            conversation_id=conversation_id,
                            route_result=route_result,
                            fallback_reason=f"Web search error: {str(e)}",
                            ctx=ctx,
                        )
                else:
                    logger.warning(
                        "WEB route requested but web_search module not available, falling back to CHAT"
                    )
                    # Fallback to CHAT mode with user notification
                    return await self._handle_web_fallback_to_chat(
                        query=query,
                        user_id=user_id,
                        conversation_id=conversation_id,
                        route_result=route_result,
                        fallback_reason="Web search module not available",
                        ctx=ctx,
                    )

            # === RAG Route (default): Vector search + LLM synthesis ===
            # Continue with existing RAG logic below...

            # v4.3.0: Parallel loading of conversation context and RAG config
            # Both operations are independent and can run concurrently
            conversation_context_task = self._load_conversation_context(
                user_id=user_id,
                conversation_id=conversation_id,
                ctx=ctx,
                route_label="RAG",
            )
            config_task = self.config_manager.get_rag_config(
                entity_type="user",
                entity_id=user_id,
                default_config=self.config.get("default_rag_config", {}),
            )

            conversation_context, config_result = await asyncio.gather(
                conversation_context_task, config_task
            )
            rag_config = config_result["config"]

            # Override config if provided
            # FIX-SETTINGS-MERGE v4.1.1: Deep merge for nested enrichment settings
            # Request-level overrides have HIGHEST priority over all other settings
            if override_config:
                # Map flat enrichment keys to nested enrichment.* structure
                enrichment_keys_map = {
                    "hyde_enabled": "hyde_enabled",
                    "investigative_enabled": "investigative_enabled",
                    "reranking_enabled": "rerank_enabled",  # Note: API uses reranking_, config uses rerank_
                    "rerank_enabled": "rerank_enabled",
                    "query_expansion_enabled": "query_expansion_enabled",
                    "fusion_enabled": "fusion_enabled",
                    "dedup_enabled": "dedup_enabled",
                    "compression_enabled": "compression_enabled",
                }

                # Ensure enrichment dict exists
                if "enrichment" not in rag_config:
                    rag_config["enrichment"] = {}

                # Apply overrides with proper nesting
                for key, value in override_config.items():
                    if key in enrichment_keys_map:
                        # Map to nested enrichment structure
                        nested_key = enrichment_keys_map[key]
                        rag_config["enrichment"][nested_key] = value
                        logger.debug(f"[FIX-SETTINGS-MERGE] Mapped {key}={value} → enrichment.{nested_key}")
                    else:
                        # Top-level setting (top_k, max_tokens, etc.)
                        rag_config[key] = value

                logger.info(
                    "[FIX-SETTINGS-MERGE] Applied request-level overrides",
                    extra={
                        "override_keys": list(override_config.keys()),
                        "enrichment_after": rag_config.get("enrichment", {}),
                    }
                )

            # NOTE: pipeline_config is built AFTER collections are determined (4-layer merge)
            # See below after allowed_collections is set

            # Determine collections to query
            # IMPORTANT: Distinguish between:
            #   - collections=None → use all accessible collections (RAG mode), fallback to Pure LLM if none
            #   - collections=[]   → skip retrieval entirely (Pure LLM mode, explicit)
            #   - collections=[...] → use specified collections (RAG mode)
            pure_llm_mode = collections is not None and len(collections) == 0
            mode_reason = None  # Tracks why we're in a particular mode

            if pure_llm_mode:
                # Pure LLM mode: explicit request to skip retrieval
                logger.info(
                    "Pure LLM mode: collections=[] - skipping retrieval (explicit)",
                    extra={"user_id": user_id, "conversation_id": conversation_id},
                )
                allowed_collections = []
                mode_reason = "explicit_pure_llm"
            elif collections:
                target_collections = collections[
                    : self.config["security"]["max_collections_per_query"]
                ]
                # v4.3.0: Parallel ACL checks for performance optimization
                acl_check_tasks = [
                    self.acl_manager.check_access(user_id, client_id, coll)
                    for coll in target_collections
                ]
                acl_results = await asyncio.gather(*acl_check_tasks, return_exceptions=True)
                allowed_collections = [
                    coll for coll, has_access in zip(target_collections, acl_results)
                    if has_access is True  # Explicit True check to handle exceptions
                ]

                if not allowed_collections:
                    return {
                        "answer": "You don't have permission to access the requested knowledge bases.",
                        "sources": [],
                        "config_used": rag_config,
                        "permissions_checked": target_collections,
                        "conversation_id": conversation_id,
                        "mode": "error",
                        "mode_reason": "access_denied",
                        "error": "access_denied",
                    }
            else:
                # Get all accessible collections (default RAG mode)
                accessible = await self.acl_manager.get_accessible_collections(
                    user_id, client_id
                )
                target_collections = accessible

                if not target_collections:
                    # STRICT RAG v2.2.3: No fallback - if RAG mode has no KB access, return error
                    logger.warning(
                        "RAG mode requested but no Knowledge Base is accessible",
                        extra={"user_id": user_id, "conversation_id": conversation_id},
                    )
                    return {
                        "answer": "RAG mode requested but no Knowledge Base is accessible for this user.",
                        "sources": [],
                        "config_used": rag_config,
                        "permissions_checked": [],
                        "conversation_id": conversation_id,
                        "mode": "error",
                        "mode_reason": "no_kb_access",
                        "error": "no_kb_access",
                    }

                # User has accessible collections - proceed with ACL check
                # v4.3.0: Parallel ACL checks for performance optimization
                acl_check_tasks = [
                    self.acl_manager.check_access(user_id, client_id, coll)
                    for coll in target_collections
                ]
                acl_results = await asyncio.gather(*acl_check_tasks, return_exceptions=True)
                allowed_collections = [
                    coll for coll, has_access in zip(target_collections, acl_results)
                    if has_access is True  # Explicit True check to handle exceptions
                ]

                if not allowed_collections:
                    # STRICT RAG v2.2.3: No fallback - if ACL filtered all, return error
                    logger.warning(
                        "ACL check failed for all collections",
                        extra={
                            "user_id": user_id,
                            "conversation_id": conversation_id,
                            "target_collections": target_collections,
                        },
                    )
                    return {
                        "answer": "You don't have permission to access any Knowledge Base.",
                        "sources": [],
                        "config_used": rag_config,
                        "permissions_checked": target_collections,
                        "conversation_id": conversation_id,
                        "mode": "error",
                        "mode_reason": "acl_filtered_all",
                        "error": "acl_filtered_all",
                    }

            # Build enrichment pipeline config with 4-layer merge
            # Now that allowed_collections is determined, we can apply collection-level metadata
            # FIX-PROP-002 v4.3.0: Propagate override_config enrichment keys as pipeline_options
            # so that user/architect overrides reach the 4-layer resolution as Layer 4 (highest priority)
            pipeline_options = kwargs.get("pipeline_options", {}) or {}
            propagated_count = 0
            if override_config:
                for oc_key, po_key in ENRICHMENT_PROPAGATION_KEYS.items():
                    source_val = override_config.get(oc_key)
                    if source_val is not None and po_key not in pipeline_options:
                        pipeline_options[po_key] = bool(source_val)
                        propagated_count += 1
                # Also check nested enrichment dict
                oc_enrichment = override_config.get("enrichment", {})
                if isinstance(oc_enrichment, dict):
                    for oc_key, po_key in ENRICHMENT_PROPAGATION_KEYS.items():
                        source_val = oc_enrichment.get(oc_key)
                        if source_val is not None and po_key not in pipeline_options:
                            pipeline_options[po_key] = bool(source_val)
                            propagated_count += 1
            if propagated_count:
                logger.info(
                    f"[ENRICHMENT_OVERRIDE] Propagated {propagated_count} keys "
                    f"from override_config: {list(pipeline_options.keys())}"
                )
            pipeline_config = await self._build_enrichment_pipeline_config(
                ctx=ctx,
                pipeline_options=pipeline_options,
                target_collections=allowed_collections,
            )

            # v5.0: Inject _original_user_query so _generate() uses natural query for LLM prompt
            if retrieval_query != query:
                rag_config["_original_user_query"] = query
                if rewrite_result and rewrite_result.get("metadata"):
                    rag_config["_rewrite_focus"] = rewrite_result["metadata"].get("current_focus", "")

            # FEAT-TOOL-001: Inject tool settings for RAG route (if enabled)
            # Ensure _settings is populated (lazy-load if needed)
            global _settings
            if _settings is None:
                try:
                    _settings = settings_manager.get_settings()
                except Exception:
                    pass
            tool_search = _settings.tool.search if _settings else None
            if tool_search and tool_search.enabled and tool_search.enabled_rag and allowed_collections:
                rag_config["_tool_settings"] = {
                    "enabled": True,
                    "max_iterations": tool_search.max_iterations,
                    "top_k": tool_search.top_k,
                    "similarity_threshold": tool_search.similarity_threshold,
                    "max_context_expansion_kb": tool_search.max_context_expansion_kb,
                    "timeout_ms": tool_search.timeout_ms,
                    "provider": tool_search.provider_rag,
                    "web_search_available": bool(self.web_search_module),
                }
                rag_config["_tool_collections"] = allowed_collections
                rag_config["_web_module"] = self.web_search_module

            # Execute RAG pipeline with HA fallback (v1.10.1)
            # Uses _chat_with_fallback for automatic retry on provider failure
            result = await self._chat_with_fallback(
                query=retrieval_query,  # v5.0: rewritten for Qdrant retrieval
                collections=allowed_collections,  # Empty list triggers Pure LLM mode in pipeline
                config=rag_config,
                return_debug=True,  # Always return debug info
                web_context=web_context,  # ROADMAP v1.5.0 - FEAT-WEB-001
                conversation_context=conversation_context,  # ROADMAP v1.5.0 - FEAT-MEM-001
                pipeline_config=pipeline_config,
            )

            # Track query
            self.total_queries += 1

            # Add permissions info, conversation_id, and mode indicator
            result["permissions_checked"] = allowed_collections
            result["conversation_id"] = conversation_id

            # Add mode indicator for frontend to display appropriate UI
            result["mode"] = "pure_llm" if pure_llm_mode else "rag"
            if mode_reason:
                result["mode_reason"] = mode_reason

            # FEAT-ROUTER-001: Add router info if available
            if route_result:
                result["router"] = route_result.to_dict()

            # v5.0: Add rewrite debug to RAG result
            if rewrite_result and result.get("debug"):
                result["debug"]["query_rewrite"] = rewrite_result

            # === FEAT-ROUTER-003: Smart Auto-Retry on Empty RAG Results ===
            # Check if RAG returned empty and trigger fallback if needed
            if (
                not pure_llm_mode
                and route_result
                and route_result.route == RouteType.RAG
                and self.semantic_router
            ):
                sources_count = len(result.get("sources", []))
                if self._is_empty_response(result.get("answer", ""), sources_count):
                    logger.info(
                        "RAG returned empty response, attempting auto-retry",
                        extra={
                            "user_id": user_id,
                            "conversation_id": conversation_id,
                            "sources_count": sources_count,
                        },
                    )
                    # Try auto-retry with fallback route
                    retry_result = await self._auto_retry_with_fallback(
                        query=query,
                        original_result=result,
                        original_route=RouteType.RAG,
                        route_result=route_result,
                        user_id=user_id,
                        client_id=client_id,
                        conversation_id=conversation_id,
                        rag_config=rag_config,
                        conversation_context=conversation_context,
                        pipeline_config=pipeline_config,
                    )
                    # If retry succeeded, use that result
                    if (
                        retry_result.get("mode") != "rag"
                        or len(retry_result.get("sources", [])) > 0
                    ):
                        # Sanitize and return retry result
                        if "answer" in retry_result:
                            retry_result["answer"] = SensitiveDataFilter.sanitize(
                                retry_result["answer"]
                            )
                        return retry_result

            # === FEAT-MEM-001: Save messages to conversation history ===
            if self.conversation_manager:
                try:
                    # Save user message
                    await self.conversation_manager.add_message(
                        user_id=user_id,
                        conversation_id=conversation_id,
                        role="user",
                        content=query,
                        metadata={"collections": allowed_collections},
                    )
                    # Save assistant response
                    await self.conversation_manager.add_message(
                        user_id=user_id,
                        conversation_id=conversation_id,
                        role="assistant",
                        content=result.get("answer", ""),
                        metadata={"sources_count": len(result.get("sources", []))},
                    )
                    logger.debug(f"Saved messages to conversation {conversation_id}")
                except Exception as e:
                    logger.warning(f"Could not save messages to conversation: {e}")

            # v4.1.0: Save to structured memory for eager compression
            if conversation_id and self.memory_module:
                try:
                    await self.memory_module.add_message(
                        session_id=conversation_id, role="user",
                        content=query, ctx=ctx,
                    )
                    await self.memory_module.add_message(
                        session_id=conversation_id, role="assistant",
                        content=result.get("answer", ""),
                        metadata={
                            "route": "rag",
                            "config_used": result.get("config_used"),
                        }, ctx=ctx,
                    )
                    self._invalidate_conversation_cache(user_id, conversation_id)
                except Exception as e:
                    logger.warning(f"Could not save to structured memory: {e}")

            logger.info(
                "✅ RAG chat completed",
                extra={
                    "user_id": user_id,
                    "collections": len(allowed_collections),
                    "conversation_id": conversation_id,
                },
            )

            # === SECURITY: Sanitize response to prevent data leakage (VULN-008) ===
            if "answer" in result:
                result["answer"] = SensitiveDataFilter.sanitize(result["answer"])
            if "sources" in result:
                for source in result["sources"]:
                    if "preview" in source:
                        source["preview"] = SensitiveDataFilter.sanitize(
                            source["preview"]
                        )
                    if "text" in source:
                        source["text"] = SensitiveDataFilter.sanitize(source["text"])

            return result

        except Exception as e:
            logger.error(f"Error in RAG chat: {e}")
            return {
                "answer": "I encountered an error while processing your question.",
                "sources": [],
                "config_used": {},
                "permissions_checked": [],
                "conversation_id": conversation_id,
                "mode": "error",
                "mode_reason": "exception",
                "error": str(e),
            }

    # ===== Utility Operations =====

    async def delete_personal_document(
        self,
        document_id: str,
        ctx: Any = None,
        **_,
    ) -> Dict[str, Any]:
        """Delete a document from the current user's Personal KB."""
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        if not document_id or not str(document_id).strip():
            raise HTTPException(status_code=400, detail="Document ID is required")

        user_id = ctx.user.user_id
        personal_kb_name = f"{PERSONAL_KB_PREFIX}{user_id[:8]}"

        if not self.qdrant_module:
            raise HTTPException(status_code=503, detail="Vector store not available")

        try:
            await self.qdrant_module.delete_document_internal(
                doc_id=document_id,
                collection=personal_kb_name,
            )

            logger.info(
                f"[PERSONAL_KB] User {user_id} deleted doc {document_id} from {personal_kb_name}",
                extra={
                    "user_id": user_id,
                    "doc_id": document_id,
                    "collection": personal_kb_name,
                },
            )

            return {
                "success": True,
                "message": "Document deleted successfully",
                "collection": personal_kb_name,
                "deleted_id": document_id,
            }
        except HTTPException:
            raise
        except Exception as e:
            logger.error(
                f"Failed to delete personal document: {e}",
                extra={
                    "user_id": user_id,
                    "doc_id": document_id,
                    "collection": personal_kb_name,
                },
            )
            raise HTTPException(500, f"Error deleting document: {str(e)}")

    async def list_personal_documents(
        self,
        limit: int = 100,
        offset: int = 0,
        ctx: Any = None,
        **_,
    ) -> Dict[str, Any]:
        """List documents in the current user's Personal KB.

        Args:
            limit: Maximum number of documents to return (default 100)
            offset: Number of documents to skip for pagination (default 0)
            ctx: Security context (required)

        Returns:
            Dict with documents list and pagination metadata
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id
        personal_kb_name = f"{PERSONAL_KB_PREFIX}{user_id[:8]}"

        if not self.qdrant_module:
            raise HTTPException(status_code=503, detail="Vector store not available")

        try:
            result = await self.qdrant_module.list_documents(
                collection=personal_kb_name,
                limit=limit,
                offset=offset,
                ctx=ctx,
            )

            logger.info(
                f"[PERSONAL_KB] User {user_id} listed documents from {personal_kb_name}",
                extra={
                    "user_id": user_id,
                    "collection": personal_kb_name,
                    "count": result.get("total", 0),
                },
            )

            return {
                "success": True,
                "documents": result.get("documents", []),
                "total": result.get("total", 0),
                "limit": limit,
                "offset": offset,
                "has_more": result.get("has_more", False),
                "collection": personal_kb_name,
            }

        except PermissionError as e:
            raise HTTPException(status_code=403, detail=str(e))
        except Exception as e:
            # Collection may not exist yet (no documents ingested)
            if "doesn't exist" in str(e).lower() or "not found" in str(e).lower():
                return {
                    "success": True,
                    "documents": [],
                    "total": 0,
                    "limit": limit,
                    "offset": offset,
                    "has_more": False,
                    "collection": personal_kb_name,
                    "message": "Personal KB is empty or not yet created",
                }
            logger.error(
                f"Failed to list personal documents: {e}",
                extra={"user_id": user_id, "collection": personal_kb_name},
            )
            raise HTTPException(500, f"Error listing documents: {str(e)}")

    async def clear_personal_kb(
        self,
        ctx: Any = None,
        **_,
    ) -> Dict[str, Any]:
        """Clear all documents from the current user's Personal KB.

        This is a destructive operation that removes all documents.
        The orchestrator validates ownership before calling qdrant.clear_internal().

        Args:
            ctx: Security context (required)

        Returns:
            Dict with operation result
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id
        personal_kb_name = f"{PERSONAL_KB_PREFIX}{user_id[:8]}"

        if not self.qdrant_module:
            raise HTTPException(status_code=503, detail="Vector store not available")

        try:
            # Use clear_internal which bypasses admin check
            # Ownership is validated here by constructing kb_name from user_id
            result = await self.qdrant_module.clear_internal(
                collection_name=personal_kb_name
            )

            logger.warning(
                f"[AUDIT] User {user_id} cleared personal KB {personal_kb_name}",
                extra={
                    "user_id": user_id,
                    "collection": personal_kb_name,
                    "action": "clear_personal_kb",
                    "audit": True,
                },
            )

            return {
                "success": True,
                "message": "Personal KB cleared successfully",
                "collection": personal_kb_name,
                "user_id": user_id,
            }

        except Exception as e:
            # Collection may not exist (nothing to clear)
            if "doesn't exist" in str(e).lower() or "not found" in str(e).lower():
                return {
                    "success": True,
                    "message": "Personal KB was already empty or not created",
                    "collection": personal_kb_name,
                    "user_id": user_id,
                }
            logger.error(
                f"Failed to clear personal KB: {e}",
                extra={"user_id": user_id, "collection": personal_kb_name},
            )
            raise HTTPException(500, f"Error clearing personal KB: {str(e)}")

    async def list_knowledge_bases(self, ctx: Any = None, **kwargs) -> Dict[str, Any]:
        """List all knowledge bases (filtered by user access)."""
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id
        client_id = getattr(ctx.user, "client_id", None)
        is_admin = self._is_admin(ctx)

        try:
            # Get all collections from Qdrant
            # IMPORTANT: rag_qdrant.list_collections() returns List[str], NOT a dict!
            result = await self.qdrant_module.list_collections(ctx=ctx)

            # Handle both list (correct format) and dict (legacy/error format)
            if isinstance(result, list):
                all_collections = result
            elif isinstance(result, dict):
                # Legacy format or error response
                error = result.get("error")
                if error:
                    return {"collections": [], "count": 0, "error": error}
                all_collections = result.get("collections", [])
            else:
                all_collections = []

            # Exclude system/internal collections from user-facing listing
            all_collections = [c for c in all_collections if c not in SYSTEM_COLLECTIONS]

            # Admins see all, users see only accessible
            if is_admin:
                return {
                    "collections": [{"name": c} for c in all_collections],
                    "count": len(all_collections),
                }
            else:
                acl_tasks = [
                    self.acl_manager.check_access(user_id, client_id, coll_name)
                    for coll_name in all_collections
                ]
                acl_results = await asyncio.gather(*acl_tasks, return_exceptions=True)
                accessible = [
                    {"name": coll_name}
                    for coll_name, has_access in zip(all_collections, acl_results)
                    if has_access is True
                ]

                return {"collections": accessible, "count": len(accessible)}

        except ValueError as ve:
            # Security context issues are expected in some code paths
            logger.debug(f"Knowledge bases listing skipped: {ve}")
            return {"collections": [], "count": 0}

        except Exception as e:
            logger.error(f"Error listing knowledge bases: {e}")
            return {"collections": [], "count": 0, "error": str(e)}

    # ===== Retrieve with ACL (BUG-RETRIEVE-ACL-001) =====

    async def retrieve_with_acl(
        self,
        query_text: str,
        collections: Optional[List[str]] = None,
        top_k: int = 5,
        ctx: Any = None,
        # Parameter aliases for tool-calling compatibility
        query: Optional[str] = None,
        collection: Optional[str] = None,
        collection_name: Optional[str] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """Retrieve documents from knowledge base with ACL enforcement.

        Same ACL check as rag_chat() but retrieval only — no LLM generation.
        Used by reasoning_loop tool 'retrieve' for non-admin users.

        Args:
            query_text: Search query (or 'query' alias from tool schema)
            collections: Target collection names
            top_k: Number of results per collection (default: 5)
            ctx: Security context with user identity

        Returns:
            {chunks: [...], collections_used: [...], total: int}
        """
        # Parameter aliasing: tool schema sends 'query', we accept both
        effective_query = query_text or query or ""
        if not effective_query:
            return {"chunks": [], "collections_used": [], "total": 0,
                    "error": "empty_query"}

        # Resolve collection aliases
        target_collections = collections or []
        if collection:
            target_collections = [collection]
        elif collection_name:
            target_collections = [collection_name]

        # --- ACL enforcement via get_user_collections (official ACL) ---
        try:
            ctx = self._require_ctx(ctx)
        except (ValueError, AttributeError):
            logger.warning("[RETRIEVE-ACL] No valid ctx — access denied")
            return {"chunks": [], "collections_used": [], "total": 0,
                    "error": "no_security_context"}

        user_id = ctx.user.user_id

        # Get authorized collections via official ACL system
        acl_result = await self.get_user_collections(ctx=ctx)
        authorized_collections = acl_result.get("collections", [])
        # Normalize: extract name if dict
        authorized_collections = [
            (c["name"] if isinstance(c, dict) else c)
            for c in authorized_collections
        ]

        logger.info(
            "[RETRIEVE-ACL] user=%s authorized=%s requested=%s",
            user_id, authorized_collections[:5], target_collections,
        )

        if target_collections:
            # Filter requested collections against authorized ones
            allowed = [c for c in target_collections if c in authorized_collections]
        else:
            # No specific collections requested — use all authorized
            allowed = authorized_collections

        if not allowed:
            logger.warning(
                "[RETRIEVE-ACL] No accessible collections: user=%s requested=%s authorized=%s",
                user_id, target_collections, authorized_collections,
            )
            return {"chunks": [], "collections_used": [], "total": 0,
                    "error": "no_accessible_collections",
                    "authorized": authorized_collections,
                    "requested": target_collections}

        # Limit collections per query
        max_cols = self.config.get("security", {}).get(
            "max_collections_per_query", 5
        )
        allowed = allowed[:max_cols]

        # --- Retrieve from each allowed collection via query_internal ---
        if not self.qdrant_module:
            logger.error("[RETRIEVE-ACL] qdrant_module not initialized")
            return {"chunks": [], "collections_used": [], "total": 0,
                    "error": "rag_qdrant_unavailable"}

        all_chunks = []
        collections_used = []

        for coll in allowed:
            try:
                result = await self.qdrant_module.query_internal(
                    query_text=effective_query,
                    top_k=top_k,
                    collection=coll,
                )
                hits = result.get("results", [])
                for hit in hits:
                    hit["collection"] = coll
                all_chunks.extend(hits)
                if hits:
                    collections_used.append(coll)
            except Exception as e:
                logger.warning(
                    "[RETRIEVE-ACL] Query failed for collection %s: %s", coll, e
                )

        # Sort by score descending
        all_chunks.sort(key=lambda c: c.get("score", 0), reverse=True)

        # Trim to top_k total
        all_chunks = all_chunks[:top_k]

        # Strip internal fields (embedding vectors, query metadata)
        all_chunks = [
            {k: v for k, v in chunk.items() if k not in _INTERNAL_ONLY_KEYS}
            for chunk in all_chunks
        ]

        logger.info(
            "[RETRIEVE-ACL] user=%s collections_allowed=%s chunks=%d",
            user_id, allowed, len(all_chunks),
        )

        return {
            "chunks": all_chunks,
            "collections_used": collections_used,
            "total": len(all_chunks),
        }

    async def retrieve_enriched(
        self,
        query_text: str,
        collections: Optional[List[str]] = None,
        top_k: int = 5,
        ctx: Any = None,
        # Parameter aliases for tool-calling compatibility
        query: Optional[str] = None,
        collection: Optional[str] = None,
        collection_name: Optional[str] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """Retrieve documents with canonical UBP enrichment (no LLM generation).

        Uses the same config resolution and enrichment pipeline as rag_chat(),
        but returns enriched chunks without generating an LLM answer.

        Config resolution: identical 4-layer merge as rag_chat
            1. System Defaults (.env)
            2. Client Policy (client_config)
            3. Collection Metadata (redis)
            4. (No user overrides in v1)

        ACL enforcement: identical to retrieve_with_acl via get_user_collections().

        Args:
            query_text: Search query (or 'query' alias from tool schema)
            collections: Target collection names (ACL-filtered)
            top_k: Number of enriched results to return (1-20, default: 5)
            ctx: Security context with user identity

        Returns:
            {chunks, collections_used, total, enrichment_applied, warnings}
        """
        warnings: List[str] = []

        # --- Parameter aliasing (same pattern as retrieve_with_acl) ---
        effective_query = query_text or query or ""
        if not effective_query:
            return {
                "chunks": [], "collections_used": [], "total": 0,
                "enrichment_applied": [], "warnings": ["empty_query"],
            }

        target_collections = collections or []
        if collection:
            target_collections = [collection]
        elif collection_name:
            target_collections = [collection_name]

        # --- ACL enforcement (reuse retrieve_with_acl pattern) ---
        try:
            ctx = self._require_ctx(ctx)
        except (ValueError, AttributeError):
            logger.warning("[RETRIEVE-ENRICHED] No valid ctx — access denied")
            return {
                "chunks": [], "collections_used": [], "total": 0,
                "enrichment_applied": [], "warnings": ["no_security_context"],
            }

        user_id = ctx.user.user_id

        acl_result = await self.get_user_collections(ctx=ctx)
        authorized_collections = [
            (c["name"] if isinstance(c, dict) else c)
            for c in acl_result.get("collections", [])
        ]

        if target_collections:
            allowed = [c for c in target_collections if c in authorized_collections]
        else:
            allowed = authorized_collections

        if not allowed:
            logger.warning(
                "[RETRIEVE-ENRICHED] No accessible collections: user=%s", user_id,
            )
            return {
                "chunks": [], "collections_used": [], "total": 0,
                "enrichment_applied": [], "warnings": ["no_accessible_collections"],
            }

        max_cols = self.config.get("security", {}).get(
            "max_collections_per_query", 5
        )
        allowed = allowed[:max_cols]

        # --- Config resolution: same canonical path as rag_chat ---
        # Step 1: Base RAG config via config_manager (user/client/default merge)
        try:
            config_result = await self.config_manager.get_rag_config(
                entity_type="user",
                entity_id=user_id,
                default_config=self.config.get("default_rag_config", {}),
            )
            rag_config = config_result["config"]
        except Exception as e:
            logger.warning("[RETRIEVE-ENRICHED] Config resolution failed: %s", e)
            rag_config = self.config.get("default_rag_config", {})
            warnings.append("config_resolution_fallback")

        rag_config["top_k"] = min(max(top_k, 1), 20)

        # Step 2: 4-layer enrichment pipeline config (same as rag_chat L4790)
        # v1: no user overrides (pipeline_options={})
        try:
            pipeline_config = await self._build_enrichment_pipeline_config(
                ctx=ctx,
                pipeline_options={},
                target_collections=allowed,
            )
        except Exception as e:
            logger.warning("[RETRIEVE-ENRICHED] Pipeline config build failed: %s", e)
            pipeline_config = None
            warnings.append("pipeline_config_fallback")

        # --- Retrieval + enrichment via _retrieve() (same as rag_chat) ---
        if not self.rag_pipeline:
            logger.error("[RETRIEVE-ENRICHED] rag_pipeline not initialized")
            return {
                "chunks": [], "collections_used": [], "total": 0,
                "enrichment_applied": [], "warnings": ["rag_pipeline_unavailable"],
            }

        try:
            retrieve_result = await self.rag_pipeline._retrieve(
                query=effective_query,
                collections=allowed,
                config=rag_config,
                pipeline_config=pipeline_config,
            )
        except Exception as e:
            logger.error("[RETRIEVE-ENRICHED] _retrieve() failed: %s", e)
            return {
                "chunks": [], "collections_used": allowed, "total": 0,
                "enrichment_applied": [], "warnings": [f"retrieval_failed: {e}"],
            }

        # --- Output mapping: flatten to v1 public contract ---
        enriched_docs = retrieve_result.get("docs", [])

        enriched_docs = [
            {k: v for k, v in doc.items() if k not in _INTERNAL_ONLY_KEYS}
            for doc in enriched_docs
        ]

        # Build enrichment_applied from enrichment_flags
        enrichment_applied = []
        flags = retrieve_result.get("enrichment_flags", {})
        if flags:
            for step_name in ["rerank", "fusion", "dedup", "compression",
                              "hyde", "query_expansion", "investigative"]:
                key = f"{step_name}_enabled"
                if flags.get(key, False):
                    enrichment_applied.append(step_name)

        # Determine collections_used from chunk metadata
        collections_used = list({
            doc.get("collection", doc.get("metadata", {}).get("collection", ""))
            for doc in enriched_docs
            if doc.get("collection") or doc.get("metadata", {}).get("collection")
        })

        logger.info(
            "[RETRIEVE-ENRICHED] user=%s collections=%s chunks=%d enrichment=%s",
            user_id, allowed, len(enriched_docs), enrichment_applied,
        )

        return {
            "chunks": enriched_docs,
            "collections_used": collections_used,
            "total": len(enriched_docs),
            "enrichment_applied": enrichment_applied,
            "warnings": warnings,
        }

    async def get_user_collections(self, ctx: Any = None, **kwargs) -> Dict[str, Any]:
        """Get collections accessible by user via official ACL system."""
        try:
            ctx = self._require_ctx(ctx)
        except ValueError as e:
            logger.warning(f"[get_user_collections] No valid ctx, returning empty: {e}")
            return {"collections": [], "count": 0}

        # Admin path: full listing via Qdrant (unchanged)
        if self._is_admin(ctx):
            result = await self.list_knowledge_bases(ctx=ctx)
            collections = [kb["name"] for kb in result.get("collections", [])]
            return {"collections": collections, "count": len(collections)}

        # Non-admin: use admin_clients.get_user_kb_assignments (official ACL)
        user_id = ctx.user.user_id
        client_id = getattr(ctx.user, "client_id", None)

        if not client_id:
            logger.warning(f"[get_user_collections] User {user_id} has no client_id")
            return {"collections": [], "count": 0}

        try:
            admin_clients = await self.di_container.resolve("admin_clients")
            assignments = await admin_clients.get_user_kb_assignments(
                client_id=client_id, user_id=user_id, ctx=ctx,
            )
            if not assignments.get("success"):
                logger.warning(f"[get_user_collections] KB assignments failed: {assignments.get('error')}")
                return {"collections": [], "count": 0}

            # effective_kbs: list of {"kb_name", "access_level", "source"}
            collections = [
                kb["kb_name"] for kb in assignments.get("effective_kbs", [])
                if kb.get("access_level") in ("read", "write", "admin")
            ]

            # Add personal KB if enabled
            if assignments.get("personal_kb_enabled", False):
                personal_kb = f"personal_{user_id[:8]}"
                if personal_kb not in collections:
                    collections.append(personal_kb)

            return {"collections": collections, "count": len(collections)}
        except Exception as e:
            logger.error(f"[get_user_collections] Error for user {user_id}: {e}")
            return {"collections": [], "count": 0}

    async def delete_knowledge_base(
        self,
        name: str = None,
        collection_id: str = None,  # Backward compatibility alias (deprecated)
        force: bool = False,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Delete a knowledge base completely.

        Safety: Deletion is blocked for populated collections with matching dimensions
        unless `force=True`. Allowed automatically when KB is empty or has a
        dimension mismatch (embedding model changed). This prevents accidental
        data loss while allowing legitimate rebuild scenarios.

        Scope:
        - Deletes Qdrant collection
        - Removes all ACL entries for collection
        - Removes collection from all clients' kb_config.universal_kbs_assigned

        Security: Admin only.

        Parameters:
            name: Collection name (preferred)
            collection_id: Deprecated alias for 'name' (backward compatibility)
            force: Skip rebuild safety check (default False)

        Note: Either 'name' or 'collection_id' must be provided.
              'collection_id' is deprecated and will be removed in v2.0.
        """
        ctx = self._require_admin(ctx, "delete_knowledge_base")
        self._require_initialized()

        # ENTERPRISE FIX: Resolve name with backward compatibility
        # Priority: name > collection_id > kwargs.get("collection_name")
        resolved_name = name or collection_id or kwargs.get("collection_name")

        # Log deprecation warning if using old parameter
        if collection_id and not name:
            logger.warning(
                "DEPRECATION WARNING: 'collection_id' parameter in delete_knowledge_base "
                "is deprecated. Use 'name' instead. Will be removed in v2.0.",
                extra={"collection_id": collection_id},
            )
        if kwargs.get("collection_name") and not name and not collection_id:
            logger.warning(
                "DEPRECATION WARNING: 'collection_name' parameter in delete_knowledge_base "
                "is deprecated. Use 'name' instead. Will be removed in v2.0.",
                extra={"collection_name": kwargs.get("collection_name")},
            )

        if not resolved_name or not str(resolved_name).strip():
            return {
                "status": "error",
                "message": "Knowledge base name is required. Use 'name' parameter (preferred) or 'collection_id' (deprecated).",
            }

        kb_name = str(resolved_name).strip()

        # Block deletion of system/internal collections
        if kb_name in SYSTEM_COLLECTIONS:
            return {
                "status": "error",
                "message": f"Cannot delete system collection '{kb_name}'. System collections are protected.",
            }

        # =====================================================================
        # KB DELETE SAFETY CHECK (v2.1.0)
        # =====================================================================
        # Prevents accidental deletion of ANY populated, healthy KB.
        # Deletion is allowed ONLY when:
        #   1. KB is empty (0 chunks) — safe to recreate
        #   2. Dimension mismatch — embedding model changed, rebuild required
        #   3. Collection not found — nothing to protect
        #   4. force=True — explicit admin override
        # Uses centralized check_rebuild_needed from rag_qdrant.
        if not force and self.qdrant_module:
            try:
                rebuild_check = await self.qdrant_module.check_rebuild_needed(kb_name)

                if not rebuild_check.get("allow_delete", False):
                    details = rebuild_check.get("details", {})
                    logger.warning(
                        f"KB delete blocked for '{kb_name}': KB has {details.get('current_chunks', '?')} chunks "
                        f"and dimensions match (dim={details.get('current_dim', '?')}). Use force=True to override."
                    )
                    return {
                        "status": "error",
                        "message": (
                            f"Knowledge Base '{kb_name}' contains {details.get('current_chunks', 'N/A')} chunks "
                            f"with matching dimensions. Deletion is only allowed when the KB is empty "
                            f"or has a dimension mismatch. Pass force=true to override."
                        ),
                        "error_code": "KB_DELETE_BLOCKED",
                        "collection": kb_name,
                        "rebuild_check": rebuild_check,
                    }

                logger.info(
                    f"KB delete allowed for '{kb_name}' (reason: {rebuild_check.get('reason', 'unknown')})"
                )
            except Exception as check_err:
                logger.warning(
                    f"KB delete safety check failed for '{kb_name}': {check_err}. Proceeding with delete."
                )

        details: Dict[str, Any] = {
            "collection": kb_name,
            "qdrant_deleted": False,
            "acl_keys_deleted": 0,
            "clients_updated": 0,
        }

        try:
            # 1) Delete Qdrant collection
            if not self.qdrant_module:
                raise RuntimeError("rag_qdrant module not initialized")
            qdrant_result = await self.qdrant_module.delete_collection(
                collection_name=kb_name,
                ctx=ctx,  # Pass admin context for auth
            )
            # rag_qdrant returns {status: deleted|error, ...}
            if isinstance(qdrant_result, dict) and qdrant_result.get("status") in {
                "deleted",
                "success",
            }:
                details["qdrant_deleted"] = True
            elif (
                isinstance(qdrant_result, dict)
                and qdrant_result.get("status") == "not_found"
            ):
                # Treat missing collection as already deleted
                details["qdrant_deleted"] = True
            else:
                err = (
                    qdrant_result.get("error")
                    if isinstance(qdrant_result, dict)
                    else None
                )
                raise RuntimeError(err or "Failed to delete Qdrant collection")

            # 2) Delete ACL keys for this collection (users + clients)
            if not self.acl_manager:
                raise RuntimeError("ACL manager not initialized")
            acl_pattern = self.acl_manager._build_acl_pattern(collection_id=kb_name)
            deleted_acl = 0
            if not self.redis_client:
                raise RuntimeError("Redis client not initialized")
            async for key in self.redis_client.scan_iter(match=acl_pattern):
                await self.redis_client.delete(key)
                deleted_acl += 1
            details["acl_keys_deleted"] = deleted_acl

            # 3) Remove from client assignments
            if not self.di_container:
                raise RuntimeError("DI container not available")
            admin_clients = await self.di_container.resolve("admin_clients")
            if not admin_clients:
                raise RuntimeError("admin_clients module not found in DI container")

            clients = await admin_clients.list_clients(
                filter=None, limit=None, offset=None, ctx=ctx
            )
            updated = 0
            for client in clients:
                kb_config = client.get("kb_config") or {}
                assigned = kb_config.get("universal_kbs_assigned") or []
                if isinstance(assigned, list) and kb_name in assigned:
                    new_assigned = [x for x in assigned if x != kb_name]
                    await admin_clients.update_client(
                        client_id=client["client_id"],
                        kb_config={"universal_kbs_assigned": new_assigned},
                        ctx=ctx,
                    )
                    updated += 1
            details["clients_updated"] = updated

            # 4) DKI cleanup: remove keywords and summary from Redis
            dki_keywords_cleared = False
            dki_summary_cleared = False
            try:
                dki_keywords_cleared = await self.keyword_manager.clear_keywords(kb_name)
                summary_key = self.keyword_manager._build_summary_key(kb_name)
                dki_summary_cleared = bool(await self.redis_client.delete(summary_key))
            except Exception as dki_err:
                logger.warning(f"DKI cleanup failed for '{kb_name}' (non-blocking): {dki_err}")
            details["dki_keywords_cleared"] = dki_keywords_cleared
            details["dki_summary_cleared"] = dki_summary_cleared

            logger.info(
                "✅ Knowledge base deleted",
                extra={
                    "collection": kb_name,
                    "admin_user_id": ctx.user.user_id,
                    **details,
                },
            )

            return {
                "status": "success",
                "message": f"Knowledge base '{kb_name}' deleted successfully",
                "details": details,
            }

        except Exception as e:
            logger.error(
                "Error deleting knowledge base",
                extra={
                    "collection": kb_name,
                    "admin_user_id": getattr(ctx.user, "user_id", None),
                    "error": str(e),
                },
            )
            return {
                "status": "error",
                "message": f"Error deleting knowledge base: {str(e)}",
                "details": details,
            }

    async def get_stats(self, ctx: Any = None, **kwargs) -> Dict[str, Any]:
        """
        Get module statistics.

        Security: Admin only operation.

        Args:
            ctx: Security context

        Returns:
            Module statistics
        """
        ctx = self._require_admin(ctx, "get_stats")
        self._require_initialized()

        try:
            # Get collection count
            kb_result = await self.list_knowledge_bases(ctx)
            total_kb = kb_result.get("count", 0)

            # Get permission count
            if not self.acl_manager:
                raise RuntimeError("ACL manager not initialized")
            perms_result = await self.acl_manager.get_permissions()
            total_perms = perms_result.get("count", 0)

            # Get config count
            prefix = self.config["rag_config_storage"]["redis_key_prefix"]
            config_count = 0
            if self.redis_client:
                async for _ in self.redis_client.scan_iter(match=f"{prefix}:*"):
                    config_count += 1

            return {
                "module": self.manifest.name,
                "total_knowledge_bases": total_kb,
                "total_documents": self.total_documents_ingested,
                "total_permissions": total_perms,
                "total_configs": config_count,
                "total_queries": self.total_queries,
                "dependencies_loaded": all(
                    [
                        self.qdrant_module is not None,
                        self.llm_module is not None,
                        self.redis_client is not None,
                    ]
                ),
            }

        except Exception as e:
            logger.error(f"Error getting stats: {e}")
            return {
                "module": self.manifest.name,
                "total_knowledge_bases": 0,
                "total_documents": 0,
                "total_permissions": 0,
                "total_configs": 0,
                "total_queries": self.total_queries,
                "dependencies_loaded": False,
                "error": str(e),
            }

    # ===== Hot-Reload Configuration (FEAT-HOT-RELOAD v1.9.x) =====

    async def reload_router_config(self, ctx: Any = None, **kwargs) -> Dict[str, Any]:
        """
        Hot-reload router configuration without restart.

        Reloads:
        1. Router weights from router_weights.yaml
        2. RouterConfig from environment variables (if changed)
        3. Updates semantic_router instance

        Security: Admin only operation.

        Args:
            ctx: Security context

        Returns:
            Dict with reload status and updated configuration
        """
        ctx = self._require_admin(ctx, "reload_router_config")
        self._require_initialized()

        reload_results = {
            "weights_reloaded": False,
            "config_reloaded": False,
            "errors": [],
        }

        try:
            # 1. Reload router weights from YAML
            from .router.models import routing_weights

            old_weights = {
                "commercial_fresh_boost": routing_weights.commercial_fresh_boost,
                "unknown_query_fallback": routing_weights.unknown_query_fallback,
                "confidence_threshold": routing_weights.hard_flag_confidence,
            }

            routing_weights.reload()

            new_weights = {
                "commercial_fresh_boost": routing_weights.commercial_fresh_boost,
                "unknown_query_fallback": routing_weights.unknown_query_fallback,
                "confidence_threshold": routing_weights.hard_flag_confidence,
            }

            reload_results["weights_reloaded"] = True
            reload_results["weights_changed"] = old_weights != new_weights
            reload_results["weights_before"] = old_weights
            reload_results["weights_after"] = new_weights

            logger.info(
                "Router weights reloaded",
                extra={"changed": old_weights != new_weights},
            )

        except Exception as e:
            reload_results["errors"].append(f"Weights reload failed: {str(e)}")
            logger.error(f"Failed to reload router weights: {e}")

        try:
            # 2. Rebuild RouterConfig from current environment
            new_router_config = _build_router_config_from_settings()

            if self.semantic_router:
                old_config = {
                    "confidence_threshold": self.semantic_router.router_config.confidence_threshold,
                    "fallback_strategy": self.semantic_router.router_config.fallback_strategy.value,
                    "enable_auto_retry": self.semantic_router.router_config.enable_auto_retry,
                    "empty_rag_triggers_web": self.semantic_router.router_config.empty_rag_triggers_web,
                }

                # Update router config
                self.semantic_router.router_config = new_router_config

                new_config = {
                    "confidence_threshold": new_router_config.confidence_threshold,
                    "fallback_strategy": new_router_config.fallback_strategy.value,
                    "enable_auto_retry": new_router_config.enable_auto_retry,
                    "empty_rag_triggers_web": new_router_config.empty_rag_triggers_web,
                }

                reload_results["config_reloaded"] = True
                reload_results["config_changed"] = old_config != new_config
                reload_results["config_before"] = old_config
                reload_results["config_after"] = new_config

                logger.info(
                    "Router config reloaded",
                    extra={"changed": old_config != new_config},
                )

        except Exception as e:
            reload_results["errors"].append(f"Config reload failed: {str(e)}")
            logger.error(f"Failed to reload router config: {e}")

        # Get router stats after reload
        if self.semantic_router:
            reload_results["router_stats"] = self.semantic_router.get_stats()

        return {
            "status": "success" if not reload_results["errors"] else "partial",
            "module": self.manifest.name,
            "operation": "reload_router_config",
            **reload_results,
        }

    # ===== Conversation Management (FEAT-MEM-001 - Task #15) =====

    async def list_conversations(
        self,
        limit: int = 50,
        offset: int = 0,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        List all conversations for the current user.

        Security: Users see only their own conversations.

        Args:
            limit: Maximum conversations to return (default 50)
            offset: Pagination offset (default 0)
            ctx: Security context

        Returns:
            Dict with conversations array, count, and total
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id

        if not self.conversation_manager:
            return {
                "conversations": [],
                "count": 0,
                "total": 0,
                "error": "Conversation manager not initialized",
            }

        logger.info(
            f"Listing conversations for user {user_id}",
            extra={"user_id": user_id, "limit": limit, "offset": offset},
        )

        return await self.conversation_manager.list_conversations(
            user_id=user_id, limit=limit, offset=offset
        )

    async def get_conversation(
        self,
        conversation_id: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Get a specific conversation with all messages.

        Security: Users can only access their own conversations.

        Args:
            conversation_id: Conversation UUID
            ctx: Security context

        Returns:
            Dict with conversation data including messages
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id

        if not conversation_id:
            return {"status": "error", "message": "conversation_id is required"}

        if not self.conversation_manager:
            return {
                "status": "error",
                "message": "Conversation manager not initialized",
            }

        logger.info(
            f"Getting conversation {conversation_id} for user {user_id}",
            extra={"user_id": user_id, "conversation_id": conversation_id},
        )

        return await self.conversation_manager.get_conversation(
            user_id=user_id, conversation_id=conversation_id
        )

    async def delete_conversation(
        self,
        conversation_id: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Delete a conversation.

        Security: Users can only delete their own conversations.

        Args:
            conversation_id: Conversation UUID to delete
            ctx: Security context

        Returns:
            Dict with status
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id

        if not conversation_id:
            return {"status": "error", "message": "conversation_id is required"}

        if not self.conversation_manager:
            return {
                "status": "error",
                "message": "Conversation manager not initialized",
            }

        logger.info(
            f"Deleting conversation {conversation_id} for user {user_id}",
            extra={"user_id": user_id, "conversation_id": conversation_id},
        )

        return await self.conversation_manager.delete_conversation(
            user_id=user_id, conversation_id=conversation_id
        )

    async def create_conversation(
        self,
        title: Optional[str] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Create a new conversation explicitly.

        Note: Conversations are usually auto-created on first message via rag_chat,
        but this method allows explicit creation with a custom title.

        Args:
            title: Optional conversation title
            ctx: Security context

        Returns:
            Dict with conversation_id, title, created_at, status
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id

        if not self.conversation_manager:
            return {
                "status": "error",
                "message": "Conversation manager not initialized",
            }

        logger.info(
            f"Creating new conversation for user {user_id}",
            extra={"user_id": user_id, "title": title},
        )

        return await self.conversation_manager.create_conversation(
            user_id=user_id, title=title
        )

    async def update_conversation_title(
        self,
        conversation_id: str,
        title: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Update the title of a conversation.

        Args:
            conversation_id: Conversation UUID
            title: New title
            ctx: Security context

        Returns:
            Dict with status
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        user_id = ctx.user.user_id

        if not conversation_id or not title:
            return {
                "status": "error",
                "message": "conversation_id and title are required",
            }

        if not self.conversation_manager:
            return {
                "status": "error",
                "message": "Conversation manager not initialized",
            }

        logger.info(
            f"Updating title for conversation {conversation_id}",
            extra={"user_id": user_id, "conversation_id": conversation_id},
        )

        return await self.conversation_manager.update_conversation_title(
            user_id=user_id, conversation_id=conversation_id, title=title
        )

    # =========================================================================
    # SYSTEM ARCHITECT AGENT (v2.0.0)
    # =========================================================================
    # AI-powered Lead Architect using RAG on official documentation.
    # Leverages Grok-4 with 2M token context for comprehensive analysis.

    async def ask_architect(
        self,
        query: str,
        conversation_id: Optional[str] = None,
        max_tokens: Optional[int] = None,
        top_k: Optional[int] = None,
        hyde_enabled: Optional[bool] = None,
        investigative_enabled: Optional[bool] = None,
        reranking_enabled: Optional[bool] = None,
        adaptive_budget_enabled: Optional[bool] = None,
        query_expansion_enabled: Optional[bool] = None,
        fusion_enabled: Optional[bool] = None,
        dedup_enabled: Optional[bool] = None,
        compression_enabled: Optional[bool] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Ask the System Architect Agent a question.

        The Architect Agent uses RAG on official UBP documentation to provide
        authoritative answers. Leverages Grok-4 with 2M token context.

        Args:
            query: The question to ask the architect
            conversation_id: Optional conversation ID for context continuity
            max_tokens: Optional override for response max tokens (default from config: 8000)
            top_k: Optional override for retrieval top_k (default from config: 100)
            hyde_enabled: Optional override for HyDE (Hypothetical Document Embeddings)
            investigative_enabled: Optional override for Investigative Query Decomposition (v2.2.2)
            reranking_enabled: Optional override for reranking (default: True)
            adaptive_budget_enabled: Optional override for adaptive budget manager (v3.5.1)
                                     If not provided, uses UBP_ADAPTIVE_BUDGET__ENABLED from env
            query_expansion_enabled: Optional override for Query Expansion (v4.0.0)
            fusion_enabled: Optional override for Chunk Fusion (v4.0.0)
            dedup_enabled: Optional override for Deduplication (v4.0.0)
            compression_enabled: Optional override for Context Compression (v4.0.0)
            ctx: Security context (required)

        Returns:
            Dict with:
                - answer: The architect's response
                - sources: Referenced documentation
                - model_used: The model that generated the response
                - config_used: Configuration parameters used
                - conversation_id: Conversation ID for follow-up questions

        Raises:
            HTTPException 404: If Architect Agent is disabled
            HTTPException 503: If knowledge base not available
        """
        ctx = self._require_ctx(ctx)
        self._require_initialized()

        # Import settings for Architect configuration
        try:
            from ubp_enterprise_hybrid.backend.app.api.admin_settings_routes import settings_manager
            settings = settings_manager.get_settings()
            architect_settings = settings.architect
        except ImportError:
            logger.error("Could not import settings_manager for Architect Agent")
            return {
                "status": "error",
                "message": "Architect Agent configuration not available",
                "error_code": "CONFIG_ERROR",
            }

        # Check if Architect is enabled
        if not architect_settings.enabled:
            logger.warning("Architect Agent is disabled")
            return {
                "status": "error",
                "message": "Architect Agent is disabled. Enable via UBP_ARCHITECT__ENABLED=true",
                "error_code": "ARCHITECT_DISABLED",
            }

        # Import the system prompt
        try:
            from .prompts.architect import (
                ARCHITECT_SYSTEM_PROMPT,
                build_architect_system_prompt,
            )
        except ImportError:
            logger.error("Could not import Architect system prompt")
            ARCHITECT_SYSTEM_PROMPT = "You are the Lead System Architect. Answer based on the documentation provided."
            build_architect_system_prompt = None

        user_id = ctx.user.user_id
        client_id = getattr(ctx.user, "client_id", None)

        # Generate conversation_id if not provided
        if not conversation_id:
            conversation_id = str(uuid.uuid4())
            logger.info(f"Architect: Created new conversation: {conversation_id}")

        logger.info(
            "Architect Agent query",
            extra={
                "user_id": user_id,
                "client_id": client_id,
                "conversation_id": conversation_id,
                "query_length": len(query),
                "kb_name": architect_settings.kb_name,
                "provider": architect_settings.provider,
            },
        )

        try:
            start_time = time.time()
            from ubp_enterprise_hybrid.modules.cores._shared import ProviderMapper  # v4.2.0: For provider validation

            # Get Architect-specific config
            # v2.1.0: Added context_limit_tokens and chars_per_token for dynamic context budget
            # v2.2.2: Added max_tokens for response length control (with client override)
            # v2.2.2: Added HyDE, Investigative, Reranking support with client override
            effective_hyde = (
                hyde_enabled
                if hyde_enabled is not None
                else architect_settings.hyde_enabled
            )
            effective_investigative = (
                investigative_enabled
                if investigative_enabled is not None
                else architect_settings.investigative_enabled
            )
            effective_reranking = (
                reranking_enabled
                if reranking_enabled is not None
                else architect_settings.reranking_enabled
            )
            # v3.5.1: Adaptive budget manager support for Architect
            # If not provided, use ENV setting (UBP_ADAPTIVE_BUDGET__ENABLED)
            effective_adaptive_budget = (
                adaptive_budget_enabled
                if adaptive_budget_enabled is not None
                else settings.adaptive_budget.enabled
            )
            # v4.0.0: Additional enrichment feature overrides
            effective_query_expansion = (
                query_expansion_enabled
                if query_expansion_enabled is not None
                else architect_settings.query_expansion_enabled
            )
            effective_fusion = (
                fusion_enabled
                if fusion_enabled is not None
                else architect_settings.fusion_enabled
            )
            effective_dedup = (
                dedup_enabled
                if dedup_enabled is not None
                else architect_settings.dedup_enabled
            )
            effective_compression = (
                compression_enabled
                if compression_enabled is not None
                else architect_settings.compression_enabled
            )
            # v6.8.x: empty default → resolve via ProviderMapper (no grok hardcode)
            architect_provider = architect_settings.provider or ""
            if not architect_provider:
                try:
                    from ubp_enterprise_hybrid.modules.cores._shared import ProviderMapper as _PM
                    chain = _PM.resolve_chain("enrichment")
                    if chain:
                        architect_provider = chain[0][1]
                except Exception:
                    pass

            tool_search = settings.tool.search
            tools_enabled = bool(tool_search.enabled and tool_search.enabled_architect)
            if build_architect_system_prompt:
                system_prompt = build_architect_system_prompt(
                    tools_enabled, tool_search.max_iterations
                )
            else:
                system_prompt = ARCHITECT_SYSTEM_PROMPT
            
            architect_config = {
                # v3.7.1: FIX-BUDGET-002 - provider MUST be in config for budget calculation
                # Without this, AdaptiveBudgetManager falls back to 4096 tokens (conservative default)
                # which causes "Document budget very tight" and drops ALL chunks
                "provider": architect_provider,
                "top_k": top_k if top_k is not None else architect_settings.top_k,
                "reranking_enabled": effective_reranking,
                "threshold": architect_settings.relevance_threshold,
                "system_prompt": system_prompt,
                # v2.1.0: Enable dynamic context budget calculation
                "context_limit_tokens": architect_settings.context_limit_tokens,
                "chars_per_token": architect_settings.chars_per_token,
                # v2.2.2: Max tokens for response (default 8000, client override supported)
                "max_tokens": max_tokens
                if max_tokens is not None
                else architect_settings.response_max_tokens,
                # v2.2.2: HyDE (Hypothetical Document Embeddings) for improved retrieval
                # v2.2.2: Investigative Query Decomposition (FEAT-INVEST-001)
                "enrichment": {
                    "enabled": True,
                    "hyde_enabled": effective_hyde,
                    "hyde_document_type": architect_settings.hyde_document_type,
                    "rerank_enabled": effective_reranking,
                    # v2.2.2: Investigative strategy with client override
                    "investigative_enabled": effective_investigative,
                    "investigative_num_questions": settings.enrichment.investigative_num_questions,
                    # v4.0.0: Additional enrichment features
                    "query_expansion_enabled": effective_query_expansion,
                    "query_expansion_variants": settings.enrichment.query_expansion_variants,
                    "fusion_enabled": effective_fusion,
                    "dedup_enabled": effective_dedup,
                    "dedup_method": settings.enrichment.dedup_method,
                    "compression_enabled": effective_compression,
                    "compression_method": settings.enrichment.compression_method,
                    "compression_ratio": settings.enrichment.compression_ratio,
                },
                # v3.5.1: Adaptive budget manager support (FIX-BUDGET-001)
                # Uses ENV UBP_ADAPTIVE_BUDGET__ENABLED as default, client can override
                "adaptive_memory_enabled": effective_adaptive_budget,
            }

            # v3.5.0: Use dedicated Architect pipeline if available (BYPASSES standard provider resolution)
            # This ensures the Architect always uses its configured providers
            # without interference from system-wide settings or overrides
            pipeline_to_use = self.architect_pipeline if self.architect_pipeline else self.rag_pipeline

            if self.architect_pipeline:
                logger.info(
                    "Using dedicated Architect pipeline (isolated provider)",
                    extra={
                        "main_provider": architect_settings.provider,
                        "enrichment_provider": architect_settings.enrichment_provider,
                    }
                )
            else:
                logger.warning(
                    "Using standard RAG pipeline for Architect (fallback mode)",
                    extra={"reason": "dedicated pipeline not initialized"}
                )

            # =====================================================================
            # v3.6.0: ARCHITECT CONVERSATION CONTEXT SUPPORT
            # =====================================================================
            # Load conversation history for context continuity in follow-up questions
            # v4.1.0: Centralized summary-first context loading
            conversation_context = await self._load_conversation_context(
                user_id=user_id,
                conversation_id=conversation_id,
                ctx=ctx,
                route_label="ARCHITECT",
            )
            structured_memory_context = conversation_context  # for debug panel compatibility
            has_conversation_context = bool(conversation_context)

            # =====================================================================
            # v5.0 FEAT-MEM-003: Memory-Aware Query Rewriting for Architect
            # =====================================================================
            rewrite_result = None
            retrieval_query = query  # default
            if conversation_id and self.memory_module and hasattr(self.memory_module, 'rewrite_query'):
                try:
                    rewrite_result = await self.memory_module.rewrite_query(
                        session_id=conversation_id, query=query, ctx=ctx
                    )
                    if rewrite_result.get("rewrite_type") not in ("none", "error"):
                        retrieval_query = rewrite_result["query"]
                        logger.info(f"[REWRITER] Architect: '{query[:40]}' -> '{retrieval_query[:60]}'")
                except Exception as e:
                    logger.warning(f"[REWRITER] Architect rewrite failed: {e}")

            # Execute RAG query on the Architect KB
            # Model must be in full format: "provider/model_name"
            # v3.5.0: The dedicated pipeline ensures enrichment uses Architect's configured provider
            # v3.6.0: Now includes conversation context for follow-up questions
            architect_collections = [architect_settings.kb_name]
            architect_rag_config = {
                **architect_config,
                # v6.0.1: provider already in architect_config, no model needed
            }
            # v5.0: Inject _original_user_query so _generate() uses natural query for LLM prompt
            if retrieval_query != query:
                architect_rag_config["_original_user_query"] = query
                if rewrite_result and rewrite_result.get("metadata"):
                    architect_rag_config["_rewrite_focus"] = rewrite_result["metadata"].get("current_focus", "")
            # FEAT-TOOL-001: Inject tool settings + collections for tool executor
            architect_rag_config["_tool_settings"] = {
                "enabled": tool_search.enabled and tool_search.enabled_architect,
                "max_iterations": tool_search.max_iterations,
                "top_k": tool_search.top_k,
                "similarity_threshold": tool_search.similarity_threshold,
                "max_context_expansion_kb": tool_search.max_context_expansion_kb,
                "timeout_ms": tool_search.timeout_ms,
                "provider": tool_search.provider_architect,
                "web_search_available": bool(self.web_search_module),
            }
            architect_rag_config["_tool_collections"] = architect_collections
            architect_rag_config["_web_module"] = self.web_search_module
            logger.info("[RAG] op=architect_tool_chat collections=%d", len(architect_collections))
            result = await pipeline_to_use.chat(
                query=retrieval_query,
                collections=architect_collections,
                config=architect_rag_config,
                return_debug=True,
                conversation_context=conversation_context if conversation_context else None,
            )

            latency_ms = (time.time() - start_time) * 1000

            # v4.0.0: Memory debug info for Context Debug Panel
            _memory_debug = {
                "has_conversation_context": has_conversation_context,
                "conversation_context_length": len(conversation_context) if conversation_context else 0,
                "has_structured_memory": bool(structured_memory_context),
                "structured_memory_length": len(structured_memory_context) if structured_memory_context else 0,
            }

            # Format response
            # v2.2.2: Use effective values (may be overridden by client)
            effective_top_k = top_k if top_k is not None else architect_settings.top_k
            effective_max_tokens = (
                max_tokens
                if max_tokens is not None
                else architect_settings.response_max_tokens
            )
            response = {
                "answer": result.get("answer", "No answer generated"),
                "sources": result.get("sources", []),
                "model_used": architect_settings.provider,  # v6.0.1: provider name
                "kb_used": architect_settings.kb_name,
                "config_used": {
                    "top_k": effective_top_k,
                    "reranking_enabled": effective_reranking,
                    "relevance_threshold": architect_settings.relevance_threshold,
                    # v2.2.2: Include HyDE and max_tokens in config_used
                    "hyde_enabled": effective_hyde,
                    "hyde_document_type": architect_settings.hyde_document_type
                    if effective_hyde
                    else None,
                    "max_tokens": effective_max_tokens,
                    # v2.2.2: Include investigative settings
                    "investigative_enabled": effective_investigative,
                    "investigative_num_questions": settings.enrichment.investigative_num_questions
                    if effective_investigative
                    else None,
                    # v3.5.0: Include Architect's dedicated enrichment provider info
                    "enrichment_provider": architect_settings.enrichment_provider,
                    # v6.0.1: enrichment_model removed — provider-only
                    "pipeline_mode": "dedicated" if self.architect_pipeline else "fallback",
                    # v3.5.1: Include adaptive budget manager status (FIX-BUDGET-001)
                    "adaptive_budget_enabled": effective_adaptive_budget,
                    # v4.0.0: Additional enrichment features status
                    "query_expansion_enabled": effective_query_expansion,
                    "fusion_enabled": effective_fusion,
                    "dedup_enabled": effective_dedup,
                    "compression_enabled": effective_compression,
                },
                "conversation_id": conversation_id,
                "latency_ms": round(latency_ms, 2),
                "status": "success",
            }

            # Add debug info if available
            if "debug" in result:
                response["debug"] = {
                    "chunks_retrieved": result["debug"].get("chunks_retrieved", 0),
                    "chunks_after_rerank": result["debug"].get(
                        "chunks_after_rerank", 0
                    ),
                    "context_tokens": result["debug"].get("context_tokens", 0),
                    # v4.0.0: Full retrieved chunks for expandable display
                    "retrieved_chunks": result["debug"].get("retrieved_chunks", []),
                    # v2.2.2: Include HyDE debug info
                    "hyde_enabled": effective_hyde,
                    "hyde_document": result["debug"].get("hyde_document"),
                    # v3.6.0: Include conversation context info
                    "has_conversation_context": has_conversation_context,
                    # v4.0.0: Include investigation/enrichment debug info
                    "investigative_enabled": effective_investigative,
                    "investigative_questions": result["debug"].get("investigative_questions"),
                    "search_queries": result["debug"].get("search_queries"),
                    "expanded_queries": result["debug"].get("expanded_queries"),
                    "optimization_applied": result["debug"].get("optimization_applied"),
                    "optimization_time_ms": result["debug"].get("optimization_time_ms"),
                    # v4.0.0: Enrichment step stats and flags
                    "rerank_stats": result["debug"].get("rerank_stats"),
                    "fusion_stats": result["debug"].get("fusion_stats"),
                    "dedup_stats": result["debug"].get("dedup_stats"),
                    "compression_stats": result["debug"].get("compression_stats"),
                    "enrichment_flags": result["debug"].get("enrichment_flags"),
                    "tool_calls": result["debug"].get("tool_calls"),
                    "filters_applied": result["debug"].get("filters_applied"),
                    "filter_entities": result["debug"].get("filter_entities"),
                    "filter_confidence": result["debug"].get("filter_confidence"),
                    # v4.0.0: Adaptive Budget Manager debug
                    "adaptive_budget_debug": result["debug"].get("adaptive_budget_debug"),
                    "context_governor": result["debug"].get("context_governor"),
                    # v4.0.0: Context Debug Panel - prompt, memory, config info
                    "prompt_debug": result["debug"].get("prompt_debug"),
                    "context_chars": result["debug"].get("context_chars", 0),
                    "memory_debug": _memory_debug,
                    "architect_config": {k: v for k, v in architect_config.items() if k != "system_prompt"},
                    # v6.8.1: Phase timing breakdown for latency RCA
                    "phase_timings": result["debug"].get("phase_timings"),
                }
                # v5.0: Add rewrite debug
                if rewrite_result:
                    response["debug"]["query_rewrite"] = rewrite_result

            # =====================================================================
            # v3.6.0: SAVE ARCHITECT CONVERSATION TO HISTORY
            # =====================================================================
            # Save both user query and assistant response for context continuity
            answer_text = response.get("answer", "")

            # 1. Save to ConversationManager (conversation history)
            if self.conversation_manager:
                try:
                    # Save user message
                    await self.conversation_manager.add_message(
                        user_id=user_id,
                        conversation_id=conversation_id,
                        role="user",
                        content=query,
                        metadata={"route": "architect", "kb": architect_settings.kb_name},
                    )
                    # Save assistant response
                    await self.conversation_manager.add_message(
                        user_id=user_id,
                        conversation_id=conversation_id,
                        role="assistant",
                        content=answer_text,
                        metadata={
                            "route": "architect",
                            "provider": architect_settings.provider,
                            "sources_count": len(response.get("sources", [])),
                        },
                    )
                    logger.debug(
                        f"Architect: Saved conversation messages",
                        extra={"conversation_id": conversation_id}
                    )
                except Exception as e:
                    logger.warning(f"Architect: Could not save conversation history: {e}")

            # 2. Save to memory_module (structured memory for topic/intent tracking)
            if conversation_id and self.memory_module:
                try:
                    # Save user message (route tag triggers auto-create with source="architect")
                    await self.memory_module.add_message(
                        session_id=conversation_id,
                        role="user",
                        content=query,
                        metadata={"route": "architect"},
                        ctx=ctx,
                    )
                    # Save assistant response (full debug for architect history)
                    await self.memory_module.add_message(
                        session_id=conversation_id,
                        role="assistant",
                        content=answer_text,
                        metadata={
                            "route": "architect",
                            "provider": architect_settings.provider,
                            "debug_info": {
                                "config_used": response.get("config_used"),
                                "debug": response.get("debug"),
                                "kb_used": response.get("kb_used"),
                                "model_used": response.get("model_used"),
                            },
                        },
                        ctx=ctx,
                    )
                    self._invalidate_conversation_cache(user_id, conversation_id)
                    logger.debug(
                        f"Architect: Saved to structured memory",
                        extra={"conversation_id": conversation_id}
                    )
                except Exception as e:
                    logger.warning(f"Architect: Could not save to structured memory: {e}")

            # v3.6.0: Add context info to response
            response["has_conversation_context"] = has_conversation_context

            logger.info(
                "Architect Agent response generated",
                extra={
                    "conversation_id": conversation_id,
                    "latency_ms": latency_ms,
                    "sources_count": len(response.get("sources", [])),
                    "has_conversation_context": has_conversation_context,
                },
            )

            return response

        except Exception as e:
            logger.error(
                f"Architect Agent error: {e}",
                extra={
                    "conversation_id": conversation_id,
                    "error": str(e),
                },
                exc_info=True,
            )
            return {
                "status": "error",
                "message": f"Architect Agent error: {str(e)}",
                "error_code": "ARCHITECT_ERROR",
                "conversation_id": conversation_id,
            }

    # =========================================================================
    # v2.3: INTERACTIVE ANALYST - REPORT OPERATIONS
    # =========================================================================

    async def start_report(
        self,
        query: str,
        template_id: Optional[str] = None,
        collections: Optional[List[str]] = None,
        force_dynamic: bool = False,
        artifact_type: Optional[str] = None,
        artifact_formats: Optional[List[str]] = None,
        format_settings: Optional[Dict[str, Any]] = None,
        execution_mode: str = "report",
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.3: Start an interactive report generation session.

        Detects template from query and proposes a structured plan.
        v2.4: Supports dynamic planning with force_dynamic flag.
        v2.5.1: Supports multi-format artifact generation.
        v5.1.0: Supports execution_mode (report, insight, exploratory).

        Args:
            query: Report request query
            template_id: Optional template ID to force
            collections: Optional collections to use for RAG
            force_dynamic: Force dynamic planning (v2.4)
            artifact_type: Output type (report, computo_metrico, capitolato, presentazione)
            artifact_formats: List of output formats (docx, xlsx, csv, pptx, md)
            format_settings: Format-specific settings from formats.yaml
            execution_mode: Pipeline mode — report (default), insight, exploratory
            ctx: Security context

        Returns:
            Session info with proposal for user approval
        """
        # v5.1.0: Validate execution_mode
        valid_modes = ("report", "insight", "exploratory")
        if execution_mode not in valid_modes:
            execution_mode = "report"
        # Check if report agents are available
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available. Check module initialization.",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        # Extract user info from context
        user_id = ctx.user.user_id if ctx and hasattr(ctx, "user") else "anonymous"
        client_id = ctx.user.client_id if ctx and hasattr(ctx, "user") else None

        try:
            # Create conversation for this report session
            conversation_id = str(uuid.uuid4())

            # Start session
            session = await self.report_session_manager.start_session(
                user_id=user_id,
                query=query,
                conversation_id=conversation_id,
                client_id=client_id,
                collections=collections,
                preferred_template=template_id,
                force_dynamic=force_dynamic,
                execution_mode=execution_mode,
            )

            # v2.5.1: Store artifact configuration in session metadata
            # Note: template-based values are already set in report_session.py start_session()
            # Only override if explicitly passed to this method
            if artifact_type:
                session.metadata["artifact_type"] = artifact_type
            if artifact_formats:
                session.metadata["artifact_formats"] = artifact_formats
            elif "artifact_formats" not in session.metadata:
                # Only set defaults if template didn't provide formats
                default_formats = {
                    "report": ["docx"],
                    "computo_metrico": ["xlsx"],
                    "capitolato": ["docx"],
                    "presentazione": ["pptx"],
                    "data_export": ["csv"],
                    "markdown": ["md"],
                }
                session.metadata["artifact_formats"] = default_formats.get(
                    session.metadata.get("artifact_type", "report"), ["docx"]
                )
            if format_settings:
                session.metadata["format_settings"] = format_settings

            await self.report_session_manager._save_session(session)

            # Build proposal text
            proposal_lines = [
                f"**Report Plan: {session.plan.template_name}**",
                f"Subject: {session.plan.subject}",
                "",
                "**Sections:**",
            ]
            for i, section in enumerate(session.plan.sections, 1):
                source_emoji = {
                    "rag_only": "📚",
                    "web_only": "🌐",
                    "rag_first": "📚→🌐",
                    "web_first": "🌐→📚",
                    "mixed": "📚+🌐",
                    "llm_reasoning": "🧠",
                }.get(section.source_preference.value, "📄")
                proposal_lines.append(f"{i}. {section.title} {source_emoji}")
                if section.description:
                    proposal_lines.append(f"   _{section.description}_")

            proposal_lines.extend(
                [
                    "",
                    "Reply with:",
                    "- **approve** to start research",
                    "- **modify [changes]** to adjust the plan",
                    "- **cancel** to abort",
                ]
            )

            return {
                "session_id": session.session_id,
                "state": session.state.value,
                "template_id": session.plan.template_id,
                "template_name": session.plan.template_name,
                "sections": [
                    {
                        "title": s.title,
                        "description": s.description,
                        "source_preference": s.source_preference.value,
                        "required": s.required,
                    }
                    for s in session.plan.sections
                ],
                "proposal": "\n".join(proposal_lines),
                "conversation_id": session.conversation_id,
                "dynamic_planning": session.plan.template_id == "dynamic",
                "status": "success",
            }

        except Exception as e:
            logger.error(f"start_report error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to start report: {str(e)}",
                "error_code": "START_REPORT_ERROR",
            }

    async def get_report_session(
        self,
        session_id: Optional[str] = None,
        conversation_id: Optional[str] = None,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.3: Get current state of a report session.

        Args:
            session_id: Optional specific session ID
            conversation_id: Optional conversation ID to find session
            ctx: Security context

        Returns:
            Session state and progress
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        user_id = ctx.user.user_id if ctx and hasattr(ctx, "user") else "anonymous"

        try:
            session = None

            if session_id:
                session = await self.report_session_manager.get_session(session_id)
            elif conversation_id:
                session = await self.report_session_manager.get_active_session(
                    user_id=user_id,
                    conversation_id=conversation_id,
                )
            else:
                # Get any active session for user
                session = await self.report_session_manager.get_active_session(
                    user_id=user_id,
                )

            if not session:
                return {
                    "status": "not_found",
                    "message": "No active report session found",
                }

            # Calculate progress from metadata (sections_data stored there if available)
            sections_data = session.metadata.get("sections_data", {})
            completed = sum(
                1 for s in sections_data.values() if s.get("status") == "completed"
            )
            total = len(session.plan.sections) if session.plan else 0

            # v2.5.1: Include artifacts information
            artifacts = session.metadata.get("artifacts", [])
            artifact = session.metadata.get("artifact")  # Legacy single artifact

            response = {
                "session_id": session.session_id,
                "state": session.state.value,
                "plan": {
                    "template_id": session.plan.template_id,
                    "template_name": session.plan.template_name,
                    "subject": session.plan.subject,
                    "sections": [
                        {
                            "title": s.title,
                            "description": s.description,
                            "source_preference": s.source_preference.value,
                        }
                        for s in session.plan.sections
                    ],
                },
                "progress": {
                    "completed": completed,
                    "total": total,
                    "percentage": round(completed / total * 100, 1) if total > 0 else 0,
                },
                "created_at": session.created_at,
                "expires_at": session.expires_at,
                "conversation_id": session.conversation_id,
                # v2.5.1: Artifact information
                "artifact_type": session.metadata.get("artifact_type", "report"),
                "artifact": artifact,  # Legacy single artifact (backward compat)
                "artifacts": artifacts,  # v2.5.1: All generated artifacts
                "execution_mode": session.metadata.get("execution_mode", "report"),  # v5.1.0
                "status": "success",
            }

            return response

        except Exception as e:
            logger.error(f"get_report_session error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to get session: {str(e)}",
                "error_code": "GET_SESSION_ERROR",
            }

    # =========================================================================
    # v2.6: ENRICHMENT CONFIGURATION ENDPOINTS
    # =========================================================================

    async def approve_structure(
        self,
        session_id: str,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.6: Approve report structure and transition to enrichment configuration.

        This endpoint is called after user reviews the proposed report structure.
        It transitions the session from AWAITING_APPROVAL to AWAITING_ENRICHMENT_CONFIG
        and returns the sections with default enrichment presets for UI rendering.

        Args:
            session_id: Report session ID
            ctx: Security context

        Returns:
            Dict with sections and default enrichment configs for UI
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            # Get session
            session = await self.report_session_manager.get_session(session_id)
            if not session:
                return {
                    "status": "error",
                    "message": "Session not found",
                    "error_code": "SESSION_NOT_FOUND",
                }

            # Validate state transition
            if session.state != ReportState.AWAITING_APPROVAL:
                return {
                    "status": "error",
                    "message": f"Invalid state for approval: {session.state.value}. Expected: awaiting_approval",
                    "error_code": "INVALID_STATE",
                }

            # Transition to AWAITING_ENRICHMENT_CONFIG
            session.state = ReportState.AWAITING_ENRICHMENT_CONFIG
            await self.report_session_manager._save_session(session)

            # Build sections with default enrichment configs based on source_preference
            sections_with_enrichment = []
            for i, section in enumerate(session.plan.sections):
                # Generate default enrichment based on source preference
                default_config = SectionEnrichmentConfig.from_preset(
                    preset="standard",
                    source_preference=section.source_preference,
                )

                sections_with_enrichment.append({
                    "index": i,
                    "title": section.title,
                    "description": section.description,
                    "source_preference": section.source_preference.value,
                    "suggested_queries": section.suggested_queries,
                    "enrichment_config": default_config.to_dict(),
                })

            logger.info(
                f"[REPORT] Structure approved, awaiting enrichment config: {session_id}",
                extra={"sections_count": len(sections_with_enrichment)},
            )

            return {
                "status": "success",
                "session_id": session_id,
                "state": session.state.value,
                "sections": sections_with_enrichment,
                "presets_available": ["fast", "standard", "quality", "custom"],
                "enrichment_options": {
                    "rerank_enabled": {
                        "label": "Rerank",
                        "description": "Riordina risultati per rilevanza (cross-encoder)",
                        "impact": "fast",
                    },
                    "query_expansion_enabled": {
                        "label": "Expand",
                        "description": "Espande query con sinonimi e varianti",
                        "impact": "medium",
                    },
                    "hyde_enabled": {
                        "label": "HyDE",
                        "description": "Genera documento ipotetico per migliore retrieval",
                        "impact": "slow",
                    },
                    "investigative_enabled": {
                        "label": "Investigate",
                        "description": "Decompone in sub-query investigative",
                        "impact": "medium",
                    },
                    "metadata_injection_enabled": {
                        "label": "Metadata",
                        "description": "Arricchisce contesto con metadata",
                        "impact": "fast",
                    },
                    "debug_enabled": {
                        "label": "Debug",
                        "description": "Mostra dettagli esecuzione worker in tempo reale",
                        "impact": "none",
                    },
                },
                "message": "Structure approved. Configure enrichment settings and click 'Start Research'.",
            }

        except Exception as e:
            logger.error(f"approve_structure error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to approve structure: {str(e)}",
                "error_code": "APPROVE_STRUCTURE_ERROR",
            }

    async def configure_enrichment(
        self,
        session_id: str,
        section_configs: List[Dict[str, Any]],
        global_debug: bool = False,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.6: Configure enrichment for each section and start research.

        This endpoint receives the user's enrichment configuration from the UI
        checkbox panel and applies it to each section. Then starts the swarm
        execution with per-section enrichment settings.

        Args:
            session_id: Report session ID
            section_configs: List of enrichment configs per section
                [{"index": 0, "rerank_enabled": true, ...}, ...]
            global_debug: Master debug switch for all sections
            ctx: Security context

        Returns:
            Dict with execution status and debug event stream info
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            # Get session
            session = await self.report_session_manager.get_session(session_id)
            if not session:
                return {
                    "status": "error",
                    "message": "Session not found",
                    "error_code": "SESSION_NOT_FOUND",
                }

            # Validate state
            if session.state != ReportState.AWAITING_ENRICHMENT_CONFIG:
                return {
                    "status": "error",
                    "message": f"Invalid state: {session.state.value}. Expected: awaiting_enrichment_config",
                    "error_code": "INVALID_STATE",
                }

            # Apply enrichment configs to sections
            for config_data in section_configs:
                section_index = config_data.get("index", config_data.get("section_index"))
                if section_index is not None and 0 <= section_index < len(session.plan.sections):
                    # Build enrichment config
                    enrichment_config = SectionEnrichmentConfig(
                        rerank_enabled=config_data.get("rerank_enabled", True),
                        query_expansion_enabled=config_data.get("query_expansion_enabled", False),
                        hyde_enabled=config_data.get("hyde_enabled", False),
                        investigative_enabled=config_data.get("investigative_enabled", False),
                        metadata_injection_enabled=config_data.get("metadata_injection_enabled", False),
                        debug_enabled=config_data.get("debug_enabled", False) or global_debug,
                    )
                    session.plan.sections[section_index].enrichment_config = enrichment_config

            # Store global debug flag in metadata
            session.metadata["global_debug_enabled"] = global_debug
            session.metadata["enrichment_configured_at"] = datetime.utcnow().isoformat()

            # Initialize debug event batch if any debug is enabled
            any_debug = global_debug or any(
                s.enrichment_config and s.enrichment_config.debug_enabled
                for s in session.plan.sections
            )
            if any_debug:
                session.metadata["debug_events_key"] = f"ubp:report:debug:{session_id}"

            # Transition to RESEARCHING
            session.state = ReportState.RESEARCHING
            await self.report_session_manager._save_session(session)

            logger.info(
                f"[REPORT] Enrichment configured, starting research: {session_id}",
                extra={
                    "global_debug": global_debug,
                    "sections_configured": len(section_configs),
                },
            )

            # Start swarm execution (non-blocking)
            # The swarm will use the enrichment_config from each section
            asyncio.create_task(
                self.report_session_manager._execute_swarm(session)
            )

            return {
                "status": "success",
                "session_id": session_id,
                "state": session.state.value,
                "message": "Research started with configured enrichment settings.",
                "debug_enabled": any_debug,
                "debug_events_key": session.metadata.get("debug_events_key"),
                "sections_configured": [
                    {
                        "index": i,
                        "title": s.title,
                        "enrichment_config": s.enrichment_config.to_dict() if s.enrichment_config else None,
                    }
                    for i, s in enumerate(session.plan.sections)
                ],
            }

        except Exception as e:
            logger.error(f"configure_enrichment error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to configure enrichment: {str(e)}",
                "error_code": "CONFIGURE_ENRICHMENT_ERROR",
            }

    async def get_debug_events(
        self,
        session_id: str,
        from_index: int = 0,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.6: Get debug events for a report session.

        Called by frontend to poll debug events for real-time display.

        Args:
            session_id: Report session ID
            from_index: Start index for pagination (get events after this index)
            ctx: Security context

        Returns:
            Dict with debug events and pagination info
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            # Get debug events from Redis
            debug_key = f"ubp:report:debug:{session_id}"

            if not self.redis_client:
                return {
                    "status": "error",
                    "message": "Redis not available for debug events",
                    "error_code": "REDIS_UNAVAILABLE",
                }

            # Get events from Redis list
            events_raw = await self.redis_client.lrange(debug_key, from_index, -1)
            events = []
            for event_json in events_raw:
                try:
                    event_data = json.loads(event_json)
                    events.append(event_data)
                except json.JSONDecodeError:
                    continue

            # Get session state
            session = await self.report_session_manager.get_session(session_id)
            is_complete = session and session.state in [
                ReportState.COMPLETED,
                ReportState.WRITING,
                ReportState.REVIEW,
            ]

            return {
                "status": "success",
                "session_id": session_id,
                "events": events,
                "events_count": len(events),
                "from_index": from_index,
                "next_index": from_index + len(events),
                "is_complete": is_complete,
                "session_state": session.state.value if session else "unknown",
            }

        except Exception as e:
            logger.error(f"get_debug_events error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to get debug events: {str(e)}",
                "error_code": "GET_DEBUG_EVENTS_ERROR",
            }

    async def list_report_templates(
        self,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.3: List available report templates.

        Returns:
            List of template objects with metadata
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            templates = self.report_session_manager.get_available_templates()

            return {
                "templates": templates,
                "count": len(templates),
                "status": "success",
            }

        except Exception as e:
            logger.error(f"list_report_templates error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to list templates: {str(e)}",
                "error_code": "LIST_TEMPLATES_ERROR",
            }

    async def get_report_draft(
        self,
        session_id: str,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.4: Get the generated report draft for a session.

        Args:
            session_id: Report session ID
            ctx: Security context

        Returns:
            Dict with draft content, sections, and metrics
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            draft = await self.report_session_manager.get_report_draft(session_id)

            if not draft:
                return {
                    "status": "not_ready",
                    "message": "Draft not yet available. Swarm execution may still be in progress.",
                    "session_id": session_id,
                }

            # v2.5.1: Include artifacts in draft response
            # Artifacts are stored in session metadata, fetch session to get them
            session = await self.report_session_manager.get_session(session_id)
            artifacts = []
            artifact = None
            if session:
                artifacts = session.metadata.get("artifacts", [])
                artifact = session.metadata.get("artifact")

            return {
                "status": "success",
                "session_id": session_id,
                "draft_status": draft.get("status"),
                "plan_title": draft.get("plan_title"),
                "full_draft": draft.get("full_draft"),
                "sections": draft.get("sections", []),
                "metrics": draft.get("metrics", {}),
                "output_classification": draft.get("output_classification"),  # v5.1.0 G3
                "coherence_check": draft.get("coherence_check"),  # v5.1.0 G2
                "semantic_validation": draft.get("semantic_validation"),  # v5.1.0 G4
                "generated_at": draft.get("generated_at"),
                # v2.5.1: Artifact information
                "artifact": artifact,  # Legacy single artifact (backward compat)
                "artifacts": artifacts,  # v2.5.1: All generated artifacts
                # v5.1.2: Evidence abstraction + reasoning pass
                "evidence_matrices": draft.get("evidence_matrices"),
                "reasoning_output": draft.get("reasoning_output"),
                "execution_mode": draft.get("execution_mode", "report"),
            }

        except Exception as e:
            logger.error(f"get_report_draft error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to get draft: {str(e)}",
                "error_code": "GET_DRAFT_ERROR",
            }

    async def get_swarm_status(
        self,
        session_id: str,
        ctx: Any = None,
    ) -> Dict[str, Any]:
        """
        v2.4: Get the current swarm execution status for a session.

        Args:
            session_id: Report session ID
            ctx: Security context

        Returns:
            Dict with execution status and progress
        """
        if not REPORT_AGENTS_AVAILABLE or not self.report_session_manager:
            return {
                "status": "error",
                "message": "Report agents not available",
                "error_code": "REPORT_AGENTS_UNAVAILABLE",
            }

        try:
            status = await self.report_session_manager.get_swarm_status(session_id)
            status["session_id"] = session_id
            return status

        except Exception as e:
            logger.error(f"get_swarm_status error: {e}", exc_info=True)
            return {
                "status": "error",
                "message": f"Failed to get status: {str(e)}",
                "error_code": "GET_STATUS_ERROR",
            }

    async def set_pipeline_permission(
        self,
        entity_type: str,
        entity_id: Optional[str],
        pipeline_name: str,
        access_level: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """Set access permission for pipeline - ADMIN ONLY."""
        ctx = self._require_ctx(ctx)

        if not self._is_admin(ctx):
            raise HTTPException(403, "Admin access required")

        if entity_type not in ["user", "client", "default"]:
            return {"status": "error", "message": f"Invalid entity_type: {entity_type}"}

        if access_level not in ["read", "execute", "none"]:
            return {"status": "error", "message": f"Invalid access_level: {access_level}"}

        # Build Redis key
        if entity_type == "default":
            key = f"ubp:pipeline_acl:default:{pipeline_name}"
        else:
            if not entity_id:
                return {
                    "status": "error",
                    "message": f"entity_id required for {entity_type}",
                }
            key = f"ubp:pipeline_acl:{entity_type}:{entity_id}:{pipeline_name}"

        # Set permission
        await self.redis_client.set(key, access_level)

        return {
            "status": "success",
            "message": f"Permission set: {entity_type}/{entity_id or 'default'}/{pipeline_name} = {access_level}",
        }

    async def get_user_pipelines(self, ctx: Any = None, **kwargs) -> Dict[str, Any]:
        """Get pipelines accessible by user based on ACL."""
        ctx = self._require_ctx(ctx)
        user_id = getattr(ctx.user, "user_id", None)
        client_id = getattr(ctx.user, "client_id", None)

        # Admin sees all
        if self._is_admin(ctx):
            return {"pipelines": None, "count": -1}

        # Get all pipeline templates
        # FIX: Use absolute import (relative fails when module loaded dynamically)
        try:
            from ubp_enterprise_hybrid.modules.cores.pipeline_orchestrator.adapter import PIPELINE_TEMPLATES
            all_pipelines = list(PIPELINE_TEMPLATES.keys())
        except ImportError:
            logger.warning(
                "Could not import PIPELINE_TEMPLATES, defaulting to allow all"
            )
            return {"pipelines": None, "count": -1}

        # Check permissions (user > client > default)
        allowed = []
        for pipeline in all_pipelines:
            # Check user-level permission
            user_key = f"ubp:pipeline_acl:user:{user_id}:{pipeline}"
            user_perm = await self.redis_client.get(user_key)
            if user_perm:
                user_perm = (
                    user_perm.decode("utf-8")
                    if isinstance(user_perm, bytes)
                    else user_perm
                )
                if user_perm in ["read", "execute"]:
                    allowed.append(pipeline)
                continue

            # Check client-level permission
            if client_id:
                client_key = f"ubp:pipeline_acl:client:{client_id}:{pipeline}"
                client_perm = await self.redis_client.get(client_key)
                if client_perm:
                    client_perm = (
                        client_perm.decode("utf-8")
                        if isinstance(client_perm, bytes)
                        else client_perm
                    )
                    if client_perm in ["read", "execute"]:
                        allowed.append(pipeline)
                    continue

            # Check default permission
            default_key = f"ubp:pipeline_acl:default:{pipeline}"
            default_perm = await self.redis_client.get(default_key)
            if default_perm:
                default_perm = (
                    default_perm.decode("utf-8")
                    if isinstance(default_perm, bytes)
                    else default_perm
                )

            if default_perm in ["read", "execute"] or default_perm is None:
                # Default: allow if no explicit deny
                allowed.append(pipeline)

        return {"pipelines": allowed if allowed else None, "count": len(allowed)}

    # =========================================================================
    # BATCH INGESTION - Server-side batch file ingestion
    # =========================================================================

    # Redis key patterns for batch ingestion
    INGEST_JOB_KEY = "ubp:ingest:job:{job_id}"
    INGEST_DATA_KEY = "ubp:ingest:data:{job_id}:{idx}"
    INGEST_ACTIVE_KEY = "ubp:ingest:active:{user_id}:{collection_id}"
    INGEST_JOB_TTL = 3600  # 1 hour

    async def _save_ingest_job(self, job: IngestJob) -> None:
        """Persist job state to Redis."""
        key = self.INGEST_JOB_KEY.format(job_id=job.job_id)
        job.updated_at = datetime.utcnow().isoformat()
        await self.redis_client.setex(key, self.INGEST_JOB_TTL, json.dumps(job.to_dict()))

    async def _load_ingest_job(self, job_id: str) -> Optional[IngestJob]:
        """Load job state from Redis."""
        key = self.INGEST_JOB_KEY.format(job_id=job_id)
        data = await self.redis_client.get(key)
        if not data:
            return None
        raw = data.decode("utf-8") if isinstance(data, bytes) else data
        return IngestJob.from_dict(json.loads(raw))

    async def _get_active_job_id(self, user_id: str, collection_id: str) -> Optional[str]:
        """Get active job ID for reconnect."""
        key = self.INGEST_ACTIVE_KEY.format(user_id=user_id, collection_id=collection_id)
        data = await self.redis_client.get(key)
        if not data:
            return None
        return data.decode("utf-8") if isinstance(data, bytes) else data

    async def start_batch_ingest(
        self,
        collection_id: str,
        files: List[Dict[str, Any]],
        chunking_config: Optional[Dict[str, Any]] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Start a server-side batch ingestion job.

        Accepts all files at once, stores them in Redis, and processes
        them in a background asyncio task. The browser can close or
        reload — the job continues server-side.

        Args:
            collection_id: Target collection
            files: List of {filename, file_content, file_type}
            chunking_config: Optional {chunk_size, chunk_overlap}
            ctx: Security context (admin required)

        Returns:
            {status: "accepted", job_id, total_files}
        """
        ctx = self._require_admin(ctx, "start_batch_ingest")
        self._require_initialized()

        if not collection_id:
            return {"status": "error", "message": "collection_id is required"}
        if not files or len(files) == 0:
            return {"status": "error", "message": "files list is required and cannot be empty"}

        user_id = ctx.user.user_id

        # Check for duplicate active job
        existing_job_id = await self._get_active_job_id(user_id, collection_id)
        if existing_job_id:
            existing_job = await self._load_ingest_job(existing_job_id)
            if existing_job and existing_job.state in (IngestJobState.PENDING, IngestJobState.PROCESSING):
                return {
                    "status": "error",
                    "message": f"Active ingestion job already exists: {existing_job_id}",
                    "existing_job_id": existing_job_id,
                }

        # Build job
        job_id = str(uuid.uuid4())
        now = datetime.utcnow().isoformat()

        file_entries = []
        for i, f in enumerate(files):
            entry = IngestFileEntry(
                filename=f.get("filename", f"file_{i}"),
                file_type=f.get("file_type", "text/plain"),
                size=len(f.get("file_content", "")),
            )
            file_entries.append(entry)

            # Store file content in Redis (separate key per file for memory efficiency)
            data_key = self.INGEST_DATA_KEY.format(job_id=job_id, idx=i)
            await self.redis_client.setex(data_key, self.INGEST_JOB_TTL, f.get("file_content", ""))

        job = IngestJob(
            job_id=job_id,
            user_id=user_id,
            collection_id=collection_id,
            state=IngestJobState.PENDING,
            files=file_entries,
            created_at=now,
            updated_at=now,
            chunking_config=chunking_config,
        )

        # Save job state
        await self._save_ingest_job(job)

        # Set active pointer for reconnect
        active_key = self.INGEST_ACTIVE_KEY.format(user_id=user_id, collection_id=collection_id)
        await self.redis_client.setex(active_key, self.INGEST_JOB_TTL, job_id)

        # Launch background task
        asyncio.create_task(self._execute_batch_ingest(job_id))

        logger.info(
            f"[BATCH_INGEST] Job {job_id} started: {len(files)} files -> {collection_id}",
            extra={"job_id": job_id, "collection_id": collection_id, "total_files": len(files)},
        )

        return {
            "status": "accepted",
            "job_id": job_id,
            "total_files": len(files),
        }

    async def get_batch_ingest_status(
        self,
        job_id: Optional[str] = None,
        collection_id: Optional[str] = None,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Get batch ingestion job status. Supports reconnect by collection_id.

        Args:
            job_id: Direct job ID lookup
            collection_id: Lookup active job for this collection (reconnect)
            ctx: Security context (admin required)

        Returns:
            Job status with per-file details
        """
        ctx = self._require_admin(ctx, "get_batch_ingest_status")
        user_id = ctx.user.user_id

        # Resolve job_id from collection_id if needed
        if not job_id and collection_id:
            job_id = await self._get_active_job_id(user_id, collection_id)

        if not job_id:
            return {"status": "not_found", "message": "No active ingestion job found"}

        job = await self._load_ingest_job(job_id)
        if not job:
            return {"status": "not_found", "message": f"Job {job_id} not found or expired"}

        return job.to_status_response()

    async def cancel_batch_ingest(
        self,
        job_id: str,
        ctx: Any = None,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Cancel an active batch ingestion job.

        Sets state to CANCELLED; the background task detects this before
        processing the next file and stops.

        Args:
            job_id: Job ID to cancel
            ctx: Security context (admin required)

        Returns:
            {status: "cancelled", job_id}
        """
        ctx = self._require_admin(ctx, "cancel_batch_ingest")

        if not job_id:
            return {"status": "error", "message": "job_id is required"}

        job = await self._load_ingest_job(job_id)
        if not job:
            return {"status": "not_found", "message": f"Job {job_id} not found or expired"}

        if job.state not in (IngestJobState.PENDING, IngestJobState.PROCESSING):
            return {
                "status": "error",
                "message": f"Job is already {job.state.value}, cannot cancel",
            }

        job.state = IngestJobState.CANCELLED
        await self._save_ingest_job(job)

        logger.info(f"[BATCH_INGEST] Job {job_id} cancelled by user")
        return {"status": "cancelled", "job_id": job_id}

    async def _execute_batch_ingest(self, job_id: str) -> None:
        """
        Background task: process all files in the batch sequentially.

        Follows the same pattern as _execute_swarm in the report system:
        - Reload state from Redis before each file
        - Check for cancellation
        - Per-file error isolation (one failure doesn't stop the batch)
        - Clean up file data from Redis after processing
        """
        try:
            job = await self._load_ingest_job(job_id)
            if not job:
                logger.error(f"[BATCH_INGEST] Job {job_id} not found at execution start")
                return

            job.state = IngestJobState.PROCESSING
            await self._save_ingest_job(job)

            # Load dedup index once for the entire batch (O(1) Redis call vs O(N))
            dedup_index = await self.qdrant_module.load_dedup_index_internal(job.collection_id)

            # FEAT-DKI-001 v2.0: Concurrent DKI pipeline buffers
            dki_text_buffer: List[str] = []      # Buffer for batch keyword extraction
            dki_all_texts: List[str] = []         # Buffer for summarize (if active)
            dki_pending_tasks: List[asyncio.Task] = []  # Concurrent tasks

            for i, file_entry in enumerate(job.files):
                # Reload job to check for cancellation
                job = await self._load_ingest_job(job_id)
                if not job or job.state == IngestJobState.CANCELLED:
                    logger.info(f"[BATCH_INGEST] Job {job_id} cancelled, stopping at file {i}")
                    # Mark remaining files as skipped
                    if job:
                        for j in range(i, len(job.files)):
                            if job.files[j].status == FileStatus.PENDING:
                                job.files[j].status = FileStatus.SKIPPED
                        await self._save_ingest_job(job)
                    break

                # Mark file as processing
                job.current_file_index = i
                job.files[i].status = FileStatus.PROCESSING
                job.files[i].started_at = datetime.utcnow().isoformat()
                await self._save_ingest_job(job)

                # Read file content from Redis
                data_key = self.INGEST_DATA_KEY.format(job_id=job_id, idx=i)
                raw_content = await self.redis_client.get(data_key)
                if not raw_content:
                    job.files[i].status = FileStatus.ERROR
                    job.files[i].error = "File content not found in Redis (expired?)"
                    job.files[i].completed_at = datetime.utcnow().isoformat()
                    job.error_count += 1
                    await self._save_ingest_job(job)
                    continue

                file_content = raw_content.decode("utf-8") if isinstance(raw_content, bytes) else raw_content

                try:
                    # Use ingest_file for full file type detection (PDF, DOCX, etc.)
                    # We call it internally, bypassing admin check since we already verified
                    result = await self._ingest_single_file(
                        collection_id=job.collection_id,
                        filename=file_entry.filename,
                        file_content=file_content,
                        file_type=file_entry.file_type,
                        user_id=job.user_id,
                        chunking_config=job.chunking_config,
                        dedup_index=dedup_index,
                    )

                    if result.get("status") == "duplicate":
                        job.files[i].status = FileStatus.DUPLICATE
                        job.files[i].document_id = result.get("document_id", "")
                        job.files[i].error = result.get("message", "Duplicate")
                        job.skip_count += 1
                    elif result.get("status") == "success":
                        job.files[i].status = FileStatus.SUCCESS
                        job.files[i].document_id = result.get("document_id", "")
                        job.files[i].chunks_count = result.get("chunks_count", 0)
                        job.success_count += 1
                        # Update dedup index for intra-batch duplicate detection
                        if result.get("content_hash"):
                            dedup_index.add(result["content_hash"])

                        # FEAT-DKI-001 v2.0: Accumulate text for concurrent DKI pipeline
                        extracted_text = result.get("extracted_text")
                        if self.keyword_manager.enabled and extracted_text:
                            dki_text_buffer.append(extracted_text)
                            if self.keyword_manager.summarize_enabled:
                                dki_all_texts.append(extracted_text)

                            # When buffer reaches batch_size, launch concurrent task
                            if len(dki_text_buffer) >= self.keyword_manager.batch_size:
                                texts_to_process = dki_text_buffer[:]  # Copy for closure safety
                                task = asyncio.create_task(
                                    self._safe_dki_extract(job.collection_id, texts_to_process)
                                )
                                dki_pending_tasks.append(task)
                                dki_text_buffer = []
                    else:
                        job.files[i].status = FileStatus.ERROR
                        job.files[i].error = result.get("message", "Unknown error")
                        job.error_count += 1

                except Exception as e:
                    logger.error(
                        f"[BATCH_INGEST] File {file_entry.filename} failed: {e}",
                        exc_info=True,
                    )
                    job.files[i].status = FileStatus.ERROR
                    job.files[i].error = str(e)
                    job.error_count += 1

                job.files[i].completed_at = datetime.utcnow().isoformat()

                # Delete file content from Redis to free memory
                await self.redis_client.delete(data_key)

                # Save progress
                await self._save_ingest_job(job)

                # Refresh TTL every 10 files
                if i > 0 and i % 10 == 0:
                    job_key = self.INGEST_JOB_KEY.format(job_id=job_id)
                    await self.redis_client.expire(job_key, self.INGEST_JOB_TTL)
                    active_key = self.INGEST_ACTIVE_KEY.format(
                        user_id=job.user_id, collection_id=job.collection_id
                    )
                    await self.redis_client.expire(active_key, self.INGEST_JOB_TTL)

            # FEAT-DKI-001 v2.0: Flush remaining DKI buffer
            if dki_text_buffer and self.keyword_manager.enabled:
                task = asyncio.create_task(
                    self._safe_dki_extract(job.collection_id, dki_text_buffer)
                )
                dki_pending_tasks.append(task)

            # Await all concurrent DKI tasks
            if dki_pending_tasks:
                results = await asyncio.gather(*dki_pending_tasks, return_exceptions=True)
                for idx, r in enumerate(results):
                    if isinstance(r, Exception):
                        logger.warning(f"DKI: Batch keyword task {idx} failed: {r}")

            # Summarize (if active) - after all keyword extraction
            if dki_all_texts and self.keyword_manager.summarize_enabled:
                try:
                    await self.keyword_manager.summarize_collection(
                        job.collection_id, dki_all_texts
                    )
                except Exception as e:
                    logger.warning(f"DKI: Summarize failed (non-blocking): {e}")

            # All files processed
            job = await self._load_ingest_job(job_id)
            if job and job.state != IngestJobState.CANCELLED:
                job.state = IngestJobState.COMPLETED if job.error_count < job.total_files else IngestJobState.FAILED
                await self._save_ingest_job(job)

                # Clear active pointer
                active_key = self.INGEST_ACTIVE_KEY.format(
                    user_id=job.user_id, collection_id=job.collection_id
                )
                await self.redis_client.delete(active_key)

            logger.info(
                f"[BATCH_INGEST] Job {job_id} finished: "
                f"{job.success_count} success, {job.error_count} errors, {job.skip_count} duplicates"
                if job else f"[BATCH_INGEST] Job {job_id} finished (job state lost)",
            )

        except Exception as e:
            logger.error(f"[BATCH_INGEST] Fatal error in job {job_id}: {e}", exc_info=True)
            try:
                job = await self._load_ingest_job(job_id)
                if job:
                    job.state = IngestJobState.FAILED
                    await self._save_ingest_job(job)
            except Exception:
                pass

    async def _safe_dki_extract(
        self, collection_id: str, texts: List[str]
    ) -> Dict[str, Any]:
        """
        FEAT-DKI-001 v2.0: Safe wrapper for concurrent keyword extraction.

        Non-blocking: catches all exceptions, returns status dict.
        Routes to single-doc or batch method based on text count.
        """
        try:
            if len(texts) == 1:
                return await self.keyword_manager.extract_and_store_keywords(
                    collection_name=collection_id, text=texts[0]
                )
            else:
                return await self.keyword_manager.extract_keywords_batch(
                    collection_name=collection_id, texts=texts
                )
        except Exception as e:
            logger.warning(f"DKI: Keyword extraction failed (non-blocking): {e}")
            return {"status": "error", "message": str(e)}

    async def _ingest_single_file(
        self,
        collection_id: str,
        filename: str,
        file_content: str,
        file_type: str,
        user_id: str,
        chunking_config: Optional[Dict[str, Any]] = None,
        dedup_index: Optional[set] = None,
    ) -> Dict[str, Any]:
        """
        Internal helper: ingest a single file without admin context check.
        Reuses the same extraction logic as ingest_file but called from
        the background task.
        """
        self._require_initialized()

        metadata = {
            "filename": filename,
            "source": "batch_ingest",
            "type": "document",
            "ingested_at": datetime.utcnow().isoformat(),
            "uploader_id": user_id,
        }

        # Detect file type and extract text (same logic as ingest_file)
        text = ""
        extraction_method = "direct"

        is_pdf = file_type == "application/pdf" or filename.lower().endswith(".pdf")
        is_json = file_type == "application/json" or filename.lower().endswith(".json")
        excel_mimes = [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ]
        is_excel = file_type in excel_mimes or filename.lower().endswith((".xlsx", ".xls"))
        is_docx = (
            file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or filename.lower().endswith(".docx")
        )
        is_csv = "csv" in file_type.lower() or filename.lower().endswith(".csv")
        is_html = "html" in file_type.lower() or filename.lower().endswith((".html", ".htm"))
        is_pptx = (
            file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or filename.lower().endswith(".pptx")
        )
        is_xml = "xml" in file_type.lower() or filename.lower().endswith(".xml")
        is_yaml = "yaml" in file_type.lower() or filename.lower().endswith((".yaml", ".yml"))

        if is_pdf:
            import base64
            import io
            from PyPDF2 import PdfReader

            pdf_bytes = base64.b64decode(file_content)
            pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
            text_parts = []
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            text = "\n\n".join(text_parts)
            extraction_method = "pypdf2"

        elif is_json:
            import json as json_module
            data = json_module.loads(file_content)
            text = json_module.dumps(data, indent=2, ensure_ascii=False)
            extraction_method = "json"

        elif is_excel:
            import base64
            import io
            from openpyxl import load_workbook

            excel_bytes = base64.b64decode(file_content)
            workbook = load_workbook(io.BytesIO(excel_bytes), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_lines = [f"=== Sheet: {sheet_name} ==="]
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_values = [str(cell) if cell is not None else "" for cell in row]
                        sheet_lines.append(" | ".join(row_values))
                text_parts.append("\n".join(sheet_lines))
            workbook.close()
            text = "\n\n".join(text_parts)
            extraction_method = "openpyxl"

        elif is_docx:
            import base64
            import io
            from docx import Document as DocxDocument

            docx_bytes = base64.b64decode(file_content)
            doc = DocxDocument(io.BytesIO(docx_bytes))
            text = "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
            extraction_method = "python-docx"

        elif is_csv:
            import base64
            import io
            import csv

            csv_content = base64.b64decode(file_content).decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(csv_content))
            lines = []
            for row in reader:
                lines.append(" | ".join(row))
            text = "\n".join(lines)
            extraction_method = "csv"

        elif is_html:
            import base64
            from bs4 import BeautifulSoup

            html_content = base64.b64decode(file_content).decode("utf-8", errors="replace")
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            extraction_method = "beautifulsoup"

        elif is_pptx:
            import base64
            import io
            from pptx import Presentation

            pptx_bytes = base64.b64decode(file_content)
            prs = Presentation(io.BytesIO(pptx_bytes))
            slide_texts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text.append(shape.text_frame.text)
                slide_texts.append("\n".join(slide_text))
            text = "\n\n".join(slide_texts)
            extraction_method = "python-pptx"

        elif is_xml:
            import base64
            from lxml import etree

            xml_bytes = base64.b64decode(file_content)
            root = etree.fromstring(xml_bytes)
            text = etree.tostring(root, method="text", encoding="unicode")
            extraction_method = "lxml"

        elif is_yaml:
            import base64
            import yaml

            yaml_content = base64.b64decode(file_content).decode("utf-8", errors="replace")
            data = yaml.safe_load(yaml_content)
            text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
            extraction_method = "pyyaml"

        else:
            # Plain text (TXT, MD, etc.)
            text = file_content
            extraction_method = "direct"

        if not text or not text.strip():
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "error",
                "message": f"No text extracted from {filename} (extraction: {extraction_method})",
            }

        metadata["extraction_method"] = extraction_method
        metadata["original_size"] = len(file_content)
        metadata["extracted_size"] = len(text)

        # --- DEDUP CHECK ---
        content_hash = hashlib.md5(text.encode()).hexdigest()
        if dedup_index is not None:
            # Batch mode: use preloaded in-memory index (O(1) lookup, no Redis call)
            is_duplicate = content_hash in dedup_index
        else:
            # Single file mode: check Redis registry
            is_duplicate = bool(await self.qdrant_module.check_duplicate_internal(
                collection=collection_id,
                content_hash=content_hash,
                filename=filename,
            ))
        if is_duplicate:
            logger.info(
                f"DEDUP: Skipping duplicate '{filename}' (hash: {content_hash[:12]}...)"
            )
            return {
                "document_id": "",
                "chunks_count": 0,
                "status": "duplicate",
                "message": f"Already ingested (content_hash: {content_hash[:8]}...)",
            }

        # Generate document ID and index
        document_id = str(uuid.uuid4())
        _enrich_ingest_metadata(metadata, collection_id)
        result = await self.qdrant_module.add_document_internal(
            doc_id=document_id,
            text=text,
            metadata=metadata,
            collection=collection_id,
        )

        if result.get("status") != "indexed":
            return {
                "document_id": document_id,
                "chunks_count": 0,
                "status": "error",
                "message": result.get("error") or "Failed to index document",
            }

        chunks_count = int(result.get("chunks_count") or 0)

        # FEAT-DKI-001 v2.0: Keyword extraction moved to _execute_batch_ingest
        # for concurrent pipeline support. Return extracted text for batch accumulation.
        return {
            "document_id": document_id,
            "chunks_count": chunks_count,
            "status": "success",
            "content_hash": content_hash,
            "message": f"Ingested {filename}: {chunks_count} chunks ({extraction_method})",
            "extracted_text": text,
        }
